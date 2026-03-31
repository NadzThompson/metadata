# Databricks notebook source
# ingest_embeddings_ADLS_DOC_metadata.py
#
# NOVA RAG Pipeline — Ingestion, OCR, Metadata & Embedding
# Implements the Three-Rule metadata architecture for both regulatory and internal documents.
#
# NOVA Three Rules Integration:
#   Rule 1 (Embed): build_semantic_header() prepends metadata to chunk text before embedding
#                    so the vector encodes regulator, section path, normative weight.
#   Rule 2 (Index): Full NOVA field mapping in create_es_index_with_mapping() for filtering/boosting.
#   Rule 3 (Prompt): render_chunk_for_prompt() injects metadata headers for LLM reasoning.
#
# New capabilities added:
#   - Three-category file discovery: regulatory JSON, internal raw, auto
#   - process_regulatory_scraped_json() for pre-scraped external docs (OSFI, PRA, etc.)
#   - Structural metadata: normative_weight, paragraph_role, cross_references
#   - OCR confidence scoring with GPT-5-mini vision model
#   - Companion metadata JSON loading for internal documents
#   - Azure Cognitive Search index with NOVA fields
#
# Add current directory to Python path for imports when running standalone
import os
import sys

# ---------------------------------------------------------------------------

USE_FALLBACK = True

# Import NER processor and helper functions from the same directory
try:
    if USE_FALLBACK:
        # Use the full NER processor (Named NER processor from ner_processor.py)
        from ner_processor import ner_extract
        USE_FALLBACK = False
    else:
        # Fallback: try manual and import from ner_processor_gk_GPT
        from ner_processor_gk_GPT import ner_extract
except ImportError as e:
    print(f"Import error: {e}")
    USE_FALLBACK = True
    ner_extract = None

# ---------------------------------------------------------------------------

# Counts the files with content
file_count = 0
files_with_entities = 0
total_entities_extracted = 0

# ---------------------------------------------------------------------------

# ingest_embeddings_ADLS_DOC_metadata.py

import os
import json
import time
import logging
import hashlib
import requests
import traceback
import tempfile
import base64
import io
import re
from pathlib import Path
from typing import List, Dict, Optional, Tuple

import fitz  # PyMuPDF
from elasticsearch import Elasticsearch, helpers

from azure.identity import DefaultAzureCredential
from azure.storage.blob import BlobServiceClient, ContainerClient
from langchain_community.document_loaders import DirectoryLoader, TextLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter

import datetime

from azure.storage.blob import BlobServiceClient

from azure.ai.formrecognizer import DocumentAnalysisClient
from azure.core.credentials import AzureKeyCredential

import docx
import openpyxl

from bs4 import BeautifulSoup

from azure.storage.blob import BlobSasPermissions
from datetime import timedelta

import datetime

from bot import BotActivityFlow, RecipientActivityLog, log
from load_config import load_config

model_info = load_config()

# Add current directory to Python path for imports when running standalone
script_dir = Path(__file__).resolve().parent
sys.path.insert(0, str(script_dir))

# ---------------------------------------------------------------------------

# ============================================================================
# NOVA METADATA ARCHITECTURE — Three-Rule Constants and Helpers
# ============================================================================

from dataclasses import dataclass, field as dc_field

# ============================================================================
# NOVA DATA MODELS — CanonicalDocument and CanonicalUnit
# ============================================================================

@dataclass
class CanonicalUnit:
    """A single parsable unit (paragraph, section, slide, table batch) with NOVA metadata."""
    unit_id: str = ""
    unit_type: str = "paragraph"  # paragraph, section, slide, table_batch, image_ocr
    heading_path: list = dc_field(default_factory=list)
    section_path: str = ""
    text: str = ""
    citation_anchor: str = ""
    page_start: int = 0
    page_end: int = 0
    # Structural metadata (computed by enrich_chunk_with_structural_metadata)
    structural_level: str = ""       # chapter, section, subsection, paragraph, appendix
    section_number: str = ""         # e.g. "3.2.1"
    depth: int = 0
    parent_section_id: str = ""
    is_appendix: bool = False
    normative_weight: str = ""       # mandatory, advisory, permissive, informational
    paragraph_role: str = ""         # definition, requirement, procedure_step, example, exception, narrative
    cross_references: list = dc_field(default_factory=list)
    # Content boolean flags
    contains_definition: bool = False
    contains_formula: bool = False
    contains_requirement: bool = False
    contains_deadline: bool = False
    contains_assignment: bool = False
    contains_parameter: bool = False


@dataclass
class CanonicalDocument:
    """Normalized document with all NOVA metadata fields and a list of CanonicalUnits."""
    doc_id: str = ""
    doc_type: str = "internal"       # 'regulatory' or 'internal'
    source_type: str = ""            # 'regulatory' or 'internal'
    title: str = ""
    short_title: str = ""
    document_class: str = ""
    # Source tracking
    raw_path: str = ""
    canonical_json_path: str = ""
    raw_sha256: str = ""
    parser_version: str = "nova-pipeline-v2.0.0"
    quality_score: float = 0.0
    # Regulatory fields
    regulator: str = ""
    regulator_acronym: str = ""
    guideline_number: str = ""
    doc_family_id: str = ""
    version_id: str = ""
    version_label: str = ""
    version_sort_key: str = ""
    status: str = "active"
    current_version_flag: bool = True
    effective_date_start: str = ""
    effective_date_end: str = ""
    authority_class: str = ""
    authority_level: int = 0
    nova_tier: str = ""
    jurisdiction: str = ""
    sector: str = ""
    supersedes_doc_id: str = ""
    superseded_by_doc_id: str = ""
    # Internal fields
    business_owner: str = ""
    document_owner: str = ""
    approval_status: str = ""
    approval_date: str = ""
    review_date: str = ""
    next_review_date: str = ""
    confidentiality: str = ""
    business_line: str = ""
    function: str = ""
    audience: str = ""
    # Units
    units: list = dc_field(default_factory=list)

    def to_dict(self) -> dict:
        """Convert to a flat dict of all non-empty fields (for metadata merging)."""
        result = {}
        for key, value in self.__dict__.items():
            if key == "units":
                continue
            if value not in (None, "", [], 0, False) or key in ("current_version_flag",):
                result[key] = value
        return result


# ============================================================================
# NOVA CONSTANTS
# ============================================================================

# Vision model for OCR fallback (upgrade from GPT-4o to GPT-5-mini)
VISION_MODEL = os.getenv("VISION_MODEL", "gpt-5-mini-2025-08-07-eastus-dz")

# Embedding model configuration
EMBEDDING_MODEL = os.getenv("EMBEDDING_MODEL", "text-embedding-3-large")
EMBEDDING_DIMS = int(os.getenv("EMBEDDING_DIMS", "1536"))

# ADLS layout constants
BRONZE_EXTERNAL_PREFIX = "bronze/external/"
BRONZE_INTERNAL_PREFIX = "bronze/internal/"
SILVER_CANONICAL_PREFIX = "silver/canonical_json/"
SILVER_METADATA_PREFIX = "silver/metadata/"
GOLD_CHUNKS_PREFIX = "gold/chunks/"

# Known regulator prefixes for three-category discovery
KNOWN_REGULATORS = ["osfi", "pra", "boe", "bis", "bcbs", "sec", "occ", "fdic", "eba"]

# Rule 1 — Fields embedded into chunk text (changes semantic meaning)
EMBEDDED_FIELDS = [
    "doc_id", "short_title", "document_class", "heading_path", "section_path",
    "regulator", "guideline_number", "business_owner",
    "structural_level", "section_number", "normative_weight", "paragraph_role",
]

# Rule 2 — Fields stored as ES index fields (gates/boosts retrieval)
INDEX_FIELDS = [
    "doc_id", "title", "short_title", "document_class", "heading_path", "section_path",
    "citation_anchor", "regulator", "regulator_acronym", "doc_family_id",
    "version_id", "version_label", "guideline_number", "status",
    "current_version_flag", "effective_date_start", "effective_date_end",
    "authority_class", "authority_level", "nova_tier", "jurisdiction", "sector",
    "supersedes_doc_id", "superseded_by_doc_id",
    "business_owner", "document_owner", "approval_status", "confidentiality",
    "business_line", "audience",
    "contains_definition", "contains_formula", "contains_requirement",
    "contains_deadline", "contains_assignment", "contains_parameter",
    "structural_level", "section_number", "depth", "parent_section_id",
    "is_appendix", "normative_weight", "paragraph_role", "cross_references",
]

# Rule 3 — Fields injected into LLM prompt at inference time
PROMPT_INJECTED_FIELDS_REGULATORY = [
    "title", "citation_anchor", "regulator", "version_id", "version_label",
    "status", "current_version_flag", "effective_date_start", "effective_date_end",
    "authority_class", "nova_tier", "jurisdiction", "normative_weight", "paragraph_role",
]

PROMPT_INJECTED_FIELDS_INTERNAL = [
    "title", "version_id", "version_label", "current_version_flag",
    "business_owner", "approval_status", "effective_date_start", "effective_date_end",
    "business_line", "jurisdiction", "audience", "normative_weight", "paragraph_role",
]

# ---------------------------------------------------------------------------

# ============================================================================
# NOVA STRUCTURAL METADATA DETECTION HELPERS
# ============================================================================

def _classify_normative_weight(text):
    """Classify text by deontic modal verbs (mandatory/advisory/permissive/informational)."""
    text_lower = text.lower()
    mandatory_kw = [r'\bmust\b', r'\bshall\b', r'\brequired\b', r'\bmandatory\b', r'\bprohibited\b', r'\bshall not\b']
    advisory_kw = [r'\bshould\b', r'\bexpected\b', r'\brecommended\b', r'\bencouraged\b', r'\bshould not\b']
    permissive_kw = [r'\bmay\b', r'\bcan\b', r'\bpermitted\b', r'\boptional\b', r'\bat the discretion\b']

    for pat in mandatory_kw:
        if re.search(pat, text_lower):
            return "mandatory"
    for pat in advisory_kw:
        if re.search(pat, text_lower):
            return "advisory"
    for pat in permissive_kw:
        if re.search(pat, text_lower):
            return "permissive"
    return "informational"


def _classify_paragraph_role(text, heading=""):
    """Classify paragraph role: definition, requirement, procedure_step, example, exception, narrative."""
    text_lower = text.lower()
    heading_lower = heading.lower() if heading else ""

    # Definition patterns
    def_patterns = [r'"[^"]+" means\b', r'\bdefined as\b', r'\bfor the purpose[s]? of\b',
                    r'\brefers to\b', r'\bis defined\b']
    for pat in def_patterns:
        if re.search(pat, text_lower):
            return "definition"

    # Check heading for "definition" keyword
    if "definition" in heading_lower or "glossary" in heading_lower:
        return "definition"

    # Requirement patterns
    req_patterns = [r'\bmust\b', r'\bshall\b', r'\brequired to\b', r'\bobligation\b',
                    r'\bmandatory\b', r'\bprohibited\b']
    for pat in req_patterns:
        if re.search(pat, text_lower):
            return "requirement"

    # Procedure step (numbered/lettered lists)
    if re.match(r'^\s*(\d+[\.\)]\s|[a-z][\.\)]\s|\(i+\)\s|\([a-z]\)\s|•\s|–\s)', text):
        return "procedure_step"

    # Example
    if re.search(r'\bexample\b|\bfor instance\b|\be\.g\.\b|\billustrat', text_lower):
        return "example"

    # Exception
    if re.search(r'\bexcept\b|\bunless\b|\bnotwithstanding\b|\bexclud', text_lower):
        return "exception"

    return "narrative"


def _extract_cross_references(text):
    """Extract cross-references like 'see Section X', 'refer to Appendix A', guideline numbers."""
    refs = set()

    # Guideline references (B-20, IFRS 9, etc.)
    for m in re.finditer(r'\b([A-Z]{1,5}[-\s]?\d{1,4}(?:\.\d+)?)\b', text):
        refs.add(m.group(1))

    # Section/chapter/paragraph references
    for m in re.finditer(r'(?:section|chapter|paragraph|appendix|annex|schedule)\s+(\d+(?:\.\d+)*|\w)', text, re.IGNORECASE):
        refs.add(f"Section {m.group(1)}")

    return list(refs)[:20]  # Cap at 20


def _extract_section_number(heading_text):
    """Extract section number from heading text like '3.2.1 Capital Requirements'."""
    m = re.match(r'^(\d+(?:\.\d+)*)\s', heading_text.strip())
    if m:
        return m.group(1)
    # Roman numeral sections
    m = re.match(r'^((?:Part\s+)?(?:[IVXivx]+|[A-Z])(?:\.\d+)*)\s', heading_text.strip())
    if m:
        return m.group(1)
    return None


def _infer_structural_level(depth, heading_text=""):
    """Infer structural level from heading depth and text."""
    heading_lower = heading_text.lower() if heading_text else ""
    if any(kw in heading_lower for kw in ["appendix", "annex", "schedule"]):
        return "appendix"
    if depth == 0:
        return "chapter"
    elif depth == 1:
        return "section"
    elif depth == 2:
        return "subsection"
    else:
        return "paragraph"


def _compute_content_flags(text):
    """Compute boolean content flags for a chunk."""
    text_lower = text.lower()
    return {
        "contains_definition": bool(re.search(r'"[^"]+" means\b|\bdefined as\b|\brefers to\b', text_lower)),
        "contains_formula": bool(re.search(r'[=×÷∑∫√±]|\bformula\b|\bcalculat', text_lower)),
        "contains_requirement": bool(re.search(r'\bmust\b|\bshall\b|\brequired\b|\bobligation\b', text_lower)),
        "contains_deadline": bool(re.search(r'\bdeadline\b|\bdue date\b|\bno later than\b|\bby\s+\w+\s+\d{4}\b', text_lower)),
        "contains_assignment": bool(re.search(r'\bresponsible\b|\bassigned\b|\baccountable\b|\bdesignated\b', text_lower)),
        "contains_parameter": bool(re.search(r'\d+(\.\d+)?%|\bparameter\b|\bthreshold\b|\blimit\b|\bratio\b', text_lower)),
    }


def _infer_regulator_from_path(filepath):
    """Infer regulator name from file path (e.g., bronze/external/osfi/json/ -> OSFI)."""
    fp_lower = filepath.lower()
    regulator_map = {
        "osfi": "OSFI", "pra": "PRA", "boe": "Bank of England",
        "bis": "BIS", "bcbs": "BCBS", "sec": "SEC", "occ": "OCC",
        "fdic": "FDIC", "eba": "EBA",
    }
    for key, name in regulator_map.items():
        if f"/{key}/" in fp_lower or f"\\{key}\\" in fp_lower:
            return name
    return ""


def _infer_document_class_from_path(filepath):
    """Infer document_class from ADLS directory path.

    Uses directory naming conventions to classify documents:
        bronze/internal/policy/        -> 'policy'
        bronze/internal/procedures/    -> 'procedure'
        bronze/internal/research/      -> 'research_paper'
        bronze/internal/presentations/ -> 'presentation'
        bronze/internal/templates/     -> 'template'
        bronze/internal/memos/         -> 'memo'
        bronze/internal/reports/       -> 'report'
        bronze/internal/data/          -> 'structured_data'
        bronze/internal/reference/     -> 'reference_document'
    """
    fp_lower = filepath.lower().replace("\\", "/")
    # Directory-to-class mapping (ordered by specificity)
    path_class_map = {
        "/policy/": "policy",
        "/policies/": "policy",
        "/procedure/": "procedure",
        "/procedures/": "procedure",
        "/guideline/": "guideline",
        "/guidelines/": "guideline",
        "/research/": "research_paper",
        "/presentation/": "presentation",
        "/presentations/": "presentation",
        "/template/": "template",
        "/templates/": "template",
        "/memo/": "memo",
        "/memos/": "memo",
        "/report/": "report",
        "/reports/": "report",
        "/data/": "structured_data",
        "/reference/": "reference_document",
        "/training/": "training_material",
        "/syllabus/": "syllabus",
    }
    for path_pattern, doc_class in path_class_map.items():
        if path_pattern in fp_lower:
            return doc_class

    # Fallback: infer from file extension
    ext = os.path.splitext(filepath)[1].lower()
    ext_class_map = {
        ".pptx": "presentation",
        ".xlsx": "spreadsheet",
        ".xls": "spreadsheet",
        ".csv": "structured_data",
    }
    return ext_class_map.get(ext, "")


# ---------------------------------------------------------------------------

# ============================================================================
# NOVA RULE 1 — SEMANTIC HEADER (prepended to chunk text before embedding)
# ============================================================================

def build_semantic_header(chunk_metadata, source_type="auto"):
    """
    Build a compact bracket-delimited header prepended to chunk text BEFORE embedding.
    This is Rule 1: metadata that changes where the chunk sits in vector space.

    Format: [OSFI | B-20 | chapter_guideline | Capital Requirements > 3.2.1 | mandatory]

    Args:
        chunk_metadata: dict with NOVA metadata fields
        source_type: 'regulatory', 'internal', or 'auto'

    Returns:
        str: semantic header like '[OSFI | LAR Chapter 2 | mandatory]'
    """
    parts = []

    # Regulator or business owner
    reg = chunk_metadata.get("regulator", "")
    owner = chunk_metadata.get("business_owner", "")
    if reg:
        parts.append(reg)
    elif owner:
        parts.append(owner)

    # Short title or guideline number
    short = chunk_metadata.get("short_title", "")
    guideline = chunk_metadata.get("guideline_number", "")
    if short:
        parts.append(short)
    elif guideline:
        parts.append(guideline)

    # Document class
    doc_class = chunk_metadata.get("document_class", "")
    if doc_class:
        parts.append(doc_class)

    # Heading path (last 2 levels)
    heading_path = chunk_metadata.get("heading_path", [])
    if isinstance(heading_path, list) and heading_path:
        path_str = " > ".join(heading_path[-2:])
        parts.append(path_str)
    elif isinstance(heading_path, str) and heading_path:
        parts.append(heading_path)

    # Section number
    sec_num = chunk_metadata.get("section_number", "")
    if sec_num:
        parts.append(sec_num)

    # Normative weight (skip if informational — it's the default)
    nw = chunk_metadata.get("normative_weight", "")
    if nw and nw != "informational":
        parts.append(nw)

    if not parts:
        return ""

    return f"[{' | '.join(parts)}]"


# ---------------------------------------------------------------------------

# ============================================================================
# NOVA RULE 3 — PROMPT INJECTION (metadata rendered for LLM at inference time)
# ============================================================================

def render_chunk_for_prompt(chunk_metadata, chunk_text, source_type="auto"):
    """
    Render a retrieved chunk with metadata headers for LLM prompt injection.
    This is Rule 3: metadata the model needs to reason correctly.

    Args:
        chunk_metadata: dict with NOVA metadata fields
        chunk_text: the chunk content text
        source_type: 'regulatory' or 'internal'

    Returns:
        str: chunk text with metadata header for prompt
    """
    if source_type == "auto":
        source_type = chunk_metadata.get("source_type", "internal")

    if source_type == "regulatory":
        prompt_fields = PROMPT_INJECTED_FIELDS_REGULATORY
    else:
        prompt_fields = PROMPT_INJECTED_FIELDS_INTERNAL

    header_lines = []
    for field_name in prompt_fields:
        value = chunk_metadata.get(field_name, "")
        if value:
            display_name = field_name.upper().replace("_", " ")
            header_lines.append(f"{display_name}: {value}")

    if header_lines:
        header = "\n".join(header_lines)
        return f"--- SOURCE METADATA ---\n{header}\n--- CONTENT ---\n{chunk_text}"
    else:
        return chunk_text


# ---------------------------------------------------------------------------

# ============================================================================
# NOVA THREE-CATEGORY FILE DISCOVERY
# ============================================================================

def discover_paths_to_ingest(container_client, container_name):
    """
    Discover files in ADLS with three-category classification.

    Returns:
        list of tuples: (blob_path, path_category)
        where path_category is 'regulatory_json', 'internal_raw', or 'auto'
    """
    all_paths = []

    try:
        blobs = container_client.list_blobs()
        for blob in blobs:
            if blob.name.endswith("/"):
                continue

            path = blob.name
            ext = os.path.splitext(path)[1].lower()

            # Category 1: Regulatory scraped JSON
            is_regulatory_json = False
            for reg in KNOWN_REGULATORS:
                if f"external/{reg}/json/" in path.lower() and ext == ".json":
                    is_regulatory_json = True
                    break

            if is_regulatory_json:
                all_paths.append((path, "regulatory_json"))
            elif BRONZE_INTERNAL_PREFIX in path.lower():
                all_paths.append((path, "internal_raw"))
            else:
                all_paths.append((path, "auto"))

    except Exception as e:
        log_and_print(f"Error discovering paths: {e}", "error")

    return all_paths


# ---------------------------------------------------------------------------

# ============================================================================
# NOVA REGULATORY SCRAPED JSON PARSER
# ============================================================================

def process_regulatory_scraped_json(json_content, filename, filepath):
    """
    Parse pre-scraped regulatory JSON (OSFI, PRA, etc.) into document chunks.
    These JSONs already contain full metadata + section content from web scrapers.

    Args:
        json_content: bytes or dict of the scraped JSON
        filename: file name
        filepath: full ADLS path

    Returns:
        list of LangChain Document objects with NOVA metadata
    """
    from langchain.schema import Document as LCDocument

    documents = []

    try:
        if isinstance(json_content, bytes):
            data = json.loads(json_content.decode('utf-8', errors='ignore'))
        elif isinstance(json_content, str):
            data = json.loads(json_content)
        else:
            data = json_content

        # Extract NOVA metadata from scraped JSON top-level fields
        nova_metadata = {
            "doc_id": data.get("doc_id", data.get("id", "")),
            "title": data.get("title", data.get("name", "")),
            "short_title": data.get("short_title", ""),
            "regulator": data.get("regulator", data.get("authority", _infer_regulator_from_path(filepath))),
            "regulator_acronym": data.get("regulator_acronym", ""),
            "guideline_number": data.get("guideline_number", data.get("reference_number", "")),
            "jurisdiction": data.get("jurisdiction", data.get("country", "")),
            "authority_class": data.get("authority_class", ""),
            "nova_tier": data.get("nova_tier", data.get("tier", "")),
            "status": data.get("status", data.get("doc_status", "active")),
            "effective_date_start": data.get("effective_date_start", data.get("effective_date", "")),
            "effective_date_end": data.get("effective_date_end", ""),
            "document_class": data.get("document_class", data.get("doc_type", "")),
            "version_id": data.get("version_id", data.get("version", "")),
            "version_label": data.get("version_label", ""),
            "current_version_flag": data.get("current_version_flag", True),
            "sector": data.get("sector", ""),
            "doc_family_id": data.get("doc_family_id", ""),
            "supersedes_doc_id": data.get("supersedes_doc_id", ""),
            "superseded_by_doc_id": data.get("superseded_by_doc_id", ""),
            "authority_level": data.get("authority_level", ""),
            "source_type": "regulatory",
            "source": filename,
            "file_path": filepath,
            "content_type": "regulatory_json",
        }

        # Extract sections/content from the scraped JSON
        sections = data.get("sections", [])
        if not sections:
            # Try alternative content structures
            content = data.get("content", data.get("text", data.get("body", "")))
            if isinstance(content, str) and content.strip():
                sections = [{"heading": nova_metadata["title"], "content": content}]
            elif isinstance(content, list):
                sections = content

        # Build documents from sections
        heading_path = []
        for sec_idx, section in enumerate(sections):
            if isinstance(section, str):
                section = {"content": section}

            heading = section.get("heading", section.get("title", ""))
            content = section.get("content", "")

            # Handle items/paragraphs within a section
            if not content and "items" in section:
                items = section["items"]
                if isinstance(items, list):
                    content = "\n".join([str(item) for item in items])

            if not content or not content.strip():
                continue

            # Update heading path
            if heading:
                heading_path = heading_path[:1] + [heading] if heading_path else [heading]

            # Compute structural metadata
            section_number = _extract_section_number(heading) if heading else None
            depth = min(sec_idx, 3)  # Approximate depth from section order

            section_metadata = {
                **nova_metadata,
                "heading": heading,
                "heading_path": " > ".join(heading_path),
                "section_path": " > ".join(heading_path),
                "section_number": section_number or "",
                "structural_level": _infer_structural_level(depth, heading),
                "depth": depth,
                "is_appendix": any(kw in heading.lower() for kw in ["appendix", "annex", "schedule"]) if heading else False,
                "normative_weight": _classify_normative_weight(content),
                "paragraph_role": _classify_paragraph_role(content, heading),
                "cross_references": _extract_cross_references(content),
                **_compute_content_flags(content),
            }

            doc = LCDocument(
                page_content=content.strip(),
                metadata=section_metadata
            )
            documents.append(doc)

        log_and_print(f"Regulatory JSON: {filename} -> {len(documents)} sections extracted")

    except Exception as e:
        log_and_print(f"Error processing regulatory JSON {filename}: {str(e)}", "error")

    return documents


# ---------------------------------------------------------------------------

# ============================================================================
# NOVA METADATA ENRICHMENT FOR INTERNAL DOCUMENTS
# ============================================================================

def load_pre_extracted_metadata(filepath, container_client, container_name):
    """
    Load pre-extracted metadata from silver/metadata/ (written by metadata_extraction.py).

    This is the PRIMARY metadata source for internal documents. The metadata_extraction.py
    script runs BEFORE ingestion and writes a per-document metadata JSON to silver/metadata/
    with all NOVA fields resolved via three-tier logic (enrichment_registry > native > defaults).

    Falls back to companion *_metadata.json alongside the raw file if no silver/ metadata exists.

    Args:
        filepath: path to the raw file in bronze/
        container_client: Azure container client
        container_name: ADLS container name

    Returns:
        dict: resolved NOVA metadata fields, or empty dict if no metadata found
    """
    # Generate the doc_id the same way metadata_extraction.py does
    source_file = os.path.basename(filepath)
    name = os.path.splitext(source_file)[0]
    slug = re.sub(r"[^a-z0-9]+", "_", name.lower()).strip("_")

    # Try silver/metadata/ first (written by metadata_extraction.py)
    silver_candidates = [
        f"{SILVER_METADATA_PREFIX}internal.{slug}.json",   # internal doc_id pattern
        f"{SILVER_METADATA_PREFIX}ext.{slug}.json",        # external doc_id pattern
        f"{SILVER_METADATA_PREFIX}{slug}.json",            # plain slug
    ]

    for candidate in silver_candidates:
        try:
            blob_client = container_client.get_blob_client(container=container_name, blob=candidate)
            data = blob_client.download_blob().readall()
            metadata_doc = json.loads(data.decode('utf-8', errors='ignore'))
            # The metadata_extraction.py output has resolved fields nested under "resolved_metadata"
            resolved = metadata_doc.get("resolved_metadata", metadata_doc)
            log_and_print(f"Loaded pre-extracted metadata from silver: {candidate}")
            return resolved
        except:
            continue

    # Fallback: companion *_metadata.json alongside the raw file
    base_name = os.path.splitext(filepath)[0]
    companion_candidates = [
        f"{base_name}_metadata.json",
        f"{base_name}.metadata.json",
        filepath.replace("/raw/", "/json/").replace(os.path.splitext(filepath)[1], ".json"),
    ]

    for candidate in companion_candidates:
        try:
            blob_client = container_client.get_blob_client(container=container_name, blob=candidate)
            data = blob_client.download_blob().readall()
            metadata = json.loads(data.decode('utf-8', errors='ignore'))
            log_and_print(f"Found companion metadata: {candidate}")
            return metadata
        except:
            continue

    log_and_print(f"No pre-extracted metadata found for {source_file} — using parser defaults only", "warning")
    return {}


def enrich_chunk_with_structural_metadata(doc, heading_path=None, pre_extracted_metadata=None):
    """
    Enrich a LangChain Document with NOVA structural metadata.
    Called after initial parsing, before chunking/embedding.

    Two sources of enrichment:
    1. Structural detection (computed from text): normative_weight, paragraph_role,
       cross_references, content flags — always computed fresh per chunk.
    2. Pre-extracted document-level metadata (from metadata_extraction.py):
       business_owner, confidentiality, effective_date_start, jurisdiction, etc.
       — merged once from the silver/metadata/ JSON.
    """
    text = doc.page_content
    heading = doc.metadata.get("heading", "")

    # --- Merge pre-extracted metadata (document-level fields from metadata_extraction.py) ---
    # These are fields the parser can't know — they come from the enrichment registry
    # or from the file's native properties, resolved by metadata_extraction.py.
    # Priority: pre_extracted > existing parser metadata > heuristic defaults
    if pre_extracted_metadata:
        doc_level_fields = [
            "doc_id", "title", "short_title", "document_class", "source_type",
            "business_owner", "document_owner", "confidentiality", "business_line",
            "audience", "approval_status", "effective_date_start", "effective_date_end",
            "review_date", "next_review_date", "version_id", "version_label",
            "current_version_flag", "doc_family_id", "jurisdiction", "status",
            "regulator", "regulator_acronym", "guideline_number",
            "authority_class", "authority_level", "nova_tier", "sector",
            "supersedes_doc_id", "superseded_by_doc_id",
        ]
        for field in doc_level_fields:
            if field in pre_extracted_metadata and pre_extracted_metadata[field]:
                if not doc.metadata.get(field):
                    doc.metadata[field] = pre_extracted_metadata[field]

    # --- Directory-path heuristic for document_class (fallback if still empty) ---
    if not doc.metadata.get("document_class"):
        file_path = doc.metadata.get("file_path", "")
        inferred_class = _infer_document_class_from_path(file_path)
        if inferred_class:
            doc.metadata["document_class"] = inferred_class

    # --- Structural detection (computed fresh per chunk from text content) ---
    doc.metadata["normative_weight"] = _classify_normative_weight(text)
    doc.metadata["paragraph_role"] = _classify_paragraph_role(text, heading)
    doc.metadata["cross_references"] = _extract_cross_references(text)

    if heading:
        doc.metadata["section_number"] = _extract_section_number(heading) or ""

    # Content flags
    flags = _compute_content_flags(text)
    doc.metadata.update(flags)

    # Heading path
    if heading_path:
        doc.metadata["heading_path"] = " > ".join(heading_path)
        doc.metadata["section_path"] = " > ".join(heading_path)

    # Structural level from heading depth
    depth = doc.metadata.get("depth", len(heading_path) - 1 if heading_path else 0)
    doc.metadata["structural_level"] = _infer_structural_level(depth, heading)
    doc.metadata["is_appendix"] = any(
        kw in h.lower() for h in (heading_path or []) for kw in ["appendix", "annex", "schedule"]
    )

    return doc


# ---------------------------------------------------------------------------

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),  # Also print to console
    ]
)

# Override any existing logging configuration
logger = logging.getLogger(__name__)

# Set root logger to ensure all logs are captured

# Suppress verbose HTTP request logs from Azure and other HTTP libraries
logging.getLogger("azure").setLevel(logging.WARNING)
logging.getLogger("azure.core.pipeline.policies.http_logging_policy").setLevel(logging.WARNING)

# Suppress verbose Azure blob logs
logging.getLogger("azure.storage.blob").setLevel(logging.WARNING)
logging.getLogger("azure.identity").setLevel(logging.WARNING)

# Suppress verbose Elasticsearch logs
logging.getLogger("elasticsearch").setLevel(logging.WARNING)
logging.getLogger("elastic_transport").setLevel(logging.WARNING)

# Suppress verbose urllib3 and requests logs
logging.getLogger("urllib3").setLevel(logging.WARNING)
logging.getLogger("requests").setLevel(logging.WARNING)

print(f"Logging initialized. Log file: {log_file}")

logger, LOG_FILE_PATH = setup_logger()

# ---------------------------------------------------------------------------

def log_and_print(message, level="info"):
    """Print message with timestamp"""
    timestamp = datetime.datetime.utcnow().strftime("%Y-%m-%d %H:%M:%S")
    print(f"[{timestamp}] {message}")
    if level == "info":
        logger.info(message)
    elif level == "warning":
        logger.warning(message)
    elif level == "error":
        logger.error(message)

# ---------------------------------------------------------------------------

# Helper function to help map the format

def extract_metadata_from_json_content(json_text):
    """Extract metadata fields from JSON content"""
    if not json_text:
        return {}

    if isinstance(json_text, dict):
        return json_text

    # Read as "Code"
    if '"head"' in json_text:
        return {}

    try:
        return json.loads(json_text)
    except (json.JSONDecodeError, TypeError):
        return {}

# ---------------------------------------------------------------------------

def test_metadata_extraction_from_json_content(json_metadata_text) -> Dict:
    """Extract metadata from JSON content from ADLS blob"""
    metadata = {}

    # A list of metadata fields to extract from JSON
    metadata_fields = [
        'title',
        'author',
        'date',
        'category',
        'subcategory',
        'source',
        'language',
        'version',
        'description',
        'keywords',
        'doc_type',
        'access_level',
        'created_date',
        'modified_date',
        'source_file',
        'document_type',
        'jurisdiction',
        'effective_date',
        'regulator',
        'sector',
        'content_type',
    ]

    # Extract specified fields
    try:
        parsed = json.loads(json_metadata_text) if isinstance(json_metadata_text, str) else json_metadata_text

        for field in metadata_fields:
            if field in parsed:
                metadata[field] = parsed[field]
    except (json.JSONDecodeError, TypeError, AttributeError) as e:
        log_and_print(f"JSON file found but no recognized metadata fields!", "warning")

    return metadata

# ---------------------------------------------------------------------------

logger.info(f"metadata extraction completed")

# ---------------------------------------------------------------------------

# Extract metadata fields from JSON text (if they exist)

def extract_metadata_from_json_content_enhanced(json_text):
    """Enhanced metadata extraction from JSON content"""
    if not json_text:
        return {}

    metadata_parsed = extract_metadata_from_json_content(json_text)

    if not metadata_parsed:
        # JSON file found but no recognized metadata fields
        return {}

    # Parse as dict
    metadata_result = {}

    for key, value in metadata_parsed.items():
        if value and str(value).strip():
            metadata_result[key] = value

    log_and_print(f"Parsed: created {len(metadata_result)} metadata fields from JSON")

    return metadata_result

# ---------------------------------------------------------------------------

# Build clean-based metadata parsing string to be passed to chain section

def parse_metadata_embedding_and_concatenated_metadata(text, chunk_metadata=None) -> str:
    """
    Build clean-based metadata parsing string to be passed to chain section
    """
    if not text:
        return ""

    metadata_sections = []

    # Start with metadata header
    metadata_header = f"Document metadata:\n"

    # Process metadata or empty string if no metadata available

    # Prepare heading/title section (use the heading text without "NLI" prefix)
    heading_text = ""
    if chunk_metadata:
        heading_text = chunk_metadata.get("title", "")
        if not heading_text:
            heading_text = chunk_metadata.get("heading", "")

    heading_prefix = ""
    if heading_text:
        heading_prefix = f"heading: {heading_text}\n"

    # Build heading with or without heading text
    heading_str = f"{heading_prefix}" if heading_text else ""

    heading_all = f"{metadata_header}{heading_str}" if heading_text else metadata_header

    building_str = f"{heading_all}{', '.join(metadata_sections)}" if metadata_sections else heading_all

    return building_str

# ---------------------------------------------------------------------------

# Prepare heading/section text for heading text without "NLI" prefix
def build_heading_text_with_metadata_prefix(heading_text, heading_path, section_number=None, metadata=None):
    """Build heading text with metadata for embedding prefix"""

    heading_str = ""
    if heading_text:
        heading_str = heading_text

    # Prepare heading/section point join for heading text without "NLI" prefix
    heading_path_str = " > ".join(heading_path) if heading_path else ""

    # If a cross headings entry with parent heading for embedding prefix
    if section_number:
        heading_str = f"Section {section_number}: {heading_str}"

    heading_all = f"{heading_path_str} > {heading_str}" if heading_path_str else heading_str

    return heading_all

# ---------------------------------------------------------------------------

# Counts the files with content entry fields > 0
# if current_entry_fields > 0:
#     pass

# ---------------------------------------------------------------------------

# LOCAL DESKTOP ENV verification
# Verify environment variables is LOCAL DESKTOP ENV

# ---------------------------------------------------------------------------

class ProcessingConfig:
    """Processing configuration class"""
    proxy_env = ["HTTP_PROXY", "HTTPS_PROXY", "http_proxy", "https_proxy"]

    cls_original_proxies = {}

    @classmethod
    def setup_proxy(cls):
        """Clear proxy settings for ADLS initialization phase"""

        # Clear proxy env vars for ADLS initialization (ADLS REQUIRES NO PROXY)
        cls.cls_original_proxies = {}

        for var in ["HTTP_PROXY", "HTTPS_PROXY", "http_proxy", "https_proxy"]:
            if var in os.environ:
                cls.cls_original_proxies[var] = os.environ[var]
                del os.environ[var]

    @classmethod
    def restore_proxy(cls):
        """Restore proxy settings for Elasticsearch and other services"""
        for var, val in cls.cls_original_proxies.items():
            os.environ[var] = val

    # Note: All files now in the configured perm ALL restriction/value as secondary bytes,
    # regardless of extension, 3000 chunks in 200 entries per second (limited to 3000 entries content layout)

# ---------------------------------------------------------------------------

# Set up Azure ADLS blob storage connection

try:
    credential = DefaultAzureCredential()

    tenant_id = os.getenv("TENANT_ID")
    client_id = os.getenv("CLIENT_ID")

    blob_service_client = BlobServiceClient(
        account_url=f"https://{os.getenv('STORAGE_ACCOUNT_NAME')}.blob.core.windows.net",
        credential=credential
    )

except Exception as e:
    logger.error(f"Error setting up ADLS: {e}")

# ---------------------------------------------------------------------------

# Set up Elasticsearch client connection

es_host = os.getenv("ES_HOST")
es_api_key = os.getenv("ES_API_KEY")

ELASTICSEARCH_ENDPOINT = os.getenv("ELASTICSEARCH_ENDPOINT")
ELASTICSEARCH_PASSWORD = os.getenv("ELASTICSEARCH_PASSWORD")

elastic_client = Elasticsearch(
    es_host,
    api_key=es_api_key,
    request_timeout=120,
    max_retries=3,
    retry_on_timeout=True,
)

# Elasticsearch Connection test
try:
    info = elastic_client.info()
    log_and_print(f"Successfully connected to Elasticsearch!")
except Exception as e:
    log_and_print(f"Connection attempt to Elasticsearch failed. Retrying on retry_default seconds...", "warning")
    log_and_print(f"Error connecting to Elasticsearch after max_retries attempts: {e}", "error")

# ---------------------------------------------------------------------------

# ES CONFIG AND BATCH TRACKER

ADLS_TOTAL_FILES_IN_ANALYZE = 0

# ---------------------------------------------------------------------------

# VIS CLASSES AND HELPER FUNCTIONS

def check_file_exists_in_es(es_client, index_name, file_path):
    """
    Check if a file_path already exists in the Elasticsearch index.
    """
    try:
        result = es_client.search(
            index=index_name,
            body={
                "query": {
                    "term": {"file_path.keyword": file_path}
                },
                "size": 0
            }
        )
        count = result["hits"]["total"]["value"]
        return count > 0
    except Exception as e:
        logger.error(f"Error checking if file exists: {e}")
        return False

# ---------------------------------------------------------------------------

def check_file_exists_report(file_path):
    """Check file existing status in data query"""

    file_exists = check_file_exists_in_es(elastic_client, index_name, file_path)

    if file_exists:
        log_and_print(f"File already exists! '{file_path}'")

    return file_exists

# ---------------------------------------------------------------------------

def check_file_exists_in_adls(fs_client, container_name, file_path):
    """Check if a file_path already exists in the Elasticsearch index."""
    try:
        blob_client = fs_client.get_blob_client(container=container_name, blob=file_path)
        blob_client.get_blob_properties()
        return True
    except:
        return False

def save_file_exists_to_adls(fs_client, container_name, file_path, content):
    """Save file exists report to ADLS"""
    try:
        blob_client = fs_client.get_blob_client(container=container_name, blob=file_path)
        blob_client.upload_blob(content, overwrite=True)
        log_and_print(f"File saved to ADLS: {file_path}")
    except Exception as e:
        log_and_print(f"Error saving to ADLS: {e}", "error")

# ---------------------------------------------------------------------------

# DOC PROCESSING HELPER FUNCTIONS

def clean_text(text_content) -> str:
    """
    Cleans text of HTML, CSS, JavaScript, headers, and other non-content elements.
    """
    if not text_content:
        return ""

    cleaned_text = text_content

    # Remove script tags
    cleaned_text = re.sub(r'<script[^>]*>.*?</script>', '', cleaned_text, flags=re.DOTALL)

    # Remove style tags
    cleaned_text = re.sub(r'<style[^>]*>.*?</style>', '', cleaned_text, flags=re.DOTALL)

    # Remove HTML tags
    cleaned_text = re.sub(r'<[^>]+>', '', cleaned_text)

    # Remove multiple whitespace
    cleaned_text = re.sub(r'\s+', ' ', cleaned_text).strip()

    return cleaned_text

# ---------------------------------------------------------------------------

def parse_blob_metadata_to_dict(blob_properties_metadata) -> dict:
    """Parse metadata from Azure Blob properties Metadata to standard dict"""

    if not blob_properties_metadata:
        return {}

    # Convert BlobProperties to clean dict
    metadata = {}
    for key, value in blob_properties_metadata.items():
        if value and str(value).strip():
            metadata[key] = str(value)

    return metadata

# ---------------------------------------------------------------------------

def clean_html_content(html_content) -> str:
    """Clean HTML content and extract text"""
    if not html_content:
        return ""

    # Create BeautifulSoup object with cleaned content
    soup = BeautifulSoup(html_content, "html.parser")

    # Remove script and style elements
    unwanted_tags = ["script", "style", "nav", "footer", "header", "aside"]
    for tag in soup.find_all(unwanted_tags):
        tag.decompose()

    text = soup.get_text(separator="\n")
    # Remove multiple blank lines
    text = re.sub(r'\n\s*\n', '\n\n', text)
    return text.strip()

# ---------------------------------------------------------------------------

def extract_table_as_text(table_element):
    """Extract HTML table as formatted text"""
    if not table_element:
        return ""

    rows = []
    for tr in table_element.find_all('tr'):
        cells = []
        for td in tr.find_all(['td', 'th']):
            cells.append(td.get_text(strip=True))
        rows.append(' | '.join(cells))

    return '\n'.join(rows)

# ---------------------------------------------------------------------------

# HTML PROCESSING HELPER FUNCTIONS

def process_html_with_ocr(html_content, file_path, sas_key, ssl_processor, encoded_metadata=None):
    """Process HTML content with OCR for embedded images"""

    if not html_content:
        return []

    # Clean HTML and extract text
    soup = BeautifulSoup(html_content, 'html.parser')

    # Remove unwanted elements
    for tag in soup.find_all(['script', 'style', 'nav', 'footer']):
        tag.decompose()

    # Extract structured content
    structured_content = extract_text_with_structure(
        soup, file_path, sas_key, ssl_processor
    )

    return structured_content

# ---------------------------------------------------------------------------

# Extract and process images from HTML using Vision API (GPT-5-mini)

def extract_and_process_images_from_url(html_content, base_path, sas_key, ssl_processor, blob_connection_string=None):
    """
    Extract and process images from HTML using Vision API (GPT-5-mini).

    Args:
        html_content: Raw HTML content
        base_path: Path to original HTML file (for relative URLs)
        sas_key: Azure SAS key for blob access
        ssl_processor: Any existing image processor to be used
        blob_connection_string: Azure Blob Storage connection string

    Returns:
        List of dictionaries with image OCR results
    """
    image_results = []

    soup = BeautifulSoup(html_content, 'html.parser')
    img_tags = soup.find_all('img')

    if not img_tags:
        return image_results

    for img in img_tags:
        img_src = img.get('src', '')
        if not img_src:
            continue

        # Get image data
        img_file_like_obj_tag = io.BytesIO()

        try:
            log_and_print(f"Found {len(img_tags)} image(s) in HTML")

            # Handle various image source types
            img_url_processed = False

            if img_src.startswith('data:'):
                # Base64 encoded image
                img_decoded = base64.b64decode(img_src.split(',')[1])
                img_file_like_obj_tag = io.BytesIO(img_decoded)
                img_url_processed = True

            elif img_src.startswith(('http://', 'https://')):
                # Try to download from URL
                try:
                    response = requests.get(img_src, timeout=10)
                    if response.status_code == 200:
                        img_file_like_obj_tag = io.BytesIO(response.content)
                        img_url_processed = True
                except:
                    log_and_print(f"Error downloading image from {img_src}", "warning")

            else:
                # Local file path relative to HTML
                full_img_path = os.path.join(os.path.dirname(base_path), img_src)
                try:
                    img_data = read_file_from_blob(full_img_path, sas_key)
                    img_file_like_obj_tag = io.BytesIO(img_data)
                    img_url_processed = True
                except:
                    log_and_print(f"Image file not found: {img_src}", "warning")

            if not img_url_processed:
                continue

            # Process image with Vision API
            def image_to_base64():
                img_bytes = img_file_like_obj_tag.getvalue()
                return base64.b64encode(img_bytes).decode('utf-8')

            img_base64 = image_to_base64()

            # Determine image type
            img_type = "image/png"  # default
            if img_src.lower().endswith(('.jpg', '.jpeg')):
                img_type = "image/jpeg"
            elif img_src.lower().endswith('.gif'):
                img_type = "image/gif"
            elif img_src.lower().endswith('.webp'):
                img_type = "image/webp"

            # OCR with Vision API
            content_text = img.get('alt', '') or "No alt text provided"

            prompt = f"""Analyze this image in detail:

            If this is a chart, graph, or diagram:
            - Describe the type of visualization
            - List all data points, values, and labels visible
            - Describe key insights or trends
            - Note any axes labels, legends, or titles

            If this is a table:
            - Extract the table in markdown format
            - Preserve all headers and values
            - Note any totals or calculations

            If this is text or other content:
            - Extract all visible text exactly as shown

            Be comprehensive and precise."""

            tag_content_text = ssl_processor.extract_text_from_image(img_file_like_obj_tag.getvalue(), detail="high")

            image_results.append(tag_content_text)

            log_and_print(f"Image {img_src} processed successfully")

        except Exception as e:
            log_and_print(f"Error processing image {img_src}: {str(e)}")
            log_and_print(f"Traceback: {traceback.format_exc()}")

    return image_results

# ---------------------------------------------------------------------------

def extract_text_with_structure(soup, base_path=None, sas_key=None, ssl_processor=None, blob_connection_string=None):
    """
    Extract text while preserving document structure including headings.
    Also processes any embedded images to extract their content.

    Args:
        soup: BeautifulSoup object with cleaned HTML
        base_path: Optional path to HTML file (for images)
        sas_key: Optional SAS key for blob access
        ssl_processor: Optional processor for processing images

    Returns:
        A list of structured content dicts with 'heading', 'content', 'type'
    """
    structured_content = []

    # Track section hierarchy
    current_heading = None
    current_section = []

    # Find main content area
    main_content = soup.find('main') or soup.find('article') or soup.find('div', {'role': 'main'})
    if not main_content:
        main_content = soup.body if soup.body else soup

    for element in main_content.descendants:
        if isinstance(element, str):
            text = element.strip()
            if not text:
                continue

            # Check if this is a heading
            if element.parent and element.parent.name in ('h1', 'h2', 'h3', 'h4', 'h5', 'h6'):
                # Save current section if it has content
                if current_section:
                    structured_content.append({
                        'heading': current_heading,
                        'content': '\n'.join(current_section),
                        'type': 'text'
                    })
                current_heading = text
                current_section = []

            elif element.parent and element.parent.name in ('p', 'li', 'td', 'th', 'span', 'div'):
                # Clean up whitespace
                element_text = ' '.join(text.split())
                if element_text:
                    current_section.append(element_text)

        elif hasattr(element, 'name'):
            # Check for tables
            if element.name == 'table':
                table_text = extract_table_as_text(element)
                if table_text:
                    structured_content.append({
                        'heading': current_heading,
                        'content': table_text,
                        'type': 'table'
                    })

    # Don't forget the last section
    if current_section:
        structured_content.append({
            'heading': current_heading,
            'content': '\n'.join(current_section),
            'type': 'text'
        })

    # Process images if ssl_processor is available
    if ssl_processor and base_path:
        image_results = extract_and_process_images_from_url(
            str(soup), base_path, sas_key, ssl_processor, blob_connection_string
        )

        if image_results:
            for img_result in image_results:
                structured_content.append({
                    'heading': current_heading,
                    'content': img_result,
                    'type': 'image',
                    'source': 'image_ocr'
                })

    # Merge content sections
    merged_content = []
    text_passages = []

    for item in structured_content:
        if item['type'] == 'text':
            text_passages.append(item)
        else:
            # Flush text passages first
            if text_passages:
                merged_content.extend(text_passages)
                text_passages = []
            merged_content.append(item)

    # Flush remaining
    if text_passages:
        merged_content.extend(text_passages)

    return merged_content

# ---------------------------------------------------------------------------

# IMAGE OCR PROCESSING FUNCTIONS

def image_to_base64(image_bytes):
    """Convert image bytes to base64 string"""
    return base64.b64encode(image_bytes).decode('utf-8')

def get_gpt4o_description(image_bytes, detail="high"):
    """Get image description using Vision API (GPT-5-mini). Function name retained for backward compatibility."""
    img_base64 = image_to_base64(image_bytes)

    # Determine image source types
    img_type = "image/png"

    try:
        from openai import AzureOpenAI

        client = AzureOpenAI(
            api_key=os.getenv("OPENAI_API_KEY"),
            api_version=os.getenv("OPENAI_API_VERSION"),
            azure_endpoint=os.getenv("OPENAI_API_BASE")
        )

        response = client.chat.completions.create(
            model=VISION_MODEL,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": """Analyze this image in detail:

If this is a chart, graph, or diagram:
- Describe the type of visualization
- List all data points, values, and labels visible
- Describe key insights or trends
- Note any axes labels, legends, or titles

If this is a table:
- Extract the table in markdown format
- Preserve all headers and values
- Note any totals or calculations

If this is text or other content:
- Extract all visible text exactly as shown

Be comprehensive and precise.
If this contains a table, text, or other visible content, respond with content returned"""
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{img_type};base64,{img_base64}",
                                "detail": detail
                            }
                        }
                    ]
                }
            ],
            max_tokens=16384
        )

        return response.choices[0].message.content

    except Exception as e:
        log_and_print(f"Error getting vision description: {str(e)}", "error")
        return ""

# ---------------------------------------------------------------------------

def extract_text_with_ocr_enhanced(image_bytes, ssl_processor=None, detail="high"):
    """
    Extract text while preserving document structure including images.
    Uses OCR processor or Vision API (GPT-5-mini) for image analysis.
    """
    image_results = []

    try:
        if ssl_processor:
            # Use OCR processor
            text = ssl_processor.extract_text_from_image(image_bytes, detail=detail)
            if text:
                image_results.append(text)
        else:
            # Fallback to Vision API
            text = get_gpt4o_description(image_bytes, detail=detail)
            if text:
                image_results.append(text)

    except Exception as e:
        log_and_print(f"Error in OCR processing: {str(e)}", "error")

    return image_results

# ---------------------------------------------------------------------------

# Extract and process images from PDF using Vision API / OCR

def extract_images_from_pdf(pdf_content, file_path, sas_key, ssl_processor, encoded_metadata=None):
    """
    Extract and process images from PDFs using Vision API (GPT-5-mini) / OCR.

    Args:
        pdf_content: Raw PDF content as bytes
        file_path: Path to original HTML file (for relative URLs)
        sas_key: Any existing Azure SAS key for blob access
        ssl_processor: Any existing image processor to be used on content

    Returns:
        List of dictionaries with image OCR results
    """
    image_results = []

    temp_file = None
    try:
        temp_file = tempfile.NamedTemporaryFile(suffix='.pdf', delete=False)
        temp_file.write(pdf_content)
        temp_file.flush()
        temp_file.close()

        pdf_doc = fitz.open(temp_file.name)

        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            image_list = page.get_images(full=True)

            if image_list:
                log_and_print(f"Found {len(image_list)} image(s) on page {page_num + 1}")

            for img_idx, img in enumerate(image_list):
                try:
                    xref = img[0]
                    base_image = pdf_doc.extract_image(xref)

                    if base_image:
                        img_bytes = base_image["image"]
                        img_ext = base_image.get("ext", "png")

                        # Check minimum image size
                        MIN_SIZE = 100
                        if len(img_bytes) < MIN_SIZE:
                            continue

                        # Process with OCR/Vision
                        if ssl_processor:
                            text = ssl_processor.extract_text_from_image(img_bytes, detail="high")
                        else:
                            text = get_gpt4o_description(img_bytes, detail="high")

                        if text and len(text.strip()) > 10:
                            image_results.append({
                                'page': page_num + 1,
                                'image_index': img_idx,
                                'content': text,
                                'type': 'image_ocr'
                            })

                            log_and_print(f"Image {img_idx} on page {page_num + 1} processed successfully")

                except Exception as e:
                    log_and_print(f"Error processing image {img_idx} on page {page_num + 1}: {str(e)}", "warning")

            # If no embedded images, try page rendering for scanned pages
            if not image_list:
                page_text = page.get_text("text")
                if not page_text or len(page_text.strip()) < 50:
                    # Likely a scanned page - render to image
                    try:
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for better OCR
                        img_bytes = pix.tobytes("png")

                        if ssl_processor:
                            text = ssl_processor.extract_text_from_image(img_bytes, detail="high")
                        else:
                            text = get_gpt4o_description(img_bytes, detail="high")

                        if text and len(text.strip()) > 10:
                            image_results.append({
                                'page': page_num + 1,
                                'image_index': 0,
                                'content': text,
                                'type': 'scanned_page_ocr'
                            })
                    except Exception as e:
                        log_and_print(f"Error rendering page {page_num + 1}: {str(e)}", "warning")

        pdf_doc.close()

    except Exception as e:
        log_and_print(f"Error extracting images from PDF: {str(e)}", "error")
    finally:
        if temp_file and os.path.exists(temp_file.name):
            os.unlink(temp_file.name)

    return image_results

# ---------------------------------------------------------------------------

# CHUNKING FUNCTIONS

def split_text_into_semantic_chunks(text_content, max_chunk_size=1500, min_chunk_size=200):
    """
    Split text into semantically coherent chunks based on document structure.
    Respects headings and paragraph boundaries to maintain context.

    Args:
        text_content: Text to be split into chunks
        max_chunk_size: Maximum size of each chunk
        min_chunk_size: Minimum size to avoid tiny chunks

    Returns:
        List of text chunks with heading context
    """
    chunks = []
    current_chunk = ""
    current_headings = []

    # Split by lines
    lines = text_content.split('\n')

    for i, line in enumerate(lines):
        stripped = line.strip()
        if not stripped:
            continue

        # Check if it's a heading
        is_heading = False
        if stripped.startswith('#'):
            is_heading = True
        elif len(stripped) < 100 and stripped.isupper():
            is_heading = True

        if is_heading:
            # If current chunk is big enough, save it
            if current_chunk and len(current_chunk) >= min_chunk_size:
                chunk_text = ""
                if current_headings:
                    chunk_text = " > ".join([h['text'] for h in current_headings]) + "\n\n"
                chunk_text += current_chunk
                chunks.append(chunk_text)
                current_chunk = ""

            # Update heading tracking
            level = 0
            if stripped.startswith('#'):
                level = len(stripped) - len(stripped.lstrip('#'))
            else:
                level = 1  # Treat uppercase as h1

            # Keep only headings at higher or same level
            current_headings = [h for h in current_headings if h.get('level', 0) < level]
            current_headings.append({'text': stripped.lstrip('#').strip(), 'level': level})

        else:
            # Regular content
            if len(current_chunk) + len(stripped) > max_chunk_size:
                # Need to split
                if current_chunk:
                    chunk_text = ""
                    if current_headings:
                        chunk_text = " > ".join([h['text'] for h in current_headings]) + "\n\n"
                    chunk_text += current_chunk
                    chunks.append(chunk_text)
                    current_chunk = ""

            current_chunk += stripped + "\n"

    # Don't forget the last chunk
    if current_chunk and len(current_chunk) >= min_chunk_size:
        chunk_text = ""
        if current_headings:
            chunk_text = " > ".join([h['text'] for h in current_headings]) + "\n\n"
        chunk_text += current_chunk
        chunks.append(chunk_text)

    # If no chunks created (content too small), create one
    if not chunks and current_chunk:
        chunk_text = ""
        if current_headings:
            chunk_text = " > ".join([h['text'] for h in current_headings]) + "\n\n"
        chunk_text += current_chunk
        if len(chunk_text) > 50:
            chunks.append(chunk_text)

    # Merge tiny trailing chunk
    if len(chunks) > 1 and len(chunks[-1]) < min_chunk_size:
        chunks[-2] += "\n" + chunks[-1]
        chunks.pop()

    return chunks

# ---------------------------------------------------------------------------

def create_text_chunks_with_context(text_content, max_chunk_size=1500, min_chunk_size=200, chunk_overlap=200):
    """
    Create text chunks with heading context and overlap.
    Tracks all headings in the current chunk.
    """
    chunks = []
    current_chunk = ""
    current_headings = []

    lines = text_content.split('\n')

    for i, line in enumerate(lines):
        stripped = line.strip()
        if not stripped:
            continue

        # Check if heading
        block_type = "text"
        if stripped.startswith('#'):
            block_type = "heading"

        if block_type == "heading":
            # If adding this line exceeds max size
            if current_chunk and len(current_chunk) >= min_chunk_size:
                # Save current chunk
                chunk_text = ""
                if current_headings:
                    chunk_text = " > ".join([h['text'] for h in current_headings]) + "\n\n"
                chunk_text += current_chunk
                chunks.append({
                    'text': chunk_text,
                    'headings': [h['text'] for h in current_headings],
                    'heading_path': " > ".join([h['text'] for h in current_headings]),
                    'type': block_type
                })
                current_chunk = ""

            # Update heading hierarchy - carry over headings from last block
            level = len(stripped) - len(stripped.lstrip('#'))
            heading_text = stripped.lstrip('#').strip()

            # Reset headings for next block
            current_headings = [h for h in current_headings if h.get('level', 0) < level]
            current_headings.append({'text': heading_text, 'level': level})

        else:
            # Regular content
            if len(current_chunk) + len(stripped) > max_chunk_size:
                if current_chunk:
                    chunk_text = ""
                    if current_headings:
                        chunk_text = " > ".join([h['text'] for h in current_headings]) + "\n\n"
                    chunk_text += current_chunk
                    chunks.append({
                        'text': chunk_text,
                        'headings': [h['text'] for h in current_headings],
                        'heading_path': " > ".join([h['text'] for h in current_headings]),
                        'type': 'text'
                    })
                    current_chunk = ""

            current_chunk += stripped + "\n"

    # Flush remaining content - carry over headings from last
    if current_chunk:
        chunk_text = ""
        if current_headings:
            chunk_text = " > ".join([h['text'] for h in current_headings]) + "\n\n"
        chunk_text += current_chunk
        chunks.append({
            'text': chunk_text,
            'headings': [h['text'] for h in current_headings],
            'heading_path': " > ".join([h['text'] for h in current_headings]),
            'type': 'text'
        })

    return chunks

# ---------------------------------------------------------------------------

# PDF PROCESSING HELPER FUNCTIONS

def process_pdf_with_ocr(pdf_content, filename, filepath, table_manager, encoded_metadata=None, parallel_pages=False, parallel_images=False):
    """
    Process PDF using OCR to extract text, images, and tables.

    Args:
        pdf_content: Raw PDF content as bytes
        filename: Name of the PDF file
        filepath: Full path for metadata
        table_manager: Block table manager for PDF file
        encoded_metadata: Optional metadata from parent document
        parallel_pages: Whether to process pages in parallel
        parallel_images: Whether to process images in parallel

    Returns:
        List of Document objects with extracted content
    """
    # List of Document objects with extracted content
    documents = []
    temp_file = None

    try:
        # Save to temporary file for processing
        temp_file = tempfile.NamedTemporaryFile(suffix='.pdf', delete=False)
        temp_file.write(pdf_content)
        temp_file.flush()
        temp_file_path = temp_file.name
        temp_file.close()

        # IMPORTANT: Check if file is non-empty
        if os.path.getsize(temp_file_path) == 0:
            log_and_print(f"Processing empty PDF: {filename}")
            return []

        # Try to open the document
        try:
            pdf_document = fitz.open(temp_file_path)
        except Exception as e:
            log_and_print(f"Error opening PDF {filename}: {str(e)}")
            return []

        total_pages = len(pdf_document)
        log_and_print(f"Processing PDF: {filename} with {total_pages} pages")

        # Try Azure Document Intelligence first for better extraction
        try:
            DI_ENDPOINT = os.getenv("AZURE_DI_ENDPOINT")
            DI_KEY = os.getenv("AZURE_DI_KEY")

            if DI_ENDPOINT and DI_KEY:
                document_analysis_client = DocumentAnalysisClient(
                    endpoint=DI_ENDPOINT,
                    credential=AzureKeyCredential(DI_KEY)
                )

                with open(temp_file_path, "rb") as f:
                    poller = document_analysis_client.begin_analyze_document(
                        "prebuilt-layout", f
                    )
                result = poller.result()

                # Process DI results
                for page in result.pages:
                    page_text = ""
                    for line in page.lines:
                        page_text += line.content + "\n"

                    if page_text.strip():
                        base_metadata = {
                            "source": filename,
                            "file_path": filepath,
                            "page": page.page_number,
                            "total_pages": total_pages,
                            "content_type": "pdf",
                            "extraction_method": "azure_document_intelligence"
                        }

                        from langchain.schema import Document as LCDocument
                        page_doc = LCDocument(
                            page_content=page_text.strip(),
                            metadata=base_metadata
                        )
                        documents.append(page_doc)

                log_and_print(f"DI extraction complete for {filename}: {len(documents)} pages extracted (total: {total_pages})")

        except Exception as e:
            log_and_print(f"DI extraction failed for {filename}: {str(e)}", "warning")
            log_and_print(f"Falling back to PyMuPDF extraction")

        # Fallback to PyMuPDF if DI failed or returned no results
        if not documents:
            log_and_print(f"Using PyMuPDF for {filename}")

            for page_num in range(total_pages):
                try:
                    page = pdf_document[page_num]

                    # Extract text content
                    page_text = page.get_text("text")

                    # Check for scanned/image pages
                    if not page_text or len(page_text.strip()) < 50:
                        # Likely a scanned page, use OCR
                        log_and_print(f"Scanned page detected: page {page_num + 1} of {filename}")

                        # Render page to image for OCR
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                        img_bytes = pix.tobytes("png")

                        page_text = get_gpt4o_description(img_bytes, detail="high")

                    # Extract images
                    image_list = page.get_images(full=True)

                    page_doc_text = page_text.strip() if page_text else ""

                    base_metadata = {
                        "source": filename,
                        "file_path": filepath,
                        "page": page_num + 1,
                        "total_pages": total_pages,
                        "content_type": "pdf",
                        "extraction_method": "pymupdf"
                    }

                    if encoded_metadata:
                        base_metadata.update(encoded_metadata)

                    from langchain.schema import Document as LCDocument
                    page_doc = LCDocument(
                        page_content=page_doc_text,
                        metadata=base_metadata
                    )
                    documents.append(page_doc)

                except Exception as e:
                    log_and_print(f"Error processing page {page_num + 1} of {filename}: {str(e)}")
                    continue

        pdf_document.close()

        # Assign chunk indices
        for i, doc in enumerate(documents):
            doc.metadata['chunk_index'] = i
            doc.metadata['total_chunks'] = len(documents)
            if encoded_metadata:
                doc.metadata['encoded_metadata'] = encoded_metadata
                doc.metadata['parent_document'] = encoded_metadata.get("title", "Unknown")

        text_content = " ".join([doc.page_content for doc in documents if doc.page_content])
        full_content = f"PDF document '{filename}' content:\n\n{text_content}"

        log_and_print(f"Process complete: Total {len(documents)} chunks for {filename}")

        return documents

    except Exception as e:
        log_and_print(f"Error processing PDF {filename}: {str(e)}", "error")
        return []
    finally:
        if temp_file and os.path.exists(temp_file.name):
            os.unlink(temp_file.name)

# ---------------------------------------------------------------------------

# DOCX PROCESSING

def process_docx_with_metadata(docx_content, filename, filepath, table_manager, encoded_metadata=None):
    """
    Process DOCX with heading extraction and metadata.

    Args:
        docx_content: Raw DOCX content as bytes
        filename: Name of the DOCX file
        filepath: Full file path
        table_manager: Table manager instance
        encoded_metadata: Optional parent metadata

    Returns:
        List of Document objects
    """
    documents = []

    try:
        docx_file = io.BytesIO(docx_content)
        doc = docx.Document(docx_file)

        current_heading = None
        current_section = []
        heading_path = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # Check if heading style
            if para.style and para.style.name and para.style.name.startswith('Heading'):
                # Save current section
                if current_section:
                    section_text = '\n'.join(current_section)

                    base_metadata = {
                        "source": filename,
                        "file_path": filepath,
                        "heading": current_heading,
                        "heading_path": " > ".join(heading_path),
                        "content_type": "docx",
                    }

                    if encoded_metadata:
                        base_metadata.update(encoded_metadata)

                    from langchain.schema import Document as LCDocument
                    section_doc = LCDocument(
                        page_content=section_text,
                        metadata=base_metadata
                    )
                    documents.append(section_doc)
                    current_section = []

                # Update heading hierarchy
                try:
                    level = int(para.style.name.replace('Heading ', ''))
                except:
                    level = 1

                heading_path = [h for h in heading_path if True]  # Keep all for now
                current_heading = text
                heading_path.append(text)

            else:
                current_section.append(text)

        # Don't forget last section
        if current_section:
            section_text = '\n'.join(current_section)
            base_metadata = {
                "source": filename,
                "file_path": filepath,
                "heading": current_heading,
                "heading_path": " > ".join(heading_path),
                "content_type": "docx",
            }

            if encoded_metadata:
                base_metadata.update(encoded_metadata)

            from langchain.schema import Document as LCDocument
            section_doc = LCDocument(
                page_content=section_text,
                metadata=base_metadata
            )
            documents.append(section_doc)

        # Also extract tables
        for table_idx, table in enumerate(doc.tables):
            table_text = ""
            for row in table.rows:
                row_cells = [cell.text.strip() for cell in row.cells]
                table_text += " | ".join(row_cells) + "\n"

            if table_text.strip():
                table_metadata = {
                    "source": filename,
                    "file_path": filepath,
                    "content_type": "docx_table",
                    "table_index": table_idx,
                }

                from langchain.schema import Document as LCDocument
                table_doc = LCDocument(
                    page_content=table_text,
                    metadata=table_metadata
                )
                documents.append(table_doc)

        log_and_print(f"DOCX processing complete for {filename}: {len(documents)} sections/tables")

    except Exception as e:
        log_and_print(f"Error processing DOCX {filename}: {str(e)}", "error")

    return documents

# ---------------------------------------------------------------------------

# PPTX PROCESSING

def process_pptx_with_metadata(pptx_content, filename, filepath, encoded_metadata=None):
    """Process PPTX with slide extraction and metadata"""
    documents = []

    try:
        from pptx import Presentation
        pptx_file = io.BytesIO(pptx_content)
        prs = Presentation(pptx_file)

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text.append(text)

                # Also extract table data
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells])
                        slide_text.append(row_text)

            if slide_text:
                base_metadata = {
                    "source": filename,
                    "file_path": filepath,
                    "slide_number": slide_num,
                    "total_slides": len(prs.slides),
                    "content_type": "pptx",
                }

                if encoded_metadata:
                    base_metadata.update(encoded_metadata)

                from langchain.schema import Document as LCDocument
                slide_doc = LCDocument(
                    page_content='\n'.join(slide_text),
                    metadata=base_metadata
                )
                documents.append(slide_doc)

        log_and_print(f"PPTX processing complete for {filename}: {len(documents)} slides")

    except Exception as e:
        log_and_print(f"Error processing PPTX {filename}: {str(e)}", "error")

    return documents

# ---------------------------------------------------------------------------

# XLSX / CSV PROCESSING

def process_xlsx_with_metadata(xlsx_content, filename, filepath, encoded_metadata=None):
    """Process Excel file with sheet extraction and metadata"""
    documents = []

    try:
        xlsx_file = io.BytesIO(xlsx_content)
        wb = openpyxl.load_workbook(xlsx_file, read_only=True)

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = f"Sheet: {sheet_name}\n"

            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                sheet_text += row_text + "\n"

            if sheet_text.strip():
                base_metadata = {
                    "source": filename,
                    "file_path": filepath,
                    "sheet_name": sheet_name,
                    "content_type": "xlsx",
                }

                if encoded_metadata:
                    base_metadata.update(encoded_metadata)

                from langchain.schema import Document as LCDocument
                sheet_doc = LCDocument(
                    page_content=sheet_text,
                    metadata=base_metadata
                )
                documents.append(sheet_doc)

        wb.close()
        log_and_print(f"XLSX processing complete for {filename}: {len(documents)} sheets")

    except Exception as e:
        log_and_print(f"Error processing XLSX {filename}: {str(e)}", "error")

    return documents

# ---------------------------------------------------------------------------

# MAIN DOCUMENT LOADING FUNCTION

def load_and_split_documents(files_to_process, token_manager, encoded_metadata=None, path_category="auto"):
    """
    Loads and splits documents from ADLS - reads a List of Document objects.
    Supports three-category routing and NOVA structural metadata enrichment.

    Args:
        files_to_process: List of tuples (filepath, filename, file_content)
        token_manager: Blob token manager for API calls
        encoded_metadata: Optional parent metadata extracted for JSON metadata content
        path_category: 'regulatory_json', 'internal_raw', or 'auto'

    Returns:
        docs_all: List of processed Document objects with NOVA metadata
    """
    docs_all = []
    docs_failed = 0

    for filepath, filename, file_content in files_to_process:
        try:
            file_extension = os.path.splitext(filename)[1].lower()

            log_and_print(f"Processing: {filename} ({file_extension}) [category={path_category}]")

            # --- NOVA: Load pre-extracted metadata from silver/metadata/ ---
            # metadata_extraction.py writes these BEFORE ingestion runs.
            # For regulatory docs, the scraped JSON already has all metadata.
            # For internal docs, this is the primary source of business context
            # (business_owner, confidentiality, audience, etc.)
            pre_extracted = {}
            if path_category != "regulatory_json":
                try:
                    pre_extracted = load_pre_extracted_metadata(filepath, container_client, container_name)
                except Exception as e:
                    log_and_print(f"Could not load pre-extracted metadata for {filename}: {e}", "warning")

            # --- Route regulatory scraped JSON directly ---
            if path_category == "regulatory_json" and file_extension == ".json":
                docs = process_regulatory_scraped_json(file_content, filename, filepath)
                if docs:
                    # Enrich each doc with structural metadata
                    for doc in docs:
                        enrich_chunk_with_structural_metadata(doc)
                    docs_all.extend(docs)
                    log_and_print(f"Regulatory JSON: {len(docs)} sections from {filename}")
                else:
                    docs_failed += 1
                continue

            # Check supported extensions
            if file_extension not in ['.pdf', '.docx', '.pptx', '.xlsx', '.xls', '.html', '.htm', '.md', '.txt', '.csv', '.json',
                                       '.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif']:
                log_and_print(f"Skipped directory or blob without extension: {filename}", "warning")
                continue

            # PDF files
            if file_extension == '.pdf':
                docs = process_pdf_with_ocr(file_content, filename, filepath, token_manager, encoded_metadata, parallel_pages=False, parallel_images=False)

            # DOCX files
            elif file_extension == '.docx':
                docs = process_docx_with_metadata(file_content, filename, filepath, token_manager, encoded_metadata)

            # PPTX files
            elif file_extension == '.pptx':
                docs = process_pptx_with_metadata(file_content, filename, filepath, encoded_metadata)

            # HTML files
            elif file_extension in ['.html', '.htm']:
                html_text = file_content.decode('utf-8', errors='ignore') if isinstance(file_content, bytes) else file_content
                structured = process_html_with_ocr(html_text, filepath, None, None, encoded_metadata)

                docs = []
                for item in structured:
                    from langchain.schema import Document as LCDocument
                    doc = LCDocument(
                        page_content=item.get('content', ''),
                        metadata={
                            "source": filename,
                            "file_path": filepath,
                            "heading": item.get('heading', ''),
                            "content_type": f"html_{item.get('type', 'text')}",
                        }
                    )
                    docs.append(doc)

            # Excel files
            elif file_extension in ['.xlsx', '.xls']:
                docs = process_xlsx_with_metadata(file_content, filename, filepath, encoded_metadata)

            # Text/Markdown/CSV/JSON files
            elif file_extension in ['.txt', '.md', '.csv', '.json']:
                text = file_content.decode('utf-8', errors='ignore') if isinstance(file_content, bytes) else file_content

                from langchain.schema import Document as LCDocument
                docs = [LCDocument(
                    page_content=text,
                    metadata={
                        "source": filename,
                        "file_path": filepath,
                        "content_type": file_extension.lstrip('.'),
                    }
                )]

            # Image files
            elif file_extension in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif']:
                text = get_gpt4o_description(file_content, detail="high")
                from langchain.schema import Document as LCDocument
                docs = [LCDocument(
                    page_content=text if text else "",
                    metadata={
                        "source": filename,
                        "file_path": filepath,
                        "content_type": "image_ocr",
                    }
                )] if text else []

            else:
                docs = []

            if docs:
                # --- Enrich every doc with structural + pre-extracted metadata ---
                for doc in docs:
                    enrich_chunk_with_structural_metadata(doc, pre_extracted_metadata=pre_extracted)
                    # Set source_type if not already set
                    if "source_type" not in doc.metadata:
                        doc.metadata["source_type"] = "internal" if path_category == "internal_raw" else "auto"

                docs_all.extend(docs)
                log_and_print(f"Loaded {len(docs)} documents from {filename}")
            else:
                docs_failed += 1
                log_and_print(f"No documents extracted from {filename}", "warning")

        except Exception as e:
            docs_failed += 1
            log_and_print(f"Error loading {filename}: {str(e)}", "error")
            continue

    log_and_print(f"Total documents loaded: {len(docs_all)}, Failed: {docs_failed}")
    return docs_all

# ---------------------------------------------------------------------------

# ELASTICSEARCH INDEX AND BATCH OPERATIONS

def create_es_index_with_mapping(es_client, index_name):
    """Create Elasticsearch index with NOVA metadata mapping (Rule 2).

    All NOVA fields are stored as filterable/boostable index fields.
    The content_vector is the dense_vector for cosine similarity search.
    bm25_text contains the semantic header + chunk text for BM25 keyword search.
    """

    mapping = {
        "mappings": {
            "properties": {
                # --- Core chunk fields ---
                "doc_id": {"type": "keyword"},
                "chunk_id": {"type": "keyword"},
                "chunk_index": {"type": "integer"},
                "total_chunks": {"type": "integer"},
                "text": {"type": "text"},
                "content": {"type": "text"},
                "bm25_text": {"type": "text"},  # semantic header + text for BM25
                "file_name": {"type": "keyword"},
                "file_path": {"type": "keyword"},
                "file_extension": {"type": "keyword"},
                "file_size": {"type": "long"},
                "content_type": {"type": "keyword"},
                "source_type": {"type": "keyword"},  # 'regulatory' or 'internal'

                # --- Heading / section structure ---
                "heading": {"type": "text"},
                "heading_path": {"type": "keyword"},
                "heading_text": {"type": "text"},
                "section_path": {"type": "text", "fields": {"keyword": {"type": "keyword"}}},
                "citation_anchor": {"type": "keyword"},
                "source": {"type": "keyword"},

                # --- NOVA document metadata (Rule 2 index fields) ---
                "title": {"type": "text", "fields": {"keyword": {"type": "keyword"}}},
                "short_title": {"type": "keyword"},
                "document_class": {"type": "keyword"},
                "status": {"type": "keyword"},
                "effective_date_start": {"type": "date", "format": "yyyy-MM-dd||yyyy-MM||yyyy||epoch_millis", "ignore_malformed": True},
                "effective_date_end": {"type": "date", "format": "yyyy-MM-dd||yyyy-MM||yyyy||epoch_millis", "ignore_malformed": True},
                "jurisdiction": {"type": "keyword"},
                "authority_class": {"type": "keyword"},
                "authority_level": {"type": "integer"},
                "nova_tier": {"type": "keyword"},
                "confidentiality": {"type": "keyword"},
                "approval_status": {"type": "keyword"},

                # --- Regulatory-specific ---
                "regulator": {"type": "keyword"},
                "regulator_acronym": {"type": "keyword"},
                "guideline_number": {"type": "keyword"},
                "doc_family_id": {"type": "keyword"},
                "version_id": {"type": "keyword"},
                "version_label": {"type": "keyword"},
                "current_version_flag": {"type": "boolean"},
                "sector": {"type": "keyword"},
                "supersedes_doc_id": {"type": "keyword"},
                "superseded_by_doc_id": {"type": "keyword"},

                # --- Internal-specific ---
                "business_owner": {"type": "keyword"},
                "document_owner": {"type": "keyword"},
                "business_line": {"type": "keyword"},
                "audience": {"type": "keyword"},

                # --- Structural metadata (unit-level) ---
                "structural_level": {"type": "keyword"},
                "section_number": {"type": "keyword"},
                "depth": {"type": "integer"},
                "parent_section_id": {"type": "keyword"},
                "is_appendix": {"type": "boolean"},
                "normative_weight": {"type": "keyword"},
                "paragraph_role": {"type": "keyword"},
                "cross_references": {"type": "keyword"},  # array of strings

                # --- Content boolean flags ---
                "contains_definition": {"type": "boolean"},
                "contains_formula": {"type": "boolean"},
                "contains_requirement": {"type": "boolean"},
                "contains_deadline": {"type": "boolean"},
                "contains_assignment": {"type": "boolean"},
                "contains_parameter": {"type": "boolean"},

                # --- Operational ---
                "last_modified": {"type": "date"},
                "creation_time": {"type": "date"},
                "total_pages": {"type": "integer"},
                "page_number": {"type": "integer"},
                "custom_metadata": {"type": "object"},
                "ingestion_timestamp": {"type": "date"},
                "raw_path": {"type": "keyword"},
                "canonical_path": {"type": "keyword"},
                "sha256": {"type": "keyword"},
                "parser_version": {"type": "keyword"},
                "quality_score": {"type": "float"},

                # --- Embedding vector ---
                "content_vector": {
                    "type": "dense_vector",
                    "dims": EMBEDDING_DIMS,
                    "index": True,
                    "similarity": "cosine"
                }
            }
        },
        "settings": {
            "number_of_shards": 1,
            "number_of_replicas": 1
        }
    }

    if es_client.indices.exists(index=index_name):
        log_and_print(f"Index {index_name} already exists")
        return

    es_client.indices.create(index=index_name, body=mapping)
    log_and_print(f"Created NOVA index: {index_name} with {len(mapping['mappings']['properties'])} fields")


# ---------------------------------------------------------------------------

# ============================================================================
# PGVector DUAL-STORE SUPPORT (alongside Elasticsearch)
# ============================================================================

# PGVector connection configuration
PG_HOST = os.getenv("PG_HOST", "")
PG_PORT = os.getenv("PG_PORT", "5432")
PG_DATABASE = os.getenv("PG_DATABASE", "nova")
PG_USER = os.getenv("PG_USER", "")
PG_PASSWORD = os.getenv("PG_PASSWORD", "")
PG_TABLE = os.getenv("PG_TABLE", "nova_chunks")

_pg_conn = None  # cached connection


def get_pg_conn():
    """Create or return cached PGVector connection."""
    global _pg_conn
    if _pg_conn is not None:
        try:
            # Test if connection is alive
            _pg_conn.cursor().execute("SELECT 1")
            return _pg_conn
        except Exception:
            _pg_conn = None

    if not PG_HOST:
        return None  # PGVector not configured — skip silently

    try:
        import psycopg2
        from psycopg2.extras import Json

        _pg_conn = psycopg2.connect(
            host=PG_HOST, port=PG_PORT,
            dbname=PG_DATABASE, user=PG_USER, password=PG_PASSWORD,
        )
        _pg_conn.autocommit = True
        log_and_print(f"Connected to PGVector at {PG_HOST}:{PG_PORT}/{PG_DATABASE}")
        return _pg_conn
    except ImportError:
        log_and_print("psycopg2 not installed — PGVector dual-store disabled", "warning")
        return None
    except Exception as e:
        log_and_print(f"PGVector connection failed: {e}", "warning")
        return None


def create_pgvector_table(table_name=None):
    """Create PGVector table with NOVA metadata columns if it doesn't exist.

    Schema mirrors the Elasticsearch mapping so both stores have identical
    filterable fields. The embedding column uses pgvector's vector type
    for cosine similarity search.
    """
    conn = get_pg_conn()
    if not conn:
        return

    table_name = table_name or PG_TABLE

    create_sql = f"""
    CREATE EXTENSION IF NOT EXISTS vector;

    CREATE TABLE IF NOT EXISTS {table_name} (
        -- Core chunk identity
        chunk_id        TEXT PRIMARY KEY,
        doc_id          TEXT,
        chunk_index     INTEGER,
        content         TEXT,
        bm25_text       TEXT,

        -- Embedding vector
        embedding       vector({EMBEDDING_DIMS}),

        -- Document metadata (Rule 2 index fields)
        title           TEXT,
        short_title     TEXT,
        source_type     TEXT,
        document_class  TEXT,
        file_path       TEXT,
        source          TEXT,
        heading_path    TEXT,
        section_path    TEXT,
        status          TEXT DEFAULT 'active',
        effective_date_start TEXT,
        effective_date_end   TEXT,
        jurisdiction    TEXT,
        authority_class TEXT,
        authority_level INTEGER,
        nova_tier       TEXT,

        -- Regulatory
        regulator       TEXT,
        regulator_acronym TEXT,
        guideline_number TEXT,
        version_id      TEXT,
        version_label   TEXT,
        current_version_flag BOOLEAN DEFAULT TRUE,
        sector          TEXT,
        doc_family_id   TEXT,

        -- Internal
        business_owner  TEXT,
        business_line   TEXT,
        audience        TEXT,
        approval_status TEXT,
        confidentiality TEXT,

        -- Structural metadata
        structural_level TEXT,
        section_number  TEXT,
        depth           INTEGER,
        normative_weight TEXT,
        paragraph_role  TEXT,
        is_appendix     BOOLEAN DEFAULT FALSE,
        cross_references JSONB DEFAULT '[]',

        -- Content flags
        contains_definition  BOOLEAN DEFAULT FALSE,
        contains_formula     BOOLEAN DEFAULT FALSE,
        contains_requirement BOOLEAN DEFAULT FALSE,
        contains_deadline    BOOLEAN DEFAULT FALSE,
        contains_assignment  BOOLEAN DEFAULT FALSE,
        contains_parameter   BOOLEAN DEFAULT FALSE,

        -- Operational
        ingestion_timestamp  TIMESTAMPTZ DEFAULT NOW(),
        page_number     INTEGER,
        quality_score   REAL
    );

    -- Create HNSW index for fast cosine similarity search
    CREATE INDEX IF NOT EXISTS idx_{table_name}_embedding
        ON {table_name} USING hnsw (embedding vector_cosine_ops);

    -- Create indexes for common filter fields
    CREATE INDEX IF NOT EXISTS idx_{table_name}_doc_id ON {table_name} (doc_id);
    CREATE INDEX IF NOT EXISTS idx_{table_name}_regulator ON {table_name} (regulator);
    CREATE INDEX IF NOT EXISTS idx_{table_name}_status ON {table_name} (status);
    CREATE INDEX IF NOT EXISTS idx_{table_name}_source_type ON {table_name} (source_type);
    CREATE INDEX IF NOT EXISTS idx_{table_name}_normative ON {table_name} (normative_weight);
    CREATE INDEX IF NOT EXISTS idx_{table_name}_jurisdiction ON {table_name} (jurisdiction);
    """

    try:
        with conn.cursor() as cur:
            cur.execute(create_sql)
        log_and_print(f"PGVector table '{table_name}' created/verified with NOVA schema")
    except Exception as e:
        log_and_print(f"Error creating PGVector table: {e}", "error")


def upsert_chunks_to_pgvector(chunk_dicts, embeddings, table_name=None):
    """Upsert chunk records with embeddings to PGVector.

    Called after Elasticsearch upsert to maintain dual-store consistency.
    Each chunk_dict is the same dict prepared for ES in index_chunks_and_document().

    Args:
        chunk_dicts: list of dicts (same format as ES _source)
        embeddings: list of embedding vectors (parallel to chunk_dicts)
        table_name: PGVector table name (defaults to PG_TABLE)
    """
    conn = get_pg_conn()
    if not conn:
        return 0  # PGVector not configured — skip silently

    table_name = table_name or PG_TABLE

    upsert_sql = f"""
    INSERT INTO {table_name} (
        chunk_id, doc_id, chunk_index, content, bm25_text, embedding,
        title, short_title, source_type, document_class, file_path, source,
        heading_path, section_path, status, effective_date_start, effective_date_end,
        jurisdiction, authority_class, authority_level, nova_tier,
        regulator, regulator_acronym, guideline_number, version_id, version_label,
        current_version_flag, sector, doc_family_id,
        business_owner, business_line, audience, approval_status, confidentiality,
        structural_level, section_number, depth, normative_weight, paragraph_role,
        is_appendix, cross_references,
        contains_definition, contains_formula, contains_requirement,
        contains_deadline, contains_assignment, contains_parameter,
        page_number, quality_score
    ) VALUES (
        %(chunk_id)s, %(doc_id)s, %(chunk_index)s, %(content)s, %(bm25_text)s, %(embedding)s,
        %(title)s, %(short_title)s, %(source_type)s, %(document_class)s, %(file_path)s, %(source)s,
        %(heading_path)s, %(section_path)s, %(status)s, %(effective_date_start)s, %(effective_date_end)s,
        %(jurisdiction)s, %(authority_class)s, %(authority_level)s, %(nova_tier)s,
        %(regulator)s, %(regulator_acronym)s, %(guideline_number)s, %(version_id)s, %(version_label)s,
        %(current_version_flag)s, %(sector)s, %(doc_family_id)s,
        %(business_owner)s, %(business_line)s, %(audience)s, %(approval_status)s, %(confidentiality)s,
        %(structural_level)s, %(section_number)s, %(depth)s, %(normative_weight)s, %(paragraph_role)s,
        %(is_appendix)s, %(cross_references)s,
        %(contains_definition)s, %(contains_formula)s, %(contains_requirement)s,
        %(contains_deadline)s, %(contains_assignment)s, %(contains_parameter)s,
        %(page_number)s, %(quality_score)s
    )
    ON CONFLICT (chunk_id) DO UPDATE SET
        content = EXCLUDED.content,
        bm25_text = EXCLUDED.bm25_text,
        embedding = EXCLUDED.embedding,
        ingestion_timestamp = NOW()
    """

    try:
        import psycopg2.extras
        success = 0
        with conn.cursor() as cur:
            for chunk_dict, embedding in zip(chunk_dicts, embeddings):
                try:
                    params = {
                        "chunk_id": chunk_dict.get("chunk_id", hashlib.md5(
                            f"{chunk_dict.get('file_path','')}{chunk_dict.get('chunk_index','')}".encode()
                        ).hexdigest()),
                        "doc_id": chunk_dict.get("doc_id", ""),
                        "chunk_index": chunk_dict.get("chunk_index", 0),
                        "content": chunk_dict.get("content", ""),
                        "bm25_text": chunk_dict.get("bm25_text", ""),
                        "embedding": str(embedding) if embedding else None,
                        "title": chunk_dict.get("title", ""),
                        "short_title": chunk_dict.get("short_title", ""),
                        "source_type": chunk_dict.get("source_type", ""),
                        "document_class": chunk_dict.get("document_class", ""),
                        "file_path": chunk_dict.get("file_path", ""),
                        "source": chunk_dict.get("source", ""),
                        "heading_path": chunk_dict.get("heading_path", ""),
                        "section_path": chunk_dict.get("section_path", ""),
                        "status": chunk_dict.get("status", "active"),
                        "effective_date_start": chunk_dict.get("effective_date_start") or None,
                        "effective_date_end": chunk_dict.get("effective_date_end") or None,
                        "jurisdiction": chunk_dict.get("jurisdiction", ""),
                        "authority_class": chunk_dict.get("authority_class", ""),
                        "authority_level": chunk_dict.get("authority_level", 0),
                        "nova_tier": chunk_dict.get("nova_tier", ""),
                        "regulator": chunk_dict.get("regulator", ""),
                        "regulator_acronym": chunk_dict.get("regulator_acronym", ""),
                        "guideline_number": chunk_dict.get("guideline_number", ""),
                        "version_id": chunk_dict.get("version_id", ""),
                        "version_label": chunk_dict.get("version_label", ""),
                        "current_version_flag": chunk_dict.get("current_version_flag", True),
                        "sector": chunk_dict.get("sector", ""),
                        "doc_family_id": chunk_dict.get("doc_family_id", ""),
                        "business_owner": chunk_dict.get("business_owner", ""),
                        "business_line": chunk_dict.get("business_line", ""),
                        "audience": chunk_dict.get("audience", ""),
                        "approval_status": chunk_dict.get("approval_status", ""),
                        "confidentiality": chunk_dict.get("confidentiality", ""),
                        "structural_level": chunk_dict.get("structural_level", ""),
                        "section_number": chunk_dict.get("section_number", ""),
                        "depth": chunk_dict.get("depth", 0),
                        "normative_weight": chunk_dict.get("normative_weight", ""),
                        "paragraph_role": chunk_dict.get("paragraph_role", ""),
                        "is_appendix": chunk_dict.get("is_appendix", False),
                        "cross_references": psycopg2.extras.Json(chunk_dict.get("cross_references", [])),
                        "contains_definition": chunk_dict.get("contains_definition", False),
                        "contains_formula": chunk_dict.get("contains_formula", False),
                        "contains_requirement": chunk_dict.get("contains_requirement", False),
                        "contains_deadline": chunk_dict.get("contains_deadline", False),
                        "contains_assignment": chunk_dict.get("contains_assignment", False),
                        "contains_parameter": chunk_dict.get("contains_parameter", False),
                        "page_number": chunk_dict.get("page_number", 0),
                        "quality_score": chunk_dict.get("quality_score", 0.0),
                    }
                    cur.execute(upsert_sql, params)
                    success += 1
                except Exception as e:
                    log_and_print(f"PGVector upsert error for chunk: {e}", "warning")

        log_and_print(f"PGVector: upserted {success}/{len(chunk_dicts)} chunks")
        return success

    except Exception as e:
        log_and_print(f"PGVector batch upsert failed: {e}", "error")
        return 0


# ---------------------------------------------------------------------------

def batch_upsert_to_es(es_client, index_name, documents):
    """Batch upsert documents to Elasticsearch"""

    BATCH_SIZE = 50

    total_processed = 0

    for i in range(0, len(documents), BATCH_SIZE):
        batch = documents[i:i + BATCH_SIZE]

        actions = []
        for doc in batch:
            # Create unique ID
            doc_id = hashlib.md5(
                f"{doc.get('file_path', '')}{doc.get('chunk_index', '')}".encode()
            ).hexdigest()

            action = {
                "_index": index_name,
                "_id": doc_id,
                "_source": doc
            }
            actions.append(action)

        try:
            helpers.bulk(es_client, actions, raise_on_error=False)
            total_processed += len(batch)
            logger.info(f"Processed {total_processed}/{len(documents)} documents")
        except Exception as e:
            logger.error(f"Error in batch upsert: {e}")

    return total_processed

# ---------------------------------------------------------------------------

def clean_batch(batch_docs):
    """Clean batch of documents before upsert"""
    cleaned = []
    for doc in batch_docs:
        if doc and hasattr(doc, 'page_content') and doc.page_content:
            cleaned.append(doc)
    return cleaned

# ---------------------------------------------------------------------------

def index_chunks_and_document(es_client, index_name, batch_docs, total_batched, retry_count=3):
    """Index batch of document chunks to Elasticsearch with retry.

    Prepends NOVA semantic header (Rule 1) to chunk text before embedding.
    Stores all NOVA metadata fields as index fields (Rule 2).
    """

    actions = []
    for doc in batch_docs:
        meta = doc.metadata

        # --- NOVA RULE 1: Build semantic header and prepend to text ---
        semantic_header = build_semantic_header(meta, meta.get("source_type", "auto"))
        original_text = doc.page_content
        embed_text = f"{semantic_header}\n{original_text}" if semantic_header else original_text

        # Build doc dict with ALL NOVA fields for ES index (Rule 2)
        doc_dict = {
            # Core text fields
            "content": original_text,
            "bm25_text": embed_text,  # semantic header + text for BM25 search
            "text": original_text,

            # Document identity
            "doc_id": meta.get("doc_id", ""),
            "source": meta.get("source", ""),
            "file_path": meta.get("file_path", ""),
            "content_type": meta.get("content_type", ""),
            "source_type": meta.get("source_type", ""),

            # Heading/section structure
            "heading": meta.get("heading", ""),
            "heading_path": meta.get("heading_path", ""),
            "heading_text": meta.get("heading", ""),
            "section_path": meta.get("section_path", ""),

            # NOVA document metadata
            "title": meta.get("title", ""),
            "short_title": meta.get("short_title", ""),
            "document_class": meta.get("document_class", ""),
            "status": meta.get("status", ""),
            "effective_date_start": meta.get("effective_date_start", None),
            "effective_date_end": meta.get("effective_date_end", None),
            "jurisdiction": meta.get("jurisdiction", ""),
            "authority_class": meta.get("authority_class", ""),
            "nova_tier": meta.get("nova_tier", ""),

            # Regulatory-specific
            "regulator": meta.get("regulator", ""),
            "regulator_acronym": meta.get("regulator_acronym", ""),
            "guideline_number": meta.get("guideline_number", ""),
            "version_id": meta.get("version_id", ""),
            "version_label": meta.get("version_label", ""),
            "current_version_flag": meta.get("current_version_flag", None),
            "sector": meta.get("sector", ""),
            "doc_family_id": meta.get("doc_family_id", ""),

            # Internal-specific
            "business_owner": meta.get("business_owner", ""),
            "business_line": meta.get("business_line", ""),
            "audience": meta.get("audience", ""),
            "approval_status": meta.get("approval_status", ""),
            "confidentiality": meta.get("confidentiality", ""),

            # Structural metadata
            "structural_level": meta.get("structural_level", ""),
            "section_number": meta.get("section_number", ""),
            "depth": meta.get("depth", 0),
            "is_appendix": meta.get("is_appendix", False),
            "normative_weight": meta.get("normative_weight", ""),
            "paragraph_role": meta.get("paragraph_role", ""),
            "cross_references": meta.get("cross_references", []),

            # Content boolean flags
            "contains_definition": meta.get("contains_definition", False),
            "contains_formula": meta.get("contains_formula", False),
            "contains_requirement": meta.get("contains_requirement", False),
            "contains_deadline": meta.get("contains_deadline", False),
            "contains_assignment": meta.get("contains_assignment", False),
            "contains_parameter": meta.get("contains_parameter", False),

            # Operational
            "ingestion_timestamp": datetime.datetime.utcnow().isoformat(),
            "page_number": meta.get("page", meta.get("page_number", 0)),
            "metadata": meta,  # full metadata dict for backward compat
        }

        # Remove None values to avoid ES mapping errors
        doc_dict = {k: v for k, v in doc_dict.items() if v is not None}

        doc_id = hashlib.md5(
            f"{meta.get('file_path', '')}{meta.get('chunk_index', '')}".encode()
        ).hexdigest()

        action = {
            "_index": index_name,
            "_id": doc_id,
            "_source": doc_dict
        }
        actions.append(action)

    # Retry logic
    for attempt in range(retry_count):
        try:
            success, errors = helpers.bulk(es_client, actions, raise_on_error=False)
            log_and_print(f"Successfully added {success} documents to vector store!")

            if errors:
                log_and_print(f"Batch had {len(errors)} errors", "warning")

            return success

        except Exception as e:
            log_and_print(f"Bulk index attempt {attempt + 1} failed: {e}", "error")
            if attempt < retry_count - 1:
                sample_file = batch_docs[0].metadata.get("source", "unknown") if batch_docs else "unknown"
                log_and_print(f"Batch containing {sample_file} failed. Retrying in 30 seconds...")
                time.sleep(30)

    return 0

# ---------------------------------------------------------------------------

def create_es_batch_after(doc_batch_extract_structure, es_client, index_name):
    """Helper to add batch after doc_batch_extract structure"""

    total_added = 0

    for batch_item in doc_batch_extract_structure:
        record_type = batch_item.metadata.get("content_type", "other")
        try:
            doc_id = hashlib.md5(
                f"{batch_item.metadata.get('file_path', '')}{batch_item.metadata.get('chunk_index', '')}".encode()
            ).hexdigest()

            action = {
                "_index": index_name,
                "_id": doc_id,
                "_source": {
                    "content": batch_item.page_content,
                    "metadata": batch_item.metadata,
                }
            }

            es_client.index(index=index_name, id=doc_id, document=action["_source"])
            total_added += 1
            log_and_print(f"Added document {doc_id} successfully")

        except Exception as e:
            log_and_print(f"Error adding document: {str(e)}", "error")

    # Don't write tokens to records
    return total_added

# ---------------------------------------------------------------------------

def get_indexed_files(es_client, index_name):
    """Get list of already indexed files from Elasticsearch"""
    try:
        result = es_client.search(
            index=index_name,
            body={
                "size": 0,
                "aggs": {
                    "file_paths": {
                        "terms": {
                            "field": "file_path.keyword",
                            "size": 10000
                        }
                    }
                }
            }
        )

        file_paths = [bucket["key"] for bucket in result["aggregations"]["file_paths"]["buckets"]]
        return set(file_paths)
    except Exception as e:
        log_and_print(f"Error getting indexed files: {e}")
        return set()

# ---------------------------------------------------------------------------

# MAIN EXECUTION PIPELINE

# Step 1: Find all files in ADLS and create indexing setup

index_name = "document-embeddings-metadata"

# Try to find existing tracking file
ADLS_TOTAL_METADATA_PATH = f"metadata/tracking_files.json"

try:
    all_files_metadata = {}
    log_and_print("No existing tracking file found. Starting fresh.")
except:
    all_files_metadata = {}

# ---------------------------------------------------------------------------

# Step 2: Indexing files to ADLS

# Find files in ADLS container
container_name = os.getenv("CONTAINER_NAME", "documents")

try:
    container_client = blob_service_client.get_container_client(container_name)

    all_files = []
    for blob in container_client.list_blobs():
        if not blob.name.endswith("/"):
            all_files.append(blob.name)

    log_and_print(f"Total files found in ADLS: {len(all_files)}")

except Exception as e:
    log_and_print(f"Error listing files in ADLS: {e}", "error")
    all_files = []

# ---------------------------------------------------------------------------

# Step 3: Connecting to Elasticsearch

try:
    info = elastic_client.info()
    log_and_print(f"Successfully connected to Elasticsearch!")
except Exception as e:
    log_and_print(f"Connection attempt to Elasticsearch failed. Retrying on retry_default seconds...", "warning")
    log_and_print(f"Failure connecting to Elasticsearch after max_retries attempts: {e}", "error")
    raise

# Step 3.5: Destroying Elasticsearch index
# create_es_index_with_mapping(elastic_client, index_name)

# ---------------------------------------------------------------------------

# Step 4: NOVA three-category discovery and file filtering

es_indexed_files = get_indexed_files(elastic_client, index_name)
log_and_print(f"Files already indexed: {len(es_indexed_files)}")

# Use NOVA three-category discovery
discovered_paths = discover_paths_to_ingest(container_client, container_name)
log_and_print(f"Discovered {len(discovered_paths)} files in ADLS")

# Separate by category and filter already-indexed
regulatory_json_files = []
internal_raw_files = []
auto_files = []
files_already_indexed = []

for blob_path, path_category in discovered_paths:
    if blob_path in es_indexed_files:
        files_already_indexed.append(blob_path)
        continue

    filename = os.path.basename(blob_path)

    try:
        blob_client = container_client.get_blob_client(blob_path)
        file_bytes = blob_client.download_blob().readall()

        file_tuple = (blob_path, filename, file_bytes)
        if path_category == "regulatory_json":
            regulatory_json_files.append(file_tuple)
        elif path_category == "internal_raw":
            internal_raw_files.append(file_tuple)
        else:
            auto_files.append(file_tuple)
    except Exception as e:
        log_and_print(f"Error downloading {blob_path}: {e}", "error")

# Backward compat: merge all for original pipeline flow
files_to_process = regulatory_json_files + internal_raw_files + auto_files

log_and_print(f"Files to process: {len(files_to_process)} "
              f"(regulatory_json={len(regulatory_json_files)}, "
              f"internal_raw={len(internal_raw_files)}, "
              f"auto={len(auto_files)})")
log_and_print(f"Files already indexed: {len(files_already_indexed)}")

if len(files_to_process) == 0:
    log_and_print("No new files to process. All files are already indexed.")

# ---------------------------------------------------------------------------

# Step 5: Initialize blob token manager

log_and_print("Step 5 is setting up blob token manager")

try:
    token_manager = blob_service_client
    log_and_print("Blob token manager initialized")
except Exception as e:
    log_and_print(f"Error initializing blob token manager: {e}", "error")
    raise

# ---------------------------------------------------------------------------

# Step 6: Load and split files by category (NOVA three-category processing)

log_and_print(f"Step 6 is loading and splitting {len(files_to_process)} files with OCR and metadata")

docs = []

# Process regulatory JSON files (pre-scraped metadata, route to process_regulatory_scraped_json)
if regulatory_json_files:
    log_and_print(f"Processing {len(regulatory_json_files)} regulatory JSON files...")
    reg_docs = load_and_split_documents(regulatory_json_files, token_manager, all_files_metadata, path_category="regulatory_json")
    docs.extend(reg_docs)
    log_and_print(f"Regulatory JSON: {len(reg_docs)} documents")

# Process internal raw files (DOCX, PDF, HTML, etc. with OCR)
if internal_raw_files:
    log_and_print(f"Processing {len(internal_raw_files)} internal raw files...")
    int_docs = load_and_split_documents(internal_raw_files, token_manager, all_files_metadata, path_category="internal_raw")
    docs.extend(int_docs)
    log_and_print(f"Internal raw: {len(int_docs)} documents")

# Process auto-detected files
if auto_files:
    log_and_print(f"Processing {len(auto_files)} auto-detected files...")
    auto_docs = load_and_split_documents(auto_files, token_manager, all_files_metadata, path_category="auto")
    docs.extend(auto_docs)
    log_and_print(f"Auto: {len(auto_docs)} documents")

log_and_print(f"Done loading and splitting documents: {len(docs)} total")

# ---------------------------------------------------------------------------

# Step 7: Set up fresh token and create embedding model

from openai import AzureOpenAI

embedding_client = AzureOpenAI(
    api_key=os.getenv("OPENAI_API_KEY"),
    api_version=os.getenv("OPENAI_API_VERSION"),
    azure_endpoint=os.getenv("OPENAI_API_BASE")
)

embedding_model = os.getenv("OPENAI_EMBEDDING_DEPLOYMENT", "text-embedding-3-large")

# Create vector store
from langchain_community.vectorstores import ElasticsearchStore

vector_store = ElasticsearchStore(
    es_cloud_id=os.getenv("ES_CLOUD_ID"),
    index_name=index_name,
    es_api_key=es_api_key,
)

log_and_print("Embedding model and vector store initialized successfully!")

# ---------------------------------------------------------------------------

# Step 8: Adding CsvIndex3 document chunks to index in batches

# Maximum process after initialization

filtered_docs = [x for x in docs if x is not None and hasattr(x, 'page_content') and x.page_content]

log_and_print(f"Filtered docs: {len(filtered_docs)} (removed {len(docs) - len(filtered_docs)} empty/None docs)")

# ---------------------------------------------------------------------------

# Batch indexing

batch_size = 50

for i in range(0, len(filtered_docs), batch_size):
    batch = filtered_docs[i:i + batch_size]

    # Clean batch
    batch = clean_batch(batch)
    if not batch:
        continue

    log_and_print(f"Indexing {len(batch)} documents to vector store (batch {i // batch_size + 1}/{len(filtered_docs) // batch_size + 1})")

    try:
        # Add to vector store
        vector_store.add_documents(batch)
        log_and_print(f"Successfully added {len(batch)} documents to vector store!")

    except Exception as e:
        log_and_print(f"Error adding batch: {e}", "error")

        # Retry
        sample_file = batch[0].metadata.get("source", "unknown") if batch else "unknown"
        log_and_print(f"Batch containing {sample_file} failed. Retrying in 30 seconds...")
        time.sleep(30)

        try:
            vector_store.add_documents(batch)
            log_and_print(f"Retry successful for batch containing {sample_file}")
        except Exception as e2:
            log_and_print(f"Failed to add batch after retry: {str(e2)}", "error")

            # Try individual documents
            for doc in batch:
                try:
                    vector_store.add_documents([doc])
                except Exception as e3:
                    log_and_print(f"Failed individual doc: {doc.metadata.get('source', 'unknown')}: {str(e3)}", "error")

# ---------------------------------------------------------------------------

# Final summary

log_and_print(f"Final files in ADLS: {len(all_files)}")
log_and_print(f"New files processed: {len(files_to_process)}")
log_and_print(f"Total documents/chunks indexed: {len(filtered_docs)}")
log_and_print(f"Files skipped (already indexed): {len(files_already_indexed)}")

# ---------------------------------------------------------------------------

# --- Vector Store Creation and Document Ingestion Complete ---

# ---------------------------------------------------------------------------

# Status cleanup and monitoring the system

log_and_print("Master Status:")
log_and_print(f"  Files in ADLS: {len(all_files)}")
log_and_print(f"  New files processed: {len(files_to_process)}")
log_and_print(f"  Docs created: {len(filtered_docs)}")
log_and_print(f"  ES Index: {index_name}")

reason = "Master data creation and document ingestion complete."

log_and_print(f"Processing complete!")

# ---------------------------------------------------------------------------

# Azure Cognitive Search setup (alternative to Elasticsearch)

def create_azure_search_index(index_name):
    """Create Azure Cognitive Search index with vector search configuration"""
    try:
        from azure.search.documents.indexes import SearchIndexClient
        from azure.search.documents.indexes.models import (
            SearchIndex,
            SearchableField,
            SimpleField,
            SearchFieldDataType,
            VectorSearch,
            HnswAlgorithmConfiguration,
            VectorSearchProfile,
            SearchField,
        )

        search_endpoint = os.getenv("AZURE_SEARCH_ENDPOINT")
        search_key = os.getenv("AZURE_SEARCH_KEY")

        if not search_endpoint or not search_key:
            log_and_print("Azure Search credentials not configured", "warning")
            return

        index_client = SearchIndexClient(
            endpoint=search_endpoint,
            credential=AzureKeyCredential(search_key)
        )

        fields = [
            SimpleField(name="id", type=SearchFieldDataType.String, key=True),
            SearchableField(name="content", type=SearchFieldDataType.String),
            SimpleField(name="source", type=SearchFieldDataType.String, filterable=True),
            SimpleField(name="file_path", type=SearchFieldDataType.String, filterable=True),
            SimpleField(name="content_type", type=SearchFieldDataType.String, filterable=True),
            SearchableField(name="heading", type=SearchFieldDataType.String),
            SearchableField(name="heading_path", type=SearchFieldDataType.String),
            SearchField(
                name="content_vector",
                type=SearchFieldDataType.Collection(SearchFieldDataType.Single),
                searchable=True,
                vector_search_dimensions=1536,
                vector_search_profile_name="myHnswProfile",
            ),
        ]

        vector_search = VectorSearch(
            algorithms=[
                HnswAlgorithmConfiguration(name="myHnsw"),
            ],
            profiles=[
                VectorSearchProfile(
                    name="myHnswProfile",
                    algorithm_configuration_name="myHnsw",
                ),
            ],
        )

        index = SearchIndex(
            name=index_name,
            fields=fields,
            vector_search=vector_search,
        )

        index_client.create_or_update_index(index)
        log_and_print(f"Azure Search index '{index_name}' created successfully!")

    except Exception as e:
        log_and_print(f"Error creating Azure Search index: {str(e)}", "error")

# ---------------------------------------------------------------------------

log_and_print(f"Processing complete!")
