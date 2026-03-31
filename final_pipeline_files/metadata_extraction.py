# Databricks notebook source
# ---------------------------------------------------------------------------
# NOVA RAG: Metadata Extraction Layer
# **Purpose**: Extract and normalize metadata for ALL document types, writing
# standardised per-document JSON to `silver/metadata/`.
# **Two extraction paths**:
# | Source | Input location | How metadata is obtained |
# |--------|----------------|--------------------------|
# | **External regulatory** (OSFI, PRA, BCBS, SEC …) | `bronze/external/<regulator>/json/` | **Read from scraped JSON** — the web scrapers already populate every NOVA field. This script normalizes and writes to `silver/metadata/`. |
# | **Internal documents** (policy, memos, research) | `bronze/internal/` | **Extract from raw files** — parse native metadata (docProps, EXIF, frontmatter, meta tags, etc.), merge with enrichment registry, resolve via three-tier logic. |
# **Output**: `silver/metadata/{doc_id}.json` + optional `silver/metadata/_catalog.json`
# **Field alignment**: All resolved metadata is aligned with `metadata_spec.py`
# (COMMON_FIELDS, REGULATORY_FIELDS, INTERNAL_FIELDS, UNIT_STRUCTURAL_FIELDS).
# **Three-tier resolution** (internal docs only):
# `enrichment_registry > native_metadata > heuristic_defaults`
# **ADLS layout consumed/written**:
# ```
# nova-docs/
# ├── bronze/
# │   ├── external/
# │   │   ├── osfi/json/*.json            ← scraped JSON with full metadata
# │   │   ├── pra/<category>/json/*.json   ← scraped JSON with full metadata
# │   │   └── <regulator>/json/*.json      ← scraped JSON with full metadata
# │   └── internal/                        ← raw DOCX, XLSX, PDF, PPTX, MD, images
# └── silver/
# └── metadata/
# ├── {doc_id}.json                ← written by THIS notebook
# └── _catalog.json                ← optional consolidated catalog
# ```

# ---------------------------------------------------------------------------
# azure-identity azure-storage-file-datalake -q

# ---------------------------------------------------------------------------
import json
import os
import re
import hashlib
import zipfile
import xml.etree.ElementTree as ET
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, Optional, List, Tuple
import base64
import io
import csv

import logging
import pandas as pd
from PIL import Image
from PIL.ExifTags import TAGS
import yaml
from bs4 import BeautifulSoup
from azure.identity import DefaultAzureCredential
from azure.storage.filedatalake import DataLakeServiceClient

# ---------------------------------------------------------------------------
## Logging

# ---------------------------------------------------------------------------
def setup_logging():
    _logger = logging.getLogger("nova_metadata_extraction")
    if not _logger.handlers:
        handler = logging.StreamHandler()
        handler.setFormatter(logging.Formatter("%(asctime)s %(levelname)s %(message)s"))
        _logger.addHandler(handler)
    _logger.setLevel(logging.INFO)
    return _logger

logger = setup_logging()

def log_and_print(msg: str, level: str = "info") -> None:
    """Print and log a message."""
    print(msg)
    getattr(logger, level, logger.info)(msg)

# ---------------------------------------------------------------------------
## Configuration & Widgets

# ---------------------------------------------------------------------------
try:
    dbutils.widgets.text("adls_account_url", "", "ADLS account URL")
    dbutils.widgets.text("adls_file_system", "nova-docs", "ADLS file system name")
    dbutils.widgets.text("bronze_external_prefix", "bronze/external/", "External regulatory prefix")
    dbutils.widgets.text("bronze_internal_prefix", "bronze/internal/", "Internal docs prefix")
    dbutils.widgets.text("silver_metadata_prefix", "silver/metadata/", "Output metadata prefix")
    dbutils.widgets.text("enrichment_registry_path", "silver/config/enrichment_registry.json", "Enrichment registry")
    dbutils.widgets.dropdown("mode", "extract", ["extract", "test"], "Mode")
    dbutils.widgets.dropdown("write_catalog", "true", ["true", "false"], "Write consolidated catalog?")
except Exception:
    pass


def widget(name: str, default: str = "") -> str:
    try:
        value = dbutils.widgets.get(name)
        return value if value is not None else default
    except Exception:
        return default


ADLS_ACCOUNT_URL = widget("adls_account_url")
ADLS_FILE_SYSTEM = widget("adls_file_system", "nova-docs")
bronze_external_prefix = widget("bronze_external_prefix", "bronze/external/")
bronze_internal_prefix = widget("bronze_internal_prefix", "bronze/internal/")
silver_metadata_prefix = widget("silver_metadata_prefix", "silver/metadata/")
enrichment_registry_path = widget("enrichment_registry_path", "silver/config/enrichment_registry.json")
mode = widget("mode", "extract")
write_catalog = widget("write_catalog", "true").lower() == "true"

# ---------------------------------------------------------------------------
## ADLS I/O Helpers

# ---------------------------------------------------------------------------
def get_credential() -> DefaultAzureCredential:
    return DefaultAzureCredential()


def get_adls_client() -> DataLakeServiceClient:
    return DataLakeServiceClient(account_url=ADLS_ACCOUNT_URL, credential=get_credential())


def get_fs_client():
    return get_adls_client().get_file_system_client(ADLS_FILE_SYSTEM)


def read_adls_bytes(path: str) -> bytes:
    return get_fs_client().get_file_client(path).download_file().readall()


def write_adls_json(path: str, payload: dict) -> None:
    data = json.dumps(payload, ensure_ascii=False, indent=2, default=str).encode("utf-8")
    get_fs_client().get_file_client(path).upload_data(data, overwrite=True)

# ---------------------------------------------------------------------------
## Constants: NOVA metadata spec field names
# These are the canonical field names from `metadata_spec.py`.
# Every metadata JSON we produce must resolve these fields.

# ---------------------------------------------------------------------------
# Document-level fields from the spec — common + regulatory + internal
# The extraction script resolves ALL of these for every document.
# These field lists mirror metadata_spec.py — keep them in sync.
NOVA_COMMON_FIELDS = [
    "doc_id", "source_type", "title", "short_title", "document_class",
    "heading_path", "section_path", "citation_anchor",
    "raw_path", "canonical_json_path", "raw_sha256",
    "parser_version", "quality_score",
]

NOVA_REGULATORY_FIELDS = [
    "regulator", "regulator_acronym", "doc_family_id",
    "version_id", "version_label", "version_sort_key",
    "guideline_number", "status", "current_version_flag",
    "effective_date_start", "effective_date_end",
    "authority_class", "authority_level", "nova_tier",
    "jurisdiction", "sector",
    "supersedes_doc_id", "superseded_by_doc_id",
    "contains_definition", "contains_formula", "contains_requirement",
    "contains_deadline", "contains_assignment", "contains_parameter",
]

NOVA_INTERNAL_FIELDS = [
    "doc_family_id", "version_id", "version_label", "current_version_flag",
    "business_owner", "document_owner", "approval_status", "approval_date",
    "effective_date_start", "effective_date_end",
    "review_date", "next_review_date",
    "confidentiality", "business_line", "function",
    "jurisdiction", "audience", "status",
    "contains_deadline", "contains_assignment",
    "contains_definition", "contains_formula", "contains_requirement",
    "contains_parameter",
]

# All document-level fields (de-duplicated)
ALL_DOC_FIELDS = list(dict.fromkeys(
    NOVA_COMMON_FIELDS + NOVA_REGULATORY_FIELDS + NOVA_INTERNAL_FIELDS
))

SUPPORTED_EXTENSIONS = {
    ".docx", ".pptx", ".xlsx", ".xls",
    ".csv", ".tsv", ".pdf",
    ".html", ".htm", ".md", ".json", ".txt",
    ".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".gif",
}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".gif"}

# ---------------------------------------------------------------------------
## Enrichment Registry

# ---------------------------------------------------------------------------
def load_enrichment_registry() -> Dict[str, Dict]:
    """Load enrichment registry from ADLS. Returns filename → overrides dict."""
    try:
        raw = read_adls_bytes(enrichment_registry_path)
        return json.loads(raw)
    except Exception:
        log_and_print("Enrichment registry not found — using empty registry")
        return {}

ENRICHMENT_REGISTRY = load_enrichment_registry()

# ---------------------------------------------------------------------------
# ============================================================================
# DIRECTORY-PATH HEURISTICS FOR DOCUMENT CLASSIFICATION
# ============================================================================

def infer_document_class_from_path(filepath: str) -> str:
    """Infer document_class from ADLS directory path.

    Uses directory naming conventions:
        bronze/internal/policy/        -> 'policy'
        bronze/internal/research/      -> 'research_paper'
        bronze/internal/presentations/ -> 'presentation'
        bronze/internal/memos/         -> 'memo'
    """
    fp_lower = filepath.lower().replace("\\", "/")
    path_class_map = {
        "/policy/": "policy", "/policies/": "policy",
        "/procedure/": "procedure", "/procedures/": "procedure",
        "/guideline/": "guideline", "/guidelines/": "guideline",
        "/research/": "research_paper",
        "/presentation/": "presentation", "/presentations/": "presentation",
        "/template/": "template", "/templates/": "template",
        "/memo/": "memo", "/memos/": "memo",
        "/report/": "report", "/reports/": "report",
        "/data/": "structured_data",
        "/reference/": "reference_document",
        "/training/": "training_material",
        "/syllabus/": "syllabus",
    }
    for path_pattern, doc_class in path_class_map.items():
        if path_pattern in fp_lower:
            return doc_class

    ext = os.path.splitext(filepath)[1].lower()
    ext_map = {".pptx": "presentation", ".xlsx": "spreadsheet", ".csv": "structured_data"}
    return ext_map.get(ext, "")


# ---------------------------------------------------------------------------
# ============================================================================
# LLM-ASSISTED METADATA EXTRACTION
# ============================================================================
# For internal documents where native metadata is sparse (PDFs, scanned docs),
# use the LLM to infer business-context fields from the document's first page.

LLM_EXTRACTION_MODEL = os.getenv("LLM_EXTRACTION_MODEL", "gpt-5-mini-2025-08-07-eastus-dz")
LLM_EXTRACTION_ENABLED = os.getenv("LLM_EXTRACTION_ENABLED", "false").lower() == "true"


def extract_metadata_with_llm(first_page_text: str, filename: str) -> Dict[str, Any]:
    """Use the LLM to infer metadata fields from the document's first page.

    Called when native metadata extraction yields thin results and the enrichment
    registry has no entry for this file. The LLM reads the first ~2000 chars and
    infers: title, document_class, business_owner, audience, jurisdiction,
    confidentiality, and a short summary.

    Args:
        first_page_text: First ~2000 characters of the document
        filename: Original filename for context

    Returns:
        dict: inferred NOVA metadata fields (may be partial)
    """
    if not LLM_EXTRACTION_ENABLED:
        return {}

    if not first_page_text or len(first_page_text.strip()) < 50:
        return {}

    try:
        from openai import AzureOpenAI

        client = AzureOpenAI(
            api_key=os.getenv("OPENAI_API_KEY"),
            api_version=os.getenv("OPENAI_API_VERSION", "2024-02-01"),
            azure_endpoint=os.getenv("OPENAI_API_BASE"),
        )

        # Truncate to first 2000 chars to keep costs low
        text_sample = first_page_text[:2000]

        prompt = f"""You are a metadata extraction assistant for a document management system.
Analyze the following text (first page of a document named "{filename}") and extract metadata.

Return ONLY a JSON object with these fields (use null for fields you cannot determine):
{{
  "title": "full document title",
  "short_title": "compact title (max 50 chars)",
  "document_class": "one of: policy, procedure, guideline, memo, report, research_paper, presentation, syllabus, template, reference_document, structured_data",
  "business_owner": "department or team that owns this document",
  "audience": "who this document is intended for",
  "jurisdiction": "geographic/regulatory scope (e.g., Canada, UK, Global)",
  "confidentiality": "one of: public, internal, internal_confidential, restricted",
  "subject_summary": "one-sentence summary of what this document covers"
}}

TEXT:
{text_sample}

Return ONLY the JSON object, no explanation."""

        response = client.chat.completions.create(
            model=LLM_EXTRACTION_MODEL,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=500,
            temperature=0.0,
        )

        result_text = response.choices[0].message.content.strip()

        # Parse JSON from response (handle markdown code blocks)
        if result_text.startswith("```"):
            result_text = result_text.split("```")[1]
            if result_text.startswith("json"):
                result_text = result_text[4:]

        inferred = json.loads(result_text)

        # Filter out null values and the summary (not a NOVA field)
        cleaned = {}
        for key, value in inferred.items():
            if value is not None and key != "subject_summary":
                cleaned[key] = value

        log_and_print(f"LLM inferred {len(cleaned)} metadata fields for {filename}")
        return cleaned

    except Exception as e:
        log_and_print(f"LLM metadata extraction failed for {filename}: {e}", level="warning")
        return {}


# ---------------------------------------------------------------------------
# ============================================================================
# METADATA COMPLETENESS SCORE
# ============================================================================

def compute_metadata_completeness(resolved: Dict[str, Any], doc_type: str) -> float:
    """Compute what percentage of required NOVA fields are populated.

    Returns a score from 0.0 (no fields populated) to 1.0 (all fields populated).
    This score is stored as quality_score in the metadata JSON and can be used
    to flag documents needing enrichment attention.

    Args:
        resolved: dict of resolved NOVA metadata fields
        doc_type: 'regulatory' or 'internal'

    Returns:
        float: completeness score between 0.0 and 1.0
    """
    if doc_type == "regulatory":
        # Critical fields for regulatory docs
        required_fields = [
            "doc_id", "title", "regulator", "guideline_number",
            "status", "effective_date_start", "jurisdiction",
            "authority_class", "nova_tier", "document_class",
            "version_id", "sector",
        ]
        nice_to_have = [
            "short_title", "version_label", "effective_date_end",
            "doc_family_id", "supersedes_doc_id", "citation_anchor",
        ]
    else:
        # Critical fields for internal docs
        required_fields = [
            "doc_id", "title", "document_class", "business_owner",
            "status", "confidentiality", "audience",
        ]
        nice_to_have = [
            "short_title", "business_line", "jurisdiction",
            "effective_date_start", "approval_status", "version_id",
            "document_owner", "review_date",
        ]

    # Score: required fields count 2x, nice-to-have count 1x
    total_weight = len(required_fields) * 2 + len(nice_to_have)
    earned = 0

    for field in required_fields:
        value = resolved.get(field)
        if value not in (None, "", [], 0):
            earned += 2

    for field in nice_to_have:
        value = resolved.get(field)
        if value not in (None, "", [], 0):
            earned += 1

    score = earned / total_weight if total_weight > 0 else 0.0
    return round(score, 3)


# ---------------------------------------------------------------------------
## Three-tier resolution

# ---------------------------------------------------------------------------
def resolve_field(field_name: str,
                  enrichment: Dict[str, Any],
                  native: Dict[str, Any],
                  default: Any = None) -> Any:
    """Three-tier resolution: enrichment_registry > native_metadata > heuristic_defaults."""
    if enrichment.get(field_name) not in (None, "", []):
        return enrichment[field_name]
    if native.get(field_name) not in (None, "", []):
        return native[field_name]
    return default


def resolve_all_nova_fields(
    native: Dict[str, Any],
    enrichment: Dict[str, Any],
    doc_type: str,
) -> Dict[str, Any]:
    """Resolve every NOVA spec field for a document.

    For regulatory docs, resolves COMMON + REGULATORY fields.
    For internal docs, resolves COMMON + INTERNAL fields.
    """
    if doc_type == "regulatory":
        target_fields = list(dict.fromkeys(NOVA_COMMON_FIELDS + NOVA_REGULATORY_FIELDS))
    else:
        target_fields = list(dict.fromkeys(NOVA_COMMON_FIELDS + NOVA_INTERNAL_FIELDS))

    resolved = {}
    for fname in target_fields:
        resolved[fname] = resolve_field(fname, enrichment, native)

    # Heuristic defaults for critical fields when still None
    if not resolved.get("doc_id"):
        source_file = native.get("source_file", "unknown")
        resolved["doc_id"] = _generate_doc_id(source_file, doc_type)
    if not resolved.get("title"):
        resolved["title"] = native.get("source_file", "Untitled")
    if not resolved.get("quality_score"):
        resolved["quality_score"] = 0.5
    if not resolved.get("parser_version"):
        resolved["parser_version"] = "nova-metadata-extractor-v1.0.0"
    # source_type is always set from doc_type
    if not resolved.get("source_type"):
        resolved["source_type"] = doc_type
    # Default status for internal docs
    if doc_type == "internal" and not resolved.get("status"):
        resolved["status"] = resolved.get("approval_status", "active")

    return resolved


def _generate_doc_id(source_file: str, doc_type: str) -> str:
    """Generate a stable doc_id from filename."""
    base = os.path.basename(source_file)
    name = os.path.splitext(base)[0]
    slug = re.sub(r"[^a-z0-9]+", "_", name.lower()).strip("_")
    prefix = "internal" if doc_type == "internal" else "ext"
    return f"{prefix}.{slug}"

# ---------------------------------------------------------------------------
## ══════════════════════════════════════════════════════════════
## PATH A: External Regulatory Docs (scraped JSON already has metadata)
## ══════════════════════════════════════════════════════════════
# External regulatory documents (OSFI, PRA, BCBS, SEC, etc.) have already
# been scraped by purpose-built web scrapers. The scraped JSON files contain
# every NOVA metadata field at the top level. This path simply:
# 1. Reads each scraped JSON
# 2. Maps its fields to the standard NOVA metadata envelope
# 3. Computes structural summary from sections if available
# 4. Writes to `silver/metadata/{doc_id}.json`

# ---------------------------------------------------------------------------
def extract_metadata_from_scraped_json(scraped: Dict[str, Any], json_path: str) -> Dict[str, Any]:
    """Build a NOVA metadata document from a pre-scraped regulatory JSON.

    The scraped JSON has fields like doc_id, title, authority_class, nova_tier,
    jurisdiction, sector, status, effective_date_start, etc. directly at the
    top level. We read them as-is (they are the 'native' metadata for regulatory
    docs) and apply enrichment overrides if any exist.
    """
    source_file = os.path.basename(json_path)
    enrichment = ENRICHMENT_REGISTRY.get(source_file, {})
    doc_id_from_json = scraped.get("doc_id", source_file)
    enrichment_by_id = ENRICHMENT_REGISTRY.get(doc_id_from_json, {})
    enrichment = {**enrichment_by_id, **enrichment}  # filename overrides take precedence

    # The scraped JSON IS the native metadata — read all NOVA fields from it
    native = {}
    for field_name in ALL_DOC_FIELDS:
        if field_name in scraped:
            native[field_name] = scraped[field_name]

    # Also capture scraper-specific fields that map to NOVA fields
    _SCRAPER_ALIASES = {
        "raw_html_sha256": "raw_sha256",
        "raw_pdf_sha256": "raw_sha256",
        "crawl_timestamp": "download_timestamp",
        "scraped_at": "download_timestamp",
        "source_url": "raw_path",
        "raw_html_path": "raw_path",
    }
    for scraper_key, nova_key in _SCRAPER_ALIASES.items():
        if scraper_key in scraped and not native.get(nova_key):
            native[nova_key] = scraped[scraper_key]

    native["source_file"] = json_path

    # Resolve (enrichment > native > defaults)
    resolved = resolve_all_nova_fields(native, enrichment, doc_type="regulatory")

    # Structural summary from sections if present
    structural_summary = _compute_structural_summary_from_scraped(scraped)

    # Metadata completeness score
    completeness = compute_metadata_completeness(resolved, "regulatory")
    resolved["quality_score"] = completeness

    # Compute SHA-256 of the scraped JSON itself
    sha256 = scraped.get("sha256") or scraped.get("raw_html_sha256") or scraped.get("raw_sha256")
    if not sha256:
        sha256 = hashlib.sha256(json.dumps(scraped, sort_keys=True).encode()).hexdigest()

    return {
        "doc_id": resolved["doc_id"],
        "source_file": json_path,
        "file_type": "json",
        "doc_type": "regulatory",
        "extraction_timestamp": datetime.utcnow().isoformat() + "Z",
        "extraction_method": "scraped_json",
        "sha256": sha256,
        "metadata_completeness": completeness,
        "resolved_metadata": resolved,
        "structural_summary": structural_summary,
        "scraper_fields_preserved": {
            k: scraped[k] for k in [
                "source_system", "library_row_url", "listing_publication_type",
                "listing_category", "publication_type", "category",
                "pra_metadata", "boe_pdf_urls", "pra_parent_documents",
                "bm25_text", "vector_text_prefix",
            ] if k in scraped
        },
    }


def _compute_structural_summary_from_scraped(scraped: Dict) -> Dict[str, Any]:
    """Compute structural summary from a scraped JSON's sections list."""
    sections = scraped.get("sections", [])
    section_headings = scraped.get("section_headings", [])
    toc_depth = scraped.get("toc_depth", 0)
    has_appendices = scraped.get("has_appendices", False) or scraped.get("has_appendix", False)
    has_tables = scraped.get("has_tables", False)
    table_count = scraped.get("table_count", 0)
    has_footnotes = scraped.get("has_footnotes", False)
    footnote_count = scraped.get("footnote_count", 0)

    # Count normative weights and paragraph roles if sections have them
    normative_counts = {"mandatory": 0, "advisory": 0, "permissive": 0, "informational": 0}
    role_counts: Dict[str, int] = {}

    for sec in sections:
        if isinstance(sec, dict):
            nw = sec.get("normative_weight")
            if nw and nw in normative_counts:
                normative_counts[nw] += 1
            pr = sec.get("paragraph_role")
            if pr:
                role_counts[pr] = role_counts.get(pr, 0) + 1

    # If no per-section normative data, use document-level boolean flags
    if sum(normative_counts.values()) == 0:
        if scraped.get("contains_requirement"):
            normative_counts["mandatory"] = 1
        if scraped.get("contains_definition"):
            role_counts["definition"] = role_counts.get("definition", 0) + 1

    return {
        "total_sections": len(sections),
        "total_headings": len(section_headings),
        "toc_depth": toc_depth,
        "has_appendices": has_appendices,
        "has_tables": has_tables,
        "table_count": table_count,
        "has_footnotes": has_footnotes,
        "footnote_count": footnote_count,
        "normative_weights": normative_counts,
        "paragraph_roles": role_counts,
    }

# ---------------------------------------------------------------------------
## ══════════════════════════════════════════════════════════════
## PATH B: Internal Documents (extract metadata from raw files)
## ══════════════════════════════════════════════════════════════
# Internal documents arrive as raw files (DOCX, PPTX, XLSX, PDF, etc.)
# with NO pre-existing metadata JSON. This path:
# 1. Opens the raw file
# 2. Extracts native metadata (docProps, EXIF, frontmatter, etc.)
# 3. Merges with enrichment registry
# 4. Resolves all NOVA fields via three-tier logic
# 5. Writes to `silver/metadata/{doc_id}.json`

# ---------------------------------------------------------------------------
# --- Office Open XML native metadata extraction (DOCX, PPTX, XLSX) ---

def _extract_office_core_xml(zip_bytes: bytes) -> Dict[str, Any]:
    """Extract Dublin Core metadata from docProps/core.xml in an Office Open XML file."""
    meta = {}
    try:
        with zipfile.ZipFile(io.BytesIO(zip_bytes)) as zf:
            if "docProps/core.xml" in zf.namelist():
                core_xml = zf.read("docProps/core.xml")
                root = ET.fromstring(core_xml)
                ns = {
                    "dc": "http://purl.org/dc/elements/1.1/",
                    "dcterms": "http://purl.org/dc/terms/",
                    "cp": "http://schemas.openxmlformats.org/package/2006/metadata/core-properties",
                }
                _MAP = {
                    "dc:title": "title",
                    "dc:creator": "creator",
                    "dc:subject": "subject",
                    "dc:description": "description",
                    "cp:keywords": "keywords",
                    "cp:category": "category",
                    "cp:lastModifiedBy": "last_modified_by",
                    "dcterms:created": "created",
                    "dcterms:modified": "modified",
                    "cp:revision": "revision",
                }
                for xpath, key in _MAP.items():
                    el = root.find(xpath, ns)
                    if el is not None and el.text:
                        meta[key] = el.text.strip()

            # Custom properties (org-specific metadata)
            if "docProps/custom.xml" in zf.namelist():
                custom_xml = zf.read("docProps/custom.xml")
                croot = ET.fromstring(custom_xml)
                for prop in croot:
                    name = prop.attrib.get("name", "")
                    for child in prop:
                        if child.text:
                            meta[f"custom.{name}"] = child.text.strip()
    except Exception:
        pass
    return meta


def _extract_office_app_xml(zip_bytes: bytes) -> Dict[str, Any]:
    """Extract application metadata from docProps/app.xml."""
    meta = {}
    try:
        with zipfile.ZipFile(io.BytesIO(zip_bytes)) as zf:
            if "docProps/app.xml" in zf.namelist():
                app_xml = zf.read("docProps/app.xml")
                root = ET.fromstring(app_xml)
                ns = {"ep": "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"}
                for tag in ["Application", "Company", "Manager", "Template", "Pages", "Words", "Slides"]:
                    el = root.find(f"ep:{tag}", ns)
                    if el is not None and el.text:
                        meta[f"app.{tag.lower()}"] = el.text.strip()
    except Exception:
        pass
    return meta


# --- Format-specific extractors ---

def extract_docx_metadata(path: str, raw_bytes: bytes) -> Dict[str, Any]:
    """Extract metadata from a DOCX file via Office Open XML."""
    meta = {"extraction_method": "native_office_xml", "source_file": path}
    core = _extract_office_core_xml(raw_bytes)
    app = _extract_office_app_xml(raw_bytes)
    meta.update(core)
    meta.update(app)

    # Count headings and paragraphs for structural summary
    try:
        from docx import Document
        doc = Document(io.BytesIO(raw_bytes))
        heading_count = 0
        para_count = 0
        has_appendix = False
        for p in doc.paragraphs:
            if p.style and p.style.name and p.style.name.startswith("Heading"):
                heading_count += 1
                if "appendix" in p.text.lower():
                    has_appendix = True
            if p.text.strip():
                para_count += 1
        meta["heading_count"] = heading_count
        meta["paragraph_count"] = para_count
        meta["has_appendix"] = has_appendix
        meta["word_count"] = sum(len(p.text.split()) for p in doc.paragraphs)
    except Exception:
        pass
    return meta


def extract_pptx_metadata(path: str, raw_bytes: bytes) -> Dict[str, Any]:
    """Extract metadata from a PPTX file."""
    meta = {"extraction_method": "native_office_xml", "source_file": path}
    core = _extract_office_core_xml(raw_bytes)
    app = _extract_office_app_xml(raw_bytes)
    meta.update(core)
    meta.update(app)

    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(raw_bytes))
        meta["slide_count"] = len(prs.slides)
        titles = []
        for slide in prs.slides:
            if slide.shapes.title and slide.shapes.title.text:
                titles.append(slide.shapes.title.text)
        meta["slide_titles"] = titles
    except Exception:
        pass
    return meta


def extract_xlsx_metadata(path: str, raw_bytes: bytes) -> Dict[str, Any]:
    """Extract metadata from an XLSX file."""
    meta = {"extraction_method": "native_office_xml", "source_file": path}
    core = _extract_office_core_xml(raw_bytes)
    app = _extract_office_app_xml(raw_bytes)
    meta.update(core)
    meta.update(app)

    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(raw_bytes), read_only=True, data_only=True)
        meta["sheet_names"] = wb.sheetnames
        meta["sheet_count"] = len(wb.sheetnames)
        ws = wb[wb.sheetnames[0]]
        meta["first_sheet_rows"] = ws.max_row
        meta["first_sheet_cols"] = ws.max_column
        wb.close()
    except Exception:
        pass
    return meta


def extract_xls_metadata(path: str, raw_bytes: bytes) -> Dict[str, Any]:
    """Extract metadata from a legacy .xls file."""
    meta = {"extraction_method": "pandas_xls", "source_file": path}
    try:
        xls = pd.ExcelFile(io.BytesIO(raw_bytes), engine="xlrd")
        meta["sheet_names"] = xls.sheet_names
        meta["sheet_count"] = len(xls.sheet_names)
    except Exception:
        meta["extraction_method"] = "enrichment_only"
    return meta


def extract_pdf_metadata(path: str, raw_bytes: bytes) -> Dict[str, Any]:
    """Extract metadata from a PDF file (/Info dictionary)."""
    meta = {"extraction_method": "pdf_info", "source_file": path}
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=raw_bytes, filetype="pdf")
        info = doc.metadata or {}
        for k, v in info.items():
            if v:
                meta[k.lower()] = str(v)
        meta["page_count"] = doc.page_count
        text_pages = sum(1 for page in doc if page.get_text().strip())
        meta["text_pages"] = text_pages
        meta["is_scanned"] = text_pages < doc.page_count * 0.5
        meta["has_images"] = any(page.get_images() for page in doc)
        doc.close()
    except Exception:
        meta["extraction_method"] = "enrichment_only"
    return meta


def extract_html_metadata(path: str, raw_bytes: bytes) -> Dict[str, Any]:
    """Extract metadata from HTML meta tags."""
    meta = {"extraction_method": "html_meta_tags", "source_file": path}
    try:
        text = raw_bytes.decode("utf-8", errors="replace")
        soup = BeautifulSoup(text, "html.parser")
        title_tag = soup.find("title")
        if title_tag and title_tag.string:
            meta["title"] = title_tag.string.strip()
        for tag in soup.find_all("meta"):
            name = tag.get("name", "") or tag.get("property", "")
            content = tag.get("content", "")
            if name and content:
                meta[f"meta.{name.lower()}"] = content
        heading_count = len(soup.find_all(re.compile(r"^h[1-6]$")))
        meta["heading_count"] = heading_count
        para_count = len(soup.find_all("p"))
        meta["paragraph_count"] = para_count
        meta["has_tables"] = bool(soup.find("table"))
        body_text = soup.get_text()
        meta["word_count"] = len(body_text.split())
    except Exception:
        meta["extraction_method"] = "enrichment_only"
    return meta


def extract_markdown_metadata(path: str, raw_bytes: bytes) -> Dict[str, Any]:
    """Extract metadata from Markdown YAML frontmatter + heading structure."""
    meta = {"extraction_method": "frontmatter", "source_file": path}
    try:
        text = raw_bytes.decode("utf-8", errors="replace")
        if text.startswith("---"):
            parts = text.split("---", 2)
            if len(parts) >= 3:
                try:
                    fm = yaml.safe_load(parts[1])
                    if isinstance(fm, dict):
                        for k, v in fm.items():
                            meta[k] = v
                except Exception:
                    pass

        lines = text.splitlines()
        headings = [l for l in lines if l.startswith("#")]
        meta["heading_count"] = len(headings)
        meta["word_count"] = len(text.split())
        meta["has_appendix"] = any("appendix" in h.lower() for h in headings)
    except Exception:
        meta["extraction_method"] = "enrichment_only"
    return meta


def extract_json_metadata(path: str, raw_bytes: bytes) -> Dict[str, Any]:
    """Extract metadata from a generic JSON file (non-regulatory)."""
    meta = {"extraction_method": "json_keys", "source_file": path}
    try:
        data = json.loads(raw_bytes)
        if isinstance(data, dict):
            for field_name in ALL_DOC_FIELDS:
                if field_name in data:
                    meta[field_name] = data[field_name]
            meta["top_level_keys"] = list(data.keys())[:50]
    except Exception:
        meta["extraction_method"] = "enrichment_only"
    return meta


def extract_csv_metadata(path: str, raw_bytes: bytes) -> Dict[str, Any]:
    """Extract metadata from a CSV/TSV file."""
    meta = {"extraction_method": "csv_header", "source_file": path}
    try:
        text = raw_bytes.decode("utf-8", errors="replace")
        dialect = csv.Sniffer().sniff(text[:4096])
        reader = csv.reader(io.StringIO(text), dialect)
        headers = next(reader, [])
        meta["column_headers"] = headers
        meta["column_count"] = len(headers)
        row_count = sum(1 for _ in reader)
        meta["row_count"] = row_count
    except Exception:
        meta["extraction_method"] = "enrichment_only"
    return meta


def extract_txt_metadata(path: str, raw_bytes: bytes) -> Dict[str, Any]:
    """Extract metadata from a plain text file."""
    meta = {"extraction_method": "text_heuristic", "source_file": path}
    try:
        text = raw_bytes.decode("utf-8", errors="replace")
        meta["word_count"] = len(text.split())
        meta["line_count"] = len(text.splitlines())
        meta["char_count"] = len(text)
        for line in text.splitlines():
            if line.strip():
                meta["title"] = line.strip()[:200]
                break
    except Exception:
        meta["extraction_method"] = "enrichment_only"
    return meta


def extract_image_metadata(path: str, raw_bytes: bytes) -> Dict[str, Any]:
    """Extract metadata from an image file (EXIF data)."""
    meta = {"extraction_method": "image_exif", "source_file": path}
    try:
        img = Image.open(io.BytesIO(raw_bytes))
        meta["image_width"] = img.width
        meta["image_height"] = img.height
        meta["image_format"] = img.format
        meta["image_mode"] = img.mode

        exif_data = img.getexif()
        if exif_data:
            exif_meta = {}
            for tag_id, value in exif_data.items():
                tag = TAGS.get(tag_id, str(tag_id))
                try:
                    exif_meta[tag] = str(value)
                except Exception:
                    pass
            if exif_meta:
                meta["exif"] = exif_meta
    except Exception:
        meta["extraction_method"] = "enrichment_only"
    return meta

# ---------------------------------------------------------------------------
## Format dispatcher for internal docs

# ---------------------------------------------------------------------------
def extract_native_metadata(file_path: str, raw_bytes: bytes) -> Dict[str, Any]:
    """Dispatch to the correct format-specific extractor based on file extension."""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".docx":
        return extract_docx_metadata(file_path, raw_bytes)
    elif ext == ".pptx":
        return extract_pptx_metadata(file_path, raw_bytes)
    elif ext == ".xlsx":
        return extract_xlsx_metadata(file_path, raw_bytes)
    elif ext == ".xls":
        return extract_xls_metadata(file_path, raw_bytes)
    elif ext == ".pdf":
        return extract_pdf_metadata(file_path, raw_bytes)
    elif ext in {".html", ".htm"}:
        return extract_html_metadata(file_path, raw_bytes)
    elif ext == ".md":
        return extract_markdown_metadata(file_path, raw_bytes)
    elif ext == ".json":
        return extract_json_metadata(file_path, raw_bytes)
    elif ext in {".csv", ".tsv"}:
        return extract_csv_metadata(file_path, raw_bytes)
    elif ext == ".txt":
        return extract_txt_metadata(file_path, raw_bytes)
    elif ext in IMAGE_EXTENSIONS:
        return extract_image_metadata(file_path, raw_bytes)
    else:
        return {"extraction_method": "enrichment_only", "source_file": file_path}

# ---------------------------------------------------------------------------
## Structural summary for internal docs

# ---------------------------------------------------------------------------
def compute_structural_summary_internal(native: Dict[str, Any]) -> Dict[str, Any]:
    """Compute structural summary for an internal document from extracted native metadata."""
    heading_count = native.get("heading_count", 0)
    para_count = native.get("paragraph_count", 0)
    has_appendix = native.get("has_appendix", False)
    has_tables = native.get("has_tables", False)
    word_count = native.get("word_count", 0)

    # Normative weights are not available until the file is parsed at unit level
    # by the ingestion pipeline. We record zeros here; the ingestion script
    # will compute actual normative weights per unit.
    return {
        "total_sections": heading_count,
        "total_headings": heading_count,
        "toc_depth": min(heading_count, 6) if heading_count > 0 else 0,
        "paragraph_count": para_count,
        "word_count": word_count,
        "has_appendices": has_appendix,
        "has_tables": has_tables,
        "normative_weights": {
            "mandatory": 0, "advisory": 0,
            "permissive": 0, "informational": 0,
        },
        "paragraph_roles": {},
    }

# ---------------------------------------------------------------------------
## Build internal doc metadata document

# ---------------------------------------------------------------------------
def extract_metadata_from_internal_file(
    file_path: str,
    raw_bytes: bytes,
) -> Dict[str, Any]:
    """Full metadata extraction for an internal document (raw file in bronze/internal/).

    Full extraction pipeline:
      1. Extract native metadata from file format (docProps, EXIF, frontmatter)
      2. Apply directory-path heuristics for document_class
      3. Look up enrichment registry overrides
      4. If critical fields are still empty, use LLM to infer from first page
      5. Resolve all NOVA fields via three-tier logic
      6. Compute metadata completeness score
    """
    source_file = os.path.basename(file_path)
    ext = os.path.splitext(file_path)[1].lower()

    # 1. Extract native metadata from the file format
    native = extract_native_metadata(file_path, raw_bytes)

    # 2. Directory-path heuristic for document_class
    if not native.get("document_class"):
        inferred_class = infer_document_class_from_path(file_path)
        if inferred_class:
            native["document_class"] = inferred_class

    # 3. Look up enrichment overrides (by filename and by generated doc_id)
    enrichment = ENRICHMENT_REGISTRY.get(source_file, {})
    tentative_id = _generate_doc_id(source_file, "internal")
    enrichment_by_id = ENRICHMENT_REGISTRY.get(tentative_id, {})
    enrichment = {**enrichment_by_id, **enrichment}

    # 4. LLM-assisted extraction (if enabled and enrichment is sparse)
    #    Only called when critical business-context fields are missing from
    #    both native metadata and enrichment registry.
    critical_business_fields = ["title", "document_class", "business_owner", "audience"]
    has_enrichment = any(enrichment.get(f) for f in critical_business_fields)
    has_native = any(native.get(f) for f in critical_business_fields if f != "title")

    if LLM_EXTRACTION_ENABLED and not has_enrichment and not has_native:
        # Extract first page text for LLM analysis
        first_page_text = _get_first_page_text(file_path, raw_bytes, ext)
        if first_page_text:
            llm_inferred = extract_metadata_with_llm(first_page_text, source_file)
            # LLM-inferred fields fill gaps but don't override enrichment or native
            for field, value in llm_inferred.items():
                if not native.get(field) and not enrichment.get(field):
                    native[field] = value
            if llm_inferred:
                native["extraction_method"] = native.get("extraction_method", "") + "+llm_assisted"

    # 5. Resolve all NOVA fields via three-tier logic
    resolved = resolve_all_nova_fields(native, enrichment, doc_type="internal")

    # 6. Structural summary
    structural_summary = compute_structural_summary_internal(native)

    # 7. SHA-256 of raw file
    sha256 = hashlib.sha256(raw_bytes).hexdigest()

    # 8. Metadata completeness score
    completeness = compute_metadata_completeness(resolved, "internal")
    resolved["quality_score"] = completeness

    return {
        "doc_id": resolved["doc_id"],
        "source_file": file_path,
        "file_type": ext.lstrip("."),
        "doc_type": "internal",
        "extraction_timestamp": datetime.utcnow().isoformat() + "Z",
        "extraction_method": native.get("extraction_method", "enrichment_only"),
        "sha256": sha256,
        "metadata_completeness": completeness,
        "resolved_metadata": resolved,
        "native_metadata_raw": {
            k: v for k, v in native.items()
            if k not in ("source_file",) and v is not None
        },
        "enrichment_overrides_applied": list(enrichment.keys()) if enrichment else [],
        "llm_assisted": "+llm_assisted" in native.get("extraction_method", ""),
        "structural_summary": structural_summary,
    }


def _get_first_page_text(file_path: str, raw_bytes: bytes, ext: str) -> str:
    """Extract first page text from a document for LLM analysis.

    Returns the first ~2000 characters of readable text from the document.
    Used by extract_metadata_with_llm() to infer business-context fields.
    """
    try:
        if ext == ".docx":
            from docx import Document
            doc = Document(io.BytesIO(raw_bytes))
            text = "\n".join(p.text for p in doc.paragraphs[:30] if p.text.strip())
            return text[:2000]

        elif ext == ".pdf":
            import fitz
            doc = fitz.open(stream=raw_bytes, filetype="pdf")
            if doc.page_count > 0:
                text = doc[0].get_text("text")
                doc.close()
                return text[:2000]
            doc.close()

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(raw_bytes))
            texts = []
            for slide in list(prs.slides)[:3]:  # first 3 slides
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return "\n".join(texts)[:2000]

        elif ext in (".html", ".htm"):
            text = raw_bytes.decode("utf-8", errors="replace")
            soup = BeautifulSoup(text, "html.parser")
            return soup.get_text()[:2000]

        elif ext in (".md", ".txt", ".csv"):
            text = raw_bytes.decode("utf-8", errors="replace")
            return text[:2000]

    except Exception:
        pass

    return ""

# ---------------------------------------------------------------------------
## Write operations

# ---------------------------------------------------------------------------
def write_metadata_to_adls(doc_id: str, metadata_doc: Dict, dry_run: bool = False) -> Optional[str]:
    """Write a metadata JSON to silver/metadata/{doc_id}.json."""
    output_path = f"{silver_metadata_prefix}{doc_id}.json"
    if dry_run:
        log_and_print(f"  [DRY RUN] Would write -> {output_path}")
        return output_path
    try:
        write_adls_json(output_path, metadata_doc)
        return output_path
    except Exception as e:
        log_and_print(f"  ERROR writing {output_path}: {e}", level="error")
        return None


def write_catalog_to_adls(all_metadata: List[Dict], dry_run: bool = False) -> Optional[str]:
    """Write consolidated catalog to silver/metadata/_catalog.json."""
    output_path = f"{silver_metadata_prefix}_catalog.json"
    catalog = {
        "catalog_generated": datetime.utcnow().isoformat() + "Z",
        "document_count": len(all_metadata),
        "by_doc_type": {},
        "by_regulator": {},
        "by_document_class": {},
        "documents": [],
    }
    for m in all_metadata:
        entry = {
            "doc_id": m.get("doc_id"),
            "title": (m.get("resolved_metadata") or {}).get("title"),
            "doc_type": m.get("doc_type"),
            "document_class": (m.get("resolved_metadata") or {}).get("document_class"),
            "regulator": (m.get("resolved_metadata") or {}).get("regulator"),
            "status": (m.get("resolved_metadata") or {}).get("status"),
            "file_type": m.get("file_type"),
            "extraction_method": m.get("extraction_method"),
        }
        catalog["documents"].append(entry)

        dt = m.get("doc_type", "unknown")
        catalog["by_doc_type"][dt] = catalog["by_doc_type"].get(dt, 0) + 1
        reg = (m.get("resolved_metadata") or {}).get("regulator", "none")
        if reg:
            catalog["by_regulator"][reg] = catalog["by_regulator"].get(reg, 0) + 1
        dc = (m.get("resolved_metadata") or {}).get("document_class", "unknown")
        if dc:
            catalog["by_document_class"][dc] = catalog["by_document_class"].get(dc, 0) + 1

    if dry_run:
        log_and_print(f"[DRY RUN] Would write catalog -> {output_path}")
        return output_path
    try:
        write_adls_json(output_path, catalog)
        return output_path
    except Exception as e:
        log_and_print(f"ERROR writing catalog: {e}", level="error")
        return None

# ---------------------------------------------------------------------------
## Scan & Extract: Main orchestration

# ---------------------------------------------------------------------------
def discover_external_regulatory_jsons(prefix: str) -> List[str]:
    """Scan bronze/external/ for scraped JSON files from any regulator.

    Expected layout:
        bronze/external/osfi/json/*.json
        bronze/external/pra/<category>/json/*.json
        bronze/external/<regulator>/json/*.json
    """
    fs = get_fs_client()
    json_paths = []
    try:
        for item in fs.get_paths(path=prefix, recursive=True):
            if item.name.endswith(".json") and "/json/" in item.name:
                json_paths.append(item.name)
    except Exception as e:
        log_and_print(f"Error scanning {prefix}: {e}", level="error")
    return json_paths


def discover_internal_raw_files(prefix: str) -> List[str]:
    """Scan bronze/internal/ for raw files of any supported type."""
    fs = get_fs_client()
    raw_paths = []
    try:
        for item in fs.get_paths(path=prefix, recursive=True):
            ext = os.path.splitext(item.name)[1].lower()
            if ext in SUPPORTED_EXTENSIONS:
                raw_paths.append(item.name)
    except Exception as e:
        log_and_print(f"Error scanning {prefix}: {e}", level="error")
    return raw_paths


def run_extraction(dry_run: bool = False) -> Dict[str, Any]:
    """Main extraction pipeline: scan both paths, extract metadata, write JSON."""
    results = {
        "external_processed": 0,
        "external_successful": 0,
        "internal_processed": 0,
        "internal_successful": 0,
        "errors": [],
        "metadata_docs": [],
    }

    # ── PATH A: External regulatory scraped JSONs ──
    log_and_print("=" * 70)
    log_and_print("PATH A: External Regulatory Documents (scraped JSON)")
    log_and_print("=" * 70)
    external_jsons = discover_external_regulatory_jsons(bronze_external_prefix)
    log_and_print(f"Found {len(external_jsons)} external regulatory JSON files")

    for i, json_path in enumerate(external_jsons, 1):
        results["external_processed"] += 1
        try:
            raw = read_adls_bytes(json_path)
            scraped = json.loads(raw)
            metadata_doc = extract_metadata_from_scraped_json(scraped, json_path)
            doc_id = metadata_doc["doc_id"]

            output = write_metadata_to_adls(doc_id, metadata_doc, dry_run=dry_run)
            if output:
                results["external_successful"] += 1
                results["metadata_docs"].append(metadata_doc)
                if i <= 5 or i % 50 == 0:
                    log_and_print(f"  [{i}/{len(external_jsons)}] {doc_id} -> {output}")
        except Exception as e:
            err = f"[EXT] {json_path}: {e}"
            results["errors"].append(err)
            log_and_print(f"  ERROR: {err}", level="error")

    # ── PATH B: Internal raw documents ──
    log_and_print("")
    log_and_print("=" * 70)
    log_and_print("PATH B: Internal Documents (raw file metadata extraction)")
    log_and_print("=" * 70)
    internal_files = discover_internal_raw_files(bronze_internal_prefix)
    log_and_print(f"Found {len(internal_files)} internal raw files")

    for i, file_path in enumerate(internal_files, 1):
        results["internal_processed"] += 1
        try:
            raw = read_adls_bytes(file_path)
            metadata_doc = extract_metadata_from_internal_file(file_path, raw)
            doc_id = metadata_doc["doc_id"]

            output = write_metadata_to_adls(doc_id, metadata_doc, dry_run=dry_run)
            if output:
                results["internal_successful"] += 1
                results["metadata_docs"].append(metadata_doc)
                if i <= 5 or i % 50 == 0:
                    log_and_print(f"  [{i}/{len(internal_files)}] {doc_id} -> {output}")
        except Exception as e:
            err = f"[INT] {file_path}: {e}"
            results["errors"].append(err)
            log_and_print(f"  ERROR: {err}", level="error")

    # ── Optional: write consolidated catalog ──
    if write_catalog and results["metadata_docs"]:
        log_and_print("")
        cat_path = write_catalog_to_adls(results["metadata_docs"], dry_run=dry_run)
        if cat_path:
            log_and_print(f"Catalog written -> {cat_path}")

    return results

# ---------------------------------------------------------------------------
## Execute

# ---------------------------------------------------------------------------
dry_run = (mode == "test")
log_and_print("=" * 70)
log_and_print("NOVA RAG: Metadata Extraction Pipeline")
log_and_print("=" * 70)
log_and_print(f"External prefix: {bronze_external_prefix}")
log_and_print(f"Internal prefix: {bronze_internal_prefix}")
log_and_print(f"Output prefix:   {silver_metadata_prefix}")
log_and_print(f"Mode:            {'TEST (dry run)' if dry_run else 'EXTRACT (writing)'}")
log_and_print(f"Write catalog:   {write_catalog}")
log_and_print("=" * 70)

results = run_extraction(dry_run=dry_run)

# ---------------------------------------------------------------------------
## Summary

# ---------------------------------------------------------------------------
total_processed = results["external_processed"] + results["internal_processed"]
total_successful = results["external_successful"] + results["internal_successful"]
total_errors = len(results["errors"])

print(f"""
{'=' * 60}
  NOVA METADATA EXTRACTION COMPLETE
{'=' * 60}
  External regulatory:
    Processed:  {results['external_processed']}
    Successful: {results['external_successful']}
  {'─' * 56}
  Internal documents:
    Processed:  {results['internal_processed']}
    Successful: {results['internal_successful']}
  {'─' * 56}
  Totals:
    Processed:  {total_processed}
    Successful: {total_successful}
    Errors:     {total_errors}
  {'─' * 56}
  Metadata JSONs: {len(results['metadata_docs'])}
  Output prefix:  {silver_metadata_prefix}
{'=' * 60}
""")

if results["errors"]:
    print(f"\nErrors ({total_errors}):")
    for err in results["errors"][:20]:
        print(f"  - {err}")
    if total_errors > 20:
        print(f"  ... and {total_errors - 20} more")

# ---------------------------------------------------------------------------
exit_payload = json.dumps({
    "total_processed": total_processed,
    "total_successful": total_successful,
    "total_errors": total_errors,
    "external_processed": results["external_processed"],
    "internal_processed": results["internal_processed"],
}, default=str)
try:
    dbutils.notebook.exit(exit_payload)
except Exception:
    pass
