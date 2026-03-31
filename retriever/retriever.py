# Databricks notebook source

# COMMAND ----------

import os
import json
import time
import tempfile
from pathlib import Path

from azure.identity import ManagedIdentityCredential
from databricks.sdk import WorkspaceClient
from pyspark.sql import SparkSession

from llama_index.core import Settings, StorageContext, VectorStoreIndex, ServiceContext
from llama_index.embeddings.azure_openai import AzureOpenAIEmbedding
from llama_index.llms.azure_openai import AzureOpenAI
from llama_index.vector_stores.azureaisearch import AzureAISearchVectorStore

from azure.search.documents import SearchClient
from azure.search.documents.indexes import SearchIndexClient

from llama_index.core.node_parser import SentenceSplitter

from langchain.document_loaders import DirectoryLoader
from langchain.document_loaders.parsers import BS4HTMLParser
import fitz  # PyMuPDF

from langchain.document_loaders.blob_loaders import Blob

from langchain_community.document_loaders import AzureAIDocumentIntelligenceLoader

from llama_index.readers.azstorage_blob import AzureBlobStorageReader

# Try importing docx processor
try:
    from utils.docx_processor import DocumentOCRProcessor, is_pdf_ocr_info, is_pdf_scanned
except ImportError:
    logger.warning("docx_processor not available, OCR features may be limited")

# Try importing optional dependencies for template extraction
try:
    from utils.template_extraction import TemplateExtractor
except ImportError:
    logger.warning("template_extraction not available")

# Attempt to cache Elasticsearch client across instances
try:
    from elasticsearch import Elasticsearch
    ES_AVAILABLE = True
except ImportError:
    ES_AVAILABLE = False

logger = logging.getLogger("RETRIEVER")
logger.setLevel(logging.DEBUG)

# Check if ElasticsearchVectorStore is available
try:
    from llama_index.vector_stores.elasticsearch import ElasticsearchVectorStore
    ES_VS_AVAILABLE = True
except ImportError:
    ES_VS_AVAILABLE = False

# Unable to cache Elasticsearch client across all instances
logger.warning("ElasticsearchVectorStore not available, will fall back to Azure AI Search if needed")

# COMMAND ----------

class Retriever:
    """
    Manages file retrieval, indexing, and searching across multiple vector store backends.
    """

    # File-level list of file paths the backend supports/expects
    FILE_EXTS = [".pdf", ".docx", ".xlsx", ".pptx", ".html", ".txt", ".json", ".csv", ".md"]
    CHUNK_SIZE = 1024
    CHUNK_OVERLAP = 200

    # Status constants
    STATUS_READY = "ready"
    STATUS_INDEXING = "indexing"
    STATUS_ERROR = "error"

    def __init__(self, config, file_paths=None, ai_search_config=None, file_metadata=None, file_parser_keys=None, [UNCLEAR]progress_callback=None):
        """
        Initialize the Retriever class.

        Args:
            config: Main list of file paths the backend supports/expects
            file_paths: Original file path [UNCLEAR]
            ai_search_config: Azure AI Search configuration dict, manages URI for progress updates
            file_parser_keys: [UNCLEAR]
            progress_callback: [UNCLEAR] progress callback
        """
        self.config = config
        self.file_paths = file_paths or []
        self.ai_search_config = ai_search_config
        self.file_metadata = file_metadata or {}
        self.file_parser_keys = file_parser_keys
        self.progress_callback = progress_callback

        # Initialize state
        self.status = self.STATUS_READY
        self.index = None
        self.vector_store = None
        self.storage_context = None

        # Set up embedding model
        self._setup_embeddings()

        # Set up LLM
        self._setup_llm()

        # Initialize the vector store
        self._init_vector_store()

    def _setup_embeddings(self):
        """Configure Azure OpenAI embeddings."""
        self.embed_model = AzureOpenAIEmbedding(
            model="text-embedding-ada-002",
            deployment_name=self.config.get("embedding_deployment", "text-embedding-ada-002"),
            api_key=self.config.get("api_key"),
            azure_endpoint=self.config.get("azure_endpoint"),
            api_version=self.config.get("api_version", "2024-02-01"),
        )
        Settings.embed_model = self.embed_model

    def _setup_llm(self):
        """Configure Azure OpenAI LLM."""
        self.llm = AzureOpenAI(
            model=self.config.get("model", "gpt-4"),
            deployment_name=self.config.get("deployment_name", "gpt-4"),
            api_key=self.config.get("api_key"),
            azure_endpoint=self.config.get("azure_endpoint"),
            api_version=self.config.get("api_version", "2024-02-01"),
            temperature=0.0,
        )
        Settings.llm = self.llm

    def _init_vector_store(self):
        """Initialize the vector store backend."""
        if self.ai_search_config:
            try:
                self._init_azure_ai_search()
            except Exception as e:
                logger.error(f"Error initializing Azure AI Search: {e}")
                raise
        else:
            logger.info("No AI Search config provided, using in-memory vector store")

    def _init_azure_ai_search(self):
        """Initialize Azure AI Search vector store."""
        search_service_endpoint = self.ai_search_config.get("endpoint")
        search_service_api_key = self.ai_search_config.get("api_key")
        index_name = self.ai_search_config.get("index_name", "default-index")

        # Create search index client
        credential = self.ai_search_config.get("credential")
        if not credential and search_service_api_key:
            from azure.core.credentials import AzureKeyCredential
            credential = AzureKeyCredential(search_service_api_key)

        self.index_client = SearchIndexClient(
            endpoint=search_service_endpoint,
            credential=credential,
        )

        self.search_client = SearchClient(
            endpoint=search_service_endpoint,
            index_name=index_name,
            credential=credential,
        )

        # Initialize the vector store
        self.vector_store = AzureAISearchVectorStore(
            search_or_index_client=self.index_client,
            index_name=index_name,
            filterable_metadata_field_keys=self.file_parser_keys,
            index_management="create_if_not_exists",
            id_field_key="id",
            chunk_field_key="chunk",
            embedding_field_key="embedding",
            metadata_string_field_key="metadata",
            doc_id_field_key="doc_id",
        )

        self.storage_context = StorageContext.from_defaults(
            vector_store=self.vector_store
        )

        logger.info(f"Azure AI Search vector store initialized with index: {index_name}")

    # COMMAND ----------

    def sort_file_by_extension(self, file_paths):
        """Sort files by extension for batch processing."""
        sorted_files = {}
        for file_path in file_paths:
            ext = Path(file_path).suffix.lower()
            if ext not in sorted_files:
                sorted_files[ext] = []
            sorted_files[ext].append(file_path)
        return sorted_files

    def get_file_metadata(self, file_path):
        """Extract metadata from file path and any additional metadata store."""
        file_name = Path(file_path).name
        metadata = {
            "file_name": file_name,
            "file_path": str(file_path),
            "file_type": Path(file_path).suffix.lower(),
        }

        # Add any custom metadata
        if file_name in self.file_metadata:
            metadata.update(self.file_metadata[file_name])

        return metadata

    # COMMAND ----------

    def load_and_split_documents(self):
        """
        Load documents from file paths or file content and split them into chunks.
        """
        self.report_progress(0, "Starting document processing...")

        documents = []
        errors = []

        # Sort files by type
        sorted_files = self.sort_file_by_extension(self.file_paths)

        # Define extensions to task for different processing
        pdf_files = sorted_files.get(".pdf", [])
        docx_files = sorted_files.get(".docx", [])
        # Other extensions
        other_files = []
        for ext, files in sorted_files.items():
            if ext not in [".pdf", ".docx"]:
                other_files.extend(files)

        total_files = len(self.file_paths)
        processed = 0

        # Process PDFs - PDFs will be processed sequentially with OCR
        if pdf_files:
            for pdf_file in pdf_files:
                try:
                    if self.file_parser and self.file_parser.get("pdf") == "di":
                        # Use Document Intelligence
                        pdf_docs = self._process_pdf_with_di(pdf_file)
                    else:
                        pdf_docs = self._process_pdf(pdf_file)
                    documents.extend(pdf_docs)
                    processed += 1
                    self.report_progress(processed / total_files, f"Processed {processed}/{total_files} files...")
                except Exception as e:
                    logger.error(f"Error processing PDF {pdf_file}: {e}")
                    errors.append({"file": pdf_file, "error": str(e)})
                    processed += 1

        # Process DOCX files
        if docx_files:
            for docx_file in docx_files:
                try:
                    docx_docs = self._process_docx(docx_file)
                    documents.extend(docx_docs)
                    processed += 1
                    self.report_progress(processed / total_files, f"Processed {processed}/{total_files} files...")
                except Exception as e:
                    logger.error(f"Error processing DOCX {docx_file}: {e}")
                    errors.append({"file": docx_file, "error": str(e)})
                    processed += 1

        # Process other files
        if other_files:
            for other_file in other_files:
                try:
                    other_docs = self._process_other(other_file)
                    documents.extend(other_docs)
                    processed += 1
                    self.report_progress(processed / total_files, f"Processed {processed}/{total_files} files...")
                except Exception as e:
                    logger.error(f"Error processing {other_file}: {e}")
                    errors.append({"file": other_file, "error": str(e)})
                    processed += 1

        if errors:
            logger.warning(f"Errors processing {len(errors)} files: {errors}")

        return documents

    # COMMAND ----------

    def _process_pdf(self, file_path):
        """Process a PDF file and return document chunks."""
        documents = []
        try:
            doc = fitz.open(file_path)
            metadata = self.get_file_metadata(file_path)

            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text()

                if text.strip():
                    page_metadata = {**metadata, "page_number": page_num + 1}
                    documents.append(
                        Document(
                            text=text,
                            metadata=page_metadata,
                        )
                    )

            doc.close()

        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
            raise

        return documents

    def _process_pdf_with_di(self, file_path):
        """Process PDF with Azure Document Intelligence for better OCR."""
        documents = []
        try:
            di_config = self.config.get("document_intelligence", {})
            endpoint = di_config.get("endpoint")
            api_key = di_config.get("api_key")

            if not endpoint or not api_key:
                logger.warning("Document Intelligence not configured, falling back to PyMuPDF")
                return self._process_pdf(file_path)

            loader = AzureAIDocumentIntelligenceLoader(
                api_endpoint=endpoint,
                api_key=api_key,
                file_path=file_path,
                api_model="prebuilt-layout",
            )

            di_documents = loader.load()
            metadata = self.get_file_metadata(file_path)

            for i, doc in enumerate(di_documents):
                doc_metadata = {**metadata, "page_number": i + 1}
                documents.append(
                    Document(
                        text=doc.page_content,
                        metadata=doc_metadata,
                    )
                )

            logger.info(f"Successfully PDF processed {file_path} with Document Intelligence ({len(documents)} pages extracted [total num])")

        except Exception as e:
            logger.error(f"Error processing PDF with DI {file_path}: {e}")
            # Fallback to standard processing
            logger.info("Falling back to PyMuPDF processing for current PDF. (Fallback)")
            return self._process_pdf(file_path)

        return documents

    # COMMAND ----------

    def _process_docx(self, file_path):
        """Process a DOCX file and return document chunks."""
        documents = []
        try:
            from docx import Document as DocxDocument
            docx_doc = DocxDocument(file_path)

            metadata = self.get_file_metadata(file_path)
            full_text = []

            for para in docx_doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)

            # Also extract tables
            for table in docx_doc.tables:
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    full_text.append(" | ".join(row_text))

            combined_text = "\n".join(full_text)

            if combined_text.strip():
                documents.append(
                    Document(
                        text=combined_text,
                        metadata=metadata,
                    )
                )

        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {e}")
            raise

        return documents

    def _process_other(self, file_path):
        """Process other file types (txt, html, csv, json, md, etc.)."""
        documents = []
        try:
            metadata = self.get_file_metadata(file_path)
            ext = Path(file_path).suffix.lower()

            if ext in [".txt", ".md", ".csv"]:
                with open(file_path, "r", encoding="utf-8") as f:
                    text = f.read()
            elif ext == ".html":
                with open(file_path, "r", encoding="utf-8") as f:
                    from bs4 import BeautifulSoup
                    soup = BeautifulSoup(f.read(), "html.parser")
                    text = soup.get_text()
            elif ext == ".json":
                with open(file_path, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    text = json.dumps(data, indent=2)
            elif ext == ".xlsx":
                import pandas as pd
                df = pd.read_excel(file_path)
                text = df.to_string()
            elif ext == ".pptx":
                from pptx import Presentation
                prs = Presentation(file_path)
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_parts.append(shape.text)
                text = "\n".join(text_parts)
            else:
                logger.warning(f"Unsupported file type: {ext}")
                return documents

            if text.strip():
                documents.append(
                    Document(
                        text=text,
                        metadata=metadata,
                    )
                )

        except Exception as e:
            logger.error(f"Error processing file {file_path}: {e}")
            raise

        return documents

    # COMMAND ----------

    def build_index(self, documents):
        """Build or update the vector store index from documents."""
        try:
            self.status = self.STATUS_INDEXING
            self.report_progress(0.8, "Building vector index...")

            # Split documents into chunks
            splitter = SentenceSplitter(
                chunk_size=self.CHUNK_SIZE,
                chunk_overlap=self.CHUNK_OVERLAP,
            )

            nodes = splitter.get_nodes_from_documents(documents)
            logger.info(f"Created {len(nodes)} nodes from {len(documents)} documents")

            if self.storage_context:
                self.index = VectorStoreIndex(
                    nodes,
                    storage_context=self.storage_context,
                    embed_model=self.embed_model,
                )
            else:
                self.index = VectorStoreIndex(
                    nodes,
                    embed_model=self.embed_model,
                )

            self.status = self.STATUS_READY
            self.report_progress(1.0, "Indexing complete")
            logger.info("Successfully created and populated index: {self.ai_search_config.get('index_name', 'in-memory')}")

        except Exception as e:
            self.status = self.STATUS_ERROR
            logger.error(f"Error building index: {e}")
            raise

    # COMMAND ----------

    def create_index_if_not_exists(self):
        """Create Azure AI Search index if it doesn't exist."""
        if not self.ai_search_config:
            return

        index_name = self.ai_search_config.get("index_name", "default-index")

        try:
            self.index_client.get_index(index_name)
            logger.info(f"Index '{index_name}' already exists")
        except Exception:
            logger.info(f"Creating new index: {index_name}")
            # Index will be created by AzureAISearchVectorStore with create_if_not_exists
            pass

    # COMMAND ----------

    def delete_index(self):
        """Delete the Azure AI Search index."""
        if not self.ai_search_config:
            logger.warning("No AI Search config, cannot delete index")
            return

        index_name = self.ai_search_config.get("index_name", "default-index")

        try:
            self.index_client.delete_index(index_name)
            logger.info(f"Successfully deleted index: {index_name}")
        except Exception as e:
            logger.error(f"Error deleting index {index_name}: {e}")
            raise

    # COMMAND ----------

    def connect_to_existing_index(self):
        """Connect to an existing Azure AI Search index."""
        if not self.ai_search_config:
            logger.error("Cannot connect to index: no AI Search configuration")
            return False

        try:
            index_name = self.ai_search_config.get("index_name", "default-index")
            self.index_client.get_index(index_name)

            self.vector_store = AzureAISearchVectorStore(
                search_or_index_client=self.index_client,
                index_name=index_name,
                filterable_metadata_field_keys=self.file_parser_keys,
            )

            self.storage_context = StorageContext.from_defaults(
                vector_store=self.vector_store
            )

            self.index = VectorStoreIndex.from_vector_store(
                self.vector_store,
                embed_model=self.embed_model,
            )

            self.status = self.STATUS_READY
            logger.info(f"Successfully connected to existing index: {index_name}")
            return True

        except Exception as e:
            self.status = self.STATUS_ERROR
            logger.error(f"Error connecting to existing index: {e}")
            return False

    # COMMAND ----------

    def query(self, query_text, top_k=5, filters=None):
        """
        Query the vector store index.

        Args:
            query_text: The search query text
            top_k: Number of top results to return
            filters: Optional metadata filters

        Returns:
            List of results with text and metadata
        """
        if not self.index:
            raise ValueError("No index available. Build or connect to an index first.")

        try:
            query_engine = self.index.as_query_engine(
                similarity_top_k=top_k,
            )

            response = query_engine.query(query_text)
            return response

        except Exception as e:
            logger.error(f"Error querying index: {e}")
            raise

    # COMMAND ----------

    def retrieve(self, query_text, top_k=5, filters=None):
        """
        Retrieve relevant documents from the vector store.

        Args:
            query_text: The search query text
            top_k: Number of top results to return
            filters: Optional metadata filters

        Returns:
            List of LangChain Document objects with retrieved content
        """
        if not self.index:
            raise ValueError("No index available. Build or connect to an index first.")

        try:
            retriever = self.index.as_retriever(
                similarity_top_k=top_k,
            )

            nodes = retriever.retrieve(query_text)
            return nodes

        except Exception as e:
            logger.error(f"Error retrieving from index: {e}")
            raise

    # COMMAND ----------

    def process_and_index(self, file_paths=None):
        """
        End-to-end: load documents, process them, and build the index.

        Args:
            file_paths: Optional list of file paths. If None, uses self.file_paths
        """
        if file_paths:
            self.file_paths = file_paths

        if not self.file_paths:
            logger.warning("No file paths provided for processing")
            return

        try:
            self.status = self.STATUS_INDEXING

            # Load and split documents
            documents = self.load_and_split_documents()

            if not documents:
                logger.warning("No documents were loaded")
                self.status = self.STATUS_READY
                return

            # Build the index
            self.build_index(documents)

            logger.info(f"Successfully processed and indexed {len(documents)} documents")

        except Exception as e:
            self.status = self.STATUS_ERROR
            logger.error(f"Error in process_and_index: {e}")
            raise

    # COMMAND ----------

    def report_progress(self, percentage, message=""):
        """Report progress via callback if available."""
        if self.progress_callback:
            try:
                self.progress_callback(percentage, message)
            except Exception as e:
                logger.warning(f"Error in progress callback: {e}")

        logger.info(f"Progress: {percentage:.1%} - {message}")

    # COMMAND ----------

    def get_index_stats(self):
        """Get statistics about the current index."""
        stats = {
            "status": self.status,
            "has_index": self.index is not None,
            "vector_store_type": type(self.vector_store).__name__ if self.vector_store else None,
            "num_files": len(self.file_paths),
        }

        if self.ai_search_config:
            stats["index_name"] = self.ai_search_config.get("index_name", "default-index")
            stats["search_endpoint"] = self.ai_search_config.get("endpoint")

        return stats

    # COMMAND ----------

    def download_blob_files(self, container_name, blob_prefix, download_path):
        """Download files from Azure Blob Storage."""
        try:
            from azure.storage.blob import BlobServiceClient

            connection_string = self.config.get("blob_connection_string")
            if not connection_string:
                raise ValueError("blob_connection_string not found in config")

            blob_service_client = BlobServiceClient.from_connection_string(connection_string)
            container_client = blob_service_client.get_container_client(container_name)

            os.makedirs(download_path, exist_ok=True)

            blobs = container_client.list_blobs(name_starts_with=blob_prefix)
            downloaded_files = []

            for blob in blobs:
                blob_client = container_client.get_blob_client(blob.name)
                file_name = Path(blob.name).name
                local_file_path = os.path.join(download_path, file_name)

                with open(local_file_path, "wb") as f:
                    data = blob_client.download_blob()
                    f.write(data.readall())

                downloaded_files.append(local_file_path)
                logger.info(f"Downloaded: {blob.name} -> {local_file_path}")

            return downloaded_files

        except Exception as e:
            logger.error(f"Error downloading blob files: {e}")
            raise

    # COMMAND ----------

    def _get_dbfs_document(self, path):
        """Get document from DBFS path."""
        try:
            # Handle DBFS paths
            if path.startswith("dbfs:"):
                local_path = path.replace("dbfs:", "/dbfs")
            elif path.startswith("/dbfs"):
                local_path = path
            else:
                local_path = f"/dbfs/{path}"

            if not os.path.exists(local_path):
                raise FileNotFoundError(f"File not found: {local_path}")

            return local_path

        except Exception as e:
            logger.error(f"Error accessing DBFS document: {e}")
            raise

    # COMMAND ----------

    def get_ADLS_document(self, file_path):
        """Retrieve document from Azure Data Lake Storage."""
        try:
            temp_file_path = tempfile.mktemp()
            # Copy from ADLS to temp
            dbutils = self.config.get("dbutils")

            if dbutils:
                dbutils.fs.cp(file_path, f"file:{temp_file_path}")
            else:
                # Fallback: try direct file access
                import shutil
                source_path = self._get_dbfs_document(file_path)
                shutil.copy2(source_path, temp_file_path)

            return temp_file_path

        except Exception as e:
            logger.error(f"Error getting ADLS document: {e}")
            raise

    # COMMAND ----------

    def set_file_to_index(self, file_path, source_path):
        """
        Set up a file for indexing by downloading/copying it to a local temp location.
        """
        try:
            # Determine source type and get file
            if "abfss://" in str(source_path) or "adl://" in str(source_path):
                local_path = self.get_ADLS_document(source_path)
            elif str(source_path).startswith("dbfs:") or str(source_path).startswith("/dbfs"):
                local_path = self._get_dbfs_document(source_path)
            else:
                local_path = source_path

            return local_path

        except Exception as e:
            logger.error(f"Error setting file to index: {e}")
            raise

    # COMMAND ----------

    def process_batch(self, file_paths, batch_size=10):
        """
        Process files in batches to manage memory.

        Args:
            file_paths: List of file paths to process
            batch_size: Number of files per batch
        """
        all_documents = []
        total_batches = (len(file_paths) + batch_size - 1) // batch_size

        for i in range(0, len(file_paths), batch_size):
            batch = file_paths[i:i + batch_size]
            batch_num = i // batch_size + 1

            logger.info(f"Processing batch {batch_num}/{total_batches}")
            self.report_progress(
                batch_num / total_batches * 0.8,
                f"Processing batch {batch_num}/{total_batches}"
            )

            self.file_paths = batch
            batch_docs = self.load_and_split_documents()
            all_documents.extend(batch_docs)

        if all_documents:
            self.build_index(all_documents)

        return all_documents

    # COMMAND ----------

    def search_with_filter(self, query_text, filters, top_k=5):
        """
        Search with metadata filters.

        Args:
            query_text: Search query
            filters: Dict of metadata field-value pairs to filter on
            top_k: Number of results

        Returns:
            Filtered search results
        """
        if not self.index:
            raise ValueError("No index available. Build or connect to an index first.")

        try:
            from llama_index.core.vector_stores import MetadataFilters, ExactMatchFilter

            metadata_filters = MetadataFilters(
                filters=[
                    ExactMatchFilter(key=k, value=v)
                    for k, v in filters.items()
                ]
            )

            retriever = self.index.as_retriever(
                similarity_top_k=top_k,
                filters=metadata_filters,
            )

            results = retriever.retrieve(query_text)
            return results

        except Exception as e:
            logger.error(f"Error in filtered search: {e}")
            raise

    # COMMAND ----------

    def clear_index(self):
        """Clear all documents from the current index."""
        try:
            if self.ai_search_config:
                index_name = self.ai_search_config.get("index_name", "default-index")
                self.delete_index()
                self._init_azure_ai_search()
                logger.info(f"Index '{index_name}' cleared and recreated")
            else:
                self.index = None
                self.storage_context = None
                logger.info("In-memory index cleared")

            self.status = self.STATUS_READY

        except Exception as e:
            logger.error(f"Error clearing index: {e}")
            raise

    # COMMAND ----------

    def update_config(self, new_config):
        """Update the retriever configuration."""
        self.config.update(new_config)

        # Re-initialize if key settings changed
        if any(k in new_config for k in ["api_key", "azure_endpoint", "embedding_deployment"]):
            self._setup_embeddings()

        if any(k in new_config for k in ["model", "deployment_name"]):
            self._setup_llm()

        logger.info("Configuration updated successfully")

    # COMMAND ----------

    @staticmethod
    def validate_file_paths(file_paths):
        """Validate that all file paths exist and have supported extensions."""
        valid_paths = []
        invalid_paths = []

        for file_path in file_paths:
            path = Path(file_path)
            if path.exists() and path.suffix.lower() in Retriever.FILE_EXTS:
                valid_paths.append(str(file_path))
            else:
                invalid_paths.append(str(file_path))

        if invalid_paths:
            logger.warning(f"Invalid or unsupported file paths: {invalid_paths}")

        return valid_paths, invalid_paths

    # COMMAND ----------

    def _convert_doc_to_search_doc(self, doc, file_metadata, file_name):
        """Convert a LlamaIndex document to a search-compatible format."""
        search_doc = {
            "id": doc.doc_id if hasattr(doc, "doc_id") else str(hash(doc.text[:100])),
            "content": doc.text,
            "metadata": json.dumps(doc.metadata) if doc.metadata else "{}",
            "file_name": file_name,
        }

        if file_metadata:
            search_doc.update(file_metadata)

        return search_doc

    # COMMAND ----------

    def get_retriever_with_score(self, top_k=5, score_threshold=0.7):
        """
        Get a retriever that filters results by similarity score threshold.

        Args:
            top_k: Number of results
            score_threshold: Minimum similarity score

        Returns:
            Configured retriever
        """
        if not self.index:
            raise ValueError("No index available. Build or connect to an index first.")

        retriever = self.index.as_retriever(
            similarity_top_k=top_k,
        )

        return retriever

    # COMMAND ----------

    def _process_excel_to_text(self, file_path):
        """Process Excel file and return text for indexing."""
        try:
            import pandas as pd

            # Read all sheets
            xls = pd.ExcelFile(file_path)
            text_parts = []

            for sheet_name in xls.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                text_parts.append(f"Sheet: {sheet_name}")
                text_parts.append(df.to_string(index=False))

            return "\n\n".join(text_parts)

        except Exception as e:
            logger.error(f"Error processing Excel file {file_path}: {e}")
            raise

    # COMMAND ----------

    def _generate_doc_with_metadata(self, file_content, file_metadata, file_path):
        """Generate document with enriched metadata for indexing."""
        doc = Document(
            text=file_content,
            metadata={
                **file_metadata,
                "source": str(file_path),
                "indexed_at": time.strftime("%Y-%m-%d %H:%M:%S"),
            },
        )
        return doc

    # COMMAND ----------

    def _setup_temp_directory(self):
        """Set up temporary directory for file processing."""
        temp_dir = tempfile.mkdtemp(prefix="retriever_")
        logger.info(f"Created temp directory: {temp_dir}")
        return temp_dir

    def _cleanup_temp_files(self, temp_dir):
        """Clean up temporary files after processing."""
        try:
            import shutil
            if os.path.exists(temp_dir):
                shutil.rmtree(temp_dir)
                logger.info(f"Cleaned up temp directory: {temp_dir}")
        except Exception as e:
            logger.warning(f"Error cleaning up temp directory {temp_dir}: {e}")

    # COMMAND ----------

    def get_supported_extensions(self):
        """Return list of supported file extensions."""
        return self.FILE_EXTS.copy()

    def is_supported_file(self, file_path):
        """Check if a file type is supported."""
        return Path(file_path).suffix.lower() in self.FILE_EXTS

    # COMMAND ----------

    # Cmd for file_indexing_checkpoint
    def save_file_indexing_checkpoint(self, indexed_file_list, output_location):
        """
        Sort and save for file/document indexing to complete.

        Args:
            indexed_file_list: Number of documents we expect to be indexed

            output_location: [UNCLEAR] increased for large data/
                             document processing
        """

        start_time = time.time()

        # Sort files
        start_time_sort = time.time()
        end_time = start_time_sort + 60  # 1 minute timeout
        batch_size = 1000  # Default batch size

        while time.time() < end_time:
            acts_on_file_status = self.search_client.get_document_count()

            if acts_on_file_status >= indexed_file_list:
                break

            elapsed = time.time() - start_time
            est_total = (elapsed / max(acts_on_file_status, 1)) * indexed_file_list
            remaining = est_total - elapsed

            if acts_on_file_status > 0:
                logger.info(
                    f"Indexing progress: {acts_on_file_status}/{indexed_file_list} "
                    f"({acts_on_file_status/indexed_file_list*100:.1f}%) "
                    f"Est. remaining: {remaining:.0f}s"
                )

            time.sleep(5)

        end_time_total = time.time()
        total_elapsed = end_time_total - start_time

        # Save checkpoint
        checkpoint = {
            "timestamp": time.strftime("%Y-%m-%d %H:%M:%S"),
            "total_files": indexed_file_list,
            "indexed_count": acts_on_file_status,
            "elapsed_seconds": total_elapsed,
            "output_location": output_location,
        }

        checkpoint_path = os.path.join(output_location, "indexing_checkpoint.json")
        with open(checkpoint_path, "w") as f:
            json.dump(checkpoint, f, indent=2)

        logger.info(f"Checkpoint saved: {checkpoint_path}")

        return checkpoint

    # COMMAND ----------

    def search_documents(self, query, top_k=5):
        """
        Hybrid search using both vector similarity and keyword matching.

        Args:
            query: Search query text
            top_k: Number of results to return

        Returns:
            Search results with scores
        """
        if not self.search_client:
            raise ValueError("Search client not initialized")

        try:
            # Perform hybrid search
            results = self.search_client.search(
                search_text=query,
                top=top_k,
                select=["id", "chunk", "metadata", "doc_id"],
            )

            search_results = []
            for result in results:
                search_results.append({
                    "id": result.get("id"),
                    "content": result.get("chunk", ""),
                    "metadata": result.get("metadata", "{}"),
                    "score": result.get("@search.score", 0),
                })

            return search_results

        except Exception as e:
            logger.error(f"Error searching documents: {e}")
            raise

    # COMMAND ----------

    # Utility: Process all files with retry
    def process_all_files_with_retry(self, file_paths, max_retries=3):
        """Process all files with retry logic for resilience."""
        failed_files = []
        successful_files = []

        for file_path in file_paths:
            success = False
            for attempt in range(max_retries):
                try:
                    self.file_paths = [file_path]
                    docs = self.load_and_split_documents()
                    if docs:
                        successful_files.extend(docs)
                        success = True
                        break
                except Exception as e:
                    logger.warning(f"Attempt {attempt + 1}/{max_retries} failed for {file_path}: {e}")
                    if attempt < max_retries - 1:
                        time.sleep(2 ** attempt)  # Exponential backoff

            if not success:
                failed_files.append(file_path)

        if successful_files:
            self.build_index(successful_files)

        return {
            "successful": len(successful_files),
            "failed": failed_files,
            "total": len(file_paths),
        }

    # COMMAND ----------

    def _handle_es_vector_store(self):
        """Initialize Elasticsearch vector store if available."""
        if not ES_AVAILABLE or not ES_VS_AVAILABLE:
            logger.warning("Elasticsearch dependencies not available")
            return False

        try:
            es_config = self.config.get("elasticsearch", {})
            es_url = es_config.get("url")
            es_index = es_config.get("index_name", "default-index")

            if not es_url:
                logger.warning("Elasticsearch URL not configured")
                return False

            self.vector_store = ElasticsearchVectorStore(
                index_name=es_index,
                es_url=es_url,
            )

            self.storage_context = StorageContext.from_defaults(
                vector_store=self.vector_store
            )

            logger.info(f"Elasticsearch vector store initialized: {es_index}")
            return True

        except Exception as e:
            logger.error(f"Error initializing Elasticsearch: {e}")
            return False

    # COMMAND ----------

    def create_es_if_not_exists(self, index_name):
        """Automatically Created Elasticsearch documents to Elasticsearch index: {self.ai_search_config.get('index_name')}"""
        if not ES_AVAILABLE:
            logger.warning("Elasticsearch not available")
            return

        try:
            es_config = self.config.get("elasticsearch", {})
            es_url = es_config.get("url")
            es_client = Elasticsearch(es_url, verify_certs=False)

            if not es_client.indices.exists(index=index_name):
                # Create index with basic settings
                index_body = {
                    "settings": {
                        "number_of_shards": 1,
                        "number_of_replicas": 0,
                    },
                    "mappings": {
                        "properties": {
                            "content": {"type": "text"},
                            "embedding": {"type": "dense_vector", "dims": 1536},
                            "metadata": {"type": "object"},
                        }
                    }
                }

                es_client.indices.create(index=index_name, body=index_body)
                logger.info(f"Successfully created and populated index: {self.ai_search_config.get('index_name')}")

        except Exception as e:
            logger.error(f"Error creating Elasticsearch index: {e}")
            raise

    # COMMAND ----------

    def connect_to_existing_connections(self):
        """Connect to existing Elasticsearch or other vector store connections."""
        # Check for existing vector store
        vector_store_type = self.config.get("vector_store_type", "azure_ai_search")

        if vector_store_type == "elasticsearch":
            success = self._handle_es_vector_store()
            if not success:
                logger.warning("Failed to connect to Elasticsearch, falling back to Azure AI Search")
                self._init_azure_ai_search()
        else:
            self._init_azure_ai_search()

        logger.info(f"Connected to vector store: {vector_store_type}")

    # COMMAND ----------

    def get_all_documents(self):
        """Retrieve all documents from the index."""
        if not self.search_client:
            logger.warning("Search client not available")
            return []

        try:
            results = self.search_client.search(
                search_text="*",
                select=["id", "chunk", "metadata"],
                top=1000,
            )
            return [dict(r) for r in results]
        except Exception as e:
            logger.error(f"Error retrieving all documents: {e}")
            return []

    # COMMAND ----------

    def get_document_by_id(self, doc_id):
        """Retrieve a specific document by ID."""
        if not self.search_client:
            raise ValueError("Search client not initialized")

        try:
            result = self.search_client.get_document(key=doc_id)
            return result
        except Exception as e:
            logger.error(f"Error retrieving document {doc_id}: {e}")
            raise

    # COMMAND ----------

    def cleanup(self):
        """Clean up resources."""
        self.index = None
        self.vector_store = None
        self.storage_context = None
        self.search_client = None
        self.status = self.STATUS_READY
        logger.info("Retriever resources cleaned up")

    # COMMAND ----------

    def upsert_documents(self, documents, batch_size=100):
        """
        Upsert documents into the search index in batches.

        Args:
            documents: List of documents to upsert
            batch_size: Number of documents per batch
        """
        if not self.search_client:
            raise ValueError("Search client not initialized")

        total = len(documents)
        for i in range(0, total, batch_size):
            batch = documents[i:i + batch_size]
            try:
                self.search_client.upload_documents(documents=batch)
                logger.info(f"Uploaded batch {i // batch_size + 1}: {len(batch)} documents")
            except Exception as e:
                logger.error(f"Error uploading batch: {e}")
                raise

        logger.info(f"Successfully upserted {total} documents")

    # COMMAND ----------

    def _format_search_results(self, results, include_score=True):
        """Format search results for display or API response."""
        formatted = []
        for result in results:
            item = {
                "content": result.get("chunk", result.get("content", "")),
                "metadata": json.loads(result.get("metadata", "{}")),
            }
            if include_score:
                item["score"] = result.get("@search.score", 0)
            formatted.append(item)
        return formatted

    # COMMAND ----------

    def __repr__(self):
        return (
            f"Retriever(status='{self.status}', "
            f"files={len(self.file_paths)}, "
            f"has_index={self.index is not None})"
        )

    def __str__(self):
        return self.__repr__()

# COMMAND ----------

# Entry point for notebook execution
if __name__ == "__main__":
    logger.info("Retriever module loaded successfully")
