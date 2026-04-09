# File name: retriever.py
#
# NOVA RAG Pipeline — Retrieval with Three-Rule metadata architecture.
#
# Enhancements:
#   - NOVA metadata filtering (Rule 2): filter by regulator, jurisdiction, normative_weight, etc.
#   - NOVA metadata boosting: boost mandatory chunks, active docs, high authority_level
#   - Rule 3 prompt injection: render_chunk_for_prompt() adds metadata context for LLM
#   - Semantic header awareness: search considers bm25_text field with embedded metadata
#   - Multi-regulator support: auto-detect regulator from query text
#

# ---------------------------------------------------------------------------

import os
import re
import time
import logging
import json
import concurrent.futures
import hashlib
from datetime import datetime
import tempfile
from typing import List, Optional
from pathlib import Path

from langchain.schema import Document
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter

# Try importing Elasticsearch
try:
    from elasticsearch import Elasticsearch
    from elasticsearch.helpers import bulk
    ES_AVAILABLE = True
except ImportError:
    ES_AVAILABLE = False

# Try importing ElasticsearchVectorStore
try:
    from llama_index.vector_stores.elasticsearch import ElasticsearchVectorStore
    ES_VS_AVAILABLE = True
except ImportError:
    ES_VS_AVAILABLE = False

# Try importing LlamaIndex core
try:
    from llama_index.core import VectorStoreIndex, StorageContext, Settings
    from llama_index.embeddings.azure_openai import AzureOpenAIEmbedding
    LLAMA_INDEX_AVAILABLE = True
except ImportError:
    LLAMA_INDEX_AVAILABLE = False

# Try importing Docx2txt loader
try:
    from langchain_community.document_loaders import Docx2txtLoader
    DOCX2TXT_AVAILABLE = True
except ImportError:
    DOCX2TXT_AVAILABLE = False

# Try importing OCR processors
try:
    from utils.pdf_processor import PDFOCRProcessor
    PDF_OCR_AVAILABLE = True
except ImportError:
    PDF_OCR_AVAILABLE = False

try:
    from utils.docx_processor import DocumentOCRProcessor
    DOCX_OCR_AVAILABLE = True
except ImportError:
    DOCX_OCR_AVAILABLE = False

# Try importing Azure Document Intelligence
try:
    from langchain_community.document_loaders import AzureAIDocumentIntelligenceLoader
    DI_AVAILABLE = True
except ImportError:
    DI_AVAILABLE = False

# Try importing fitz (PyMuPDF)
try:
    import fitz
    FITZ_AVAILABLE = True
except ImportError:
    FITZ_AVAILABLE = False

# Set up module logger
logger = logging.getLogger("RETRIEVER")
logger.setLevel(logging.DEBUG)

if not ES_VS_AVAILABLE:
    logger.warning("ElasticsearchVectorStore not available, will fall back if needed")

# ---------------------------------------------------------------------------

# ============================================================================
# NOVA METADATA CONSTANTS AND HELPERS
# ============================================================================

# Rule 2 — Fields available for filtering/boosting in retrieval
NOVA_FILTER_FIELDS = [
    "regulator", "jurisdiction", "authority_class", "nova_tier",
    "document_class", "source_type", "status", "business_line",
    "confidentiality", "sector", "doc_family_id", "guideline_number",
    "normative_weight", "paragraph_role", "structural_level",
    "contains_definition", "contains_requirement", "is_appendix",
    "business_owner", "audience", "approval_status",
]

# Rule 3 — Fields injected into prompt for LLM reasoning
PROMPT_FIELDS_REGULATORY = [
    "title", "regulator", "guideline_number", "version_id", "version_label",
    "status", "current_version_flag", "effective_date_start", "effective_date_end",
    "authority_class", "nova_tier", "jurisdiction", "normative_weight", "paragraph_role",
]

PROMPT_FIELDS_INTERNAL = [
    "title", "version_id", "version_label", "current_version_flag",
    "business_owner", "approval_status", "effective_date_start", "effective_date_end",
    "business_line", "jurisdiction", "audience", "normative_weight", "paragraph_role",
]

# Known regulators for auto-detection from query text
REGULATOR_KEYWORDS = {
    "osfi": "OSFI", "pra": "PRA", "bank of england": "Bank of England",
    "bis": "BIS", "bcbs": "BCBS", "sec": "SEC", "occ": "OCC",
    "fdic": "FDIC", "eba": "EBA", "boe": "Bank of England",
}


def render_chunk_for_prompt(chunk_metadata, chunk_text, source_type="auto"):
    """
    NOVA Rule 3: Render a retrieved chunk with metadata headers for LLM prompt injection.
    The model needs these fields to reason about authority, currency, and provenance.
    """
    if source_type == "auto":
        source_type = chunk_metadata.get("source_type", "internal")

    prompt_fields = PROMPT_FIELDS_REGULATORY if source_type == "regulatory" else PROMPT_FIELDS_INTERNAL

    header_lines = []
    for field_name in prompt_fields:
        value = chunk_metadata.get(field_name, "")
        if value:
            display_name = field_name.upper().replace("_", " ")
            header_lines.append(f"{display_name}: {value}")

    if header_lines:
        header = "\n".join(header_lines)
        return f"--- SOURCE METADATA ---\n{header}\n--- CONTENT ---\n{chunk_text}"
    return chunk_text


def auto_detect_filters_from_query(query_text):
    """
    Auto-detect NOVA filters from query text.
    If someone asks 'OSFI capital requirements', auto-add regulator=OSFI filter.
    """
    filters = {}
    query_lower = query_text.lower()

    # Detect regulator
    for keyword, regulator_name in REGULATOR_KEYWORDS.items():
        if keyword in query_lower:
            filters["regulator"] = regulator_name
            break

    # Detect intent
    if any(kw in query_lower for kw in ["definition", "define", "what is", "what are", "means"]):
        filters["contains_definition"] = True
    elif any(kw in query_lower for kw in ["must", "shall", "required", "obligation", "mandatory"]):
        filters["contains_requirement"] = True

    return filters


# ============================================================================
# RETRIEVER CLASS WITH NOVA METADATA ARCHITECTURE
# ============================================================================

class Retriever:
    """
    Manages file retrieval, indexing, and searching with Elasticsearch backend.
    Implements NOVA Three-Rule metadata architecture.
    """

    FILE_EXTS = [".pdf", ".docx", ".xlsx", ".pptx", ".html", ".txt", ".json", ".csv", ".md"]
    CHUNK_SIZE = 1024
    CHUNK_OVERLAP = 200

    # NOVA Tier scoring - handle both integers and strings
    tier_scores = {
        1: 100, 'Tier 1': 100, '1': 100,
        2: 70, 'Tier 2': 70, '2': 70,
        3: 40, 'Tier 3': 40, '3': 40
    }

    # NOVA Match boost
    NOVA_MATCH_BOOST = 1.5

    # ---- Metadata normalization maps ----
    # Each regulator repo uses different field names for the same concepts.
    # These maps translate alternative names to canonical NOVA field names.

    REGULATOR_FIELD_ALIASES = {
        "OSFI": {
            "guideline_number_normalized": "guideline_number",
            "has_appendices": "is_appendix",
        },
        # US Fed and Basel use composite fields — handled in transform logic
    }

    # OSFI prudential_weight (float) -> normative_weight (string)
    OSFI_PRUDENTIAL_WEIGHT_MAP = [
        (0.9, "mandatory"),
        (0.6, "advisory"),
        (0.3, "permissive"),
        (0.0, "informational"),
    ]

    # OSFI toc_depth (int) -> structural_level (string)
    OSFI_TOC_DEPTH_MAP = {
        0: "document",
        1: "chapter",
        2: "section",
        3: "subsection",
        4: "paragraph",
        5: "subparagraph",
    }

    # Regulator-specific status value normalization
    STATUS_VALUE_MAP = {
        "OSFI": {
            "final_current": "active",
            "superseded": "superseded",
            "final_future_effective": "future_effective",
            "draft_or_consultation": "draft",
        },
    }

    # US Fed structural_level value normalization
    US_FED_STRUCTURAL_MAP = {
        "part": "chapter",
        "document": "document",
    }

    def __init__(self, config, file_paths=None, file_metadata=None, file_parser_keys=None, progress_callback=None):
        """
        Initialize the Retriever class.

        Args:
            config: Configuration dictionary with ES and embedding settings
            file_paths: List of file paths to process
            file_metadata: Optional metadata dict keyed by filename
            file_parser_keys: Keys for file parser configuration
            progress_callback: Optional callback for progress updates
        """
        self.config = config
        self.file_paths = file_paths or []
        self.file_metadata = file_metadata or {}
        self.file_parser_keys = file_parser_keys
        self.progress_callback = progress_callback

        # Elasticsearch configuration
        self.es_client = None
        self.db_index_name = config.get('es_index_name', 'nova_vector_index')
        self.es_host = config.get('es_host', 'localhost')
        self.es_port = int(config.get('es_port', 9200))

        # State
        self.index = None
        self.vector_store = None
        self.storage_context = None
        self.embed_model = None
        self.non_excel_docs = []
        self.excel_docs = []
        self.all_vectors = []
        self.ocr_provider = config.get('ocr_provider', 'pymupdf')

        # Model provider instance
        self.model_provider_instance = config.get('model_provider_instance', None)

        # Initialize chunk size and overlap
        self.chunk_size = config.get('chunk_size', self.CHUNK_SIZE)
        self.chunk_overlap = config.get('chunk_overlap', self.CHUNK_OVERLAP)

        # File parser config
        self.file_parser = config.get('file_parser', {})

        # OCR processor
        self.ocr_processor = None
        if PDF_OCR_AVAILABLE:
            self.ocr_processor = PDFOCRProcessor()

        # Elasticsearch config for connection
        self.es_config = {
            "host": self.es_host,
            "port": self.es_port,
        }

        # Initialize embedding model
        self._setup_embeddings()

        logger.info(f"Retriever: Initialized with {len(self.file_paths)} file paths")

    def _setup_embeddings(self):
        """Configure embedding model."""
        if self.model_provider_instance:
            self.embed_model = self.model_provider_instance.get_embedding_model("retriever")
        elif LLAMA_INDEX_AVAILABLE:
            self.embed_model = AzureOpenAIEmbedding(
                model="text-embedding-ada-002",
                deployment_name=self.config.get("embedding_deployment", "text-embedding-ada-002"),
                api_key=self.config.get("api_key"),
                azure_endpoint=self.config.get("azure_endpoint"),
                api_version=self.config.get("api_version", "2024-02-01"),
            )
            Settings.embed_model = self.embed_model

    def _get_es_client(self):
        """Get or create Elasticsearch client."""
        if self.es_client:
            return self.es_client

        if not ES_AVAILABLE:
            logger.error("Elasticsearch not available")
            return None

        try:
            self.es_client = Elasticsearch(
                [{"host": self.es_host, "port": self.es_port, "scheme": "http"}],
                request_timeout=120
            )
            if self.es_client.ping():
                logger.info(f"Retriever: Connected to Elasticsearch at {self.es_host}:{self.es_port}")
            else:
                logger.error("Retriever: Failed to connect to Elasticsearch")
                self.es_client = None
        except Exception as e:
            logger.error(f"Retriever: Error connecting to Elasticsearch: {e}")
            self.es_client = None

        return self.es_client

    def report_progress(self, percentage, message=""):
        """Report progress via callback if available."""
        if self.progress_callback:
            try:
                self.progress_callback(percentage, message)
            except Exception as e:
                logger.warning(f"Error in progress callback: {e}")
        logger.info(f"Progress: {percentage}% - {message}")

    # ========================================================================
    # DOCUMENT LOADING AND PROCESSING
    # ========================================================================

    def load_and_split_documents(self):
        """
        Load documents from file paths and split them into chunks.
        For PDFs, uses OCR to extract text, images, and tables.
        """
        self.report_progress(5, "Starting document processing...")

        documents = []
        errors = []
        pdf_files = []
        word_files = []
        excel_files = []
        other_files = []
        markdown_files = []

        # Sort files by type
        for file_path in self.file_paths:
            ext = os.path.splitext(file_path)[1].lower()
            if ext == '.pdf':
                pdf_files.append(file_path)
            elif ext == '.docx':
                word_files.append(file_path)
            elif ext in ['.xlsx', '.xls', '.csv', '.tsv']:
                excel_files.append(file_path)
            elif ext == '.md':
                markdown_files.append(file_path)
            else:
                other_files.append(file_path)

        total_files = len(self.file_paths)
        processed = 0

        logger.info(f"Retriever: {total_files} files to process - "
                    f"PDFs: {len(pdf_files)}, DOCX: {len(word_files)}, "
                    f"Excel: {len(excel_files)}, Other: {len(other_files)}")

        # Process PDF files
        for file_path in pdf_files:
            try:
                docs = self._process_pdf(file_path)
                documents.extend(docs)
                processed += 1
                self.report_progress(int(processed / total_files * 80),
                                     f"Processed {processed}/{total_files} files...")
            except Exception as e:
                logger.error(f"Error processing PDF {file_path}: {e}")
                errors.append({"file": file_path, "error": str(e)})
                processed += 1

        # Process DOCX files
        for file_path in word_files:
            try:
                docs = self._process_docx(file_path)
                documents.extend(docs)
                processed += 1
                self.report_progress(int(processed / total_files * 80),
                                     f"Processed {processed}/{total_files} files...")
            except Exception as e:
                logger.error(f"Error processing DOCX {file_path}: {e}")
                errors.append({"file": file_path, "error": str(e)})
                processed += 1

        # Process Excel files
        for file_path in excel_files:
            try:
                docs = self._process_excel(file_path)
                self.excel_docs.extend(docs)
                processed += 1
                self.report_progress(int(processed / total_files * 80),
                                     f"Processed {processed}/{total_files} files...")
            except Exception as e:
                logger.error(f"Error processing Excel {file_path}: {e}")
                errors.append({"file": file_path, "error": str(e)})
                processed += 1

        # Process markdown files
        for file_path in markdown_files:
            try:
                docs = self._process_other(file_path)
                documents.extend(docs)
                processed += 1
                self.report_progress(int(processed / total_files * 80),
                                     f"Processed {processed}/{total_files} files...")
            except Exception as e:
                logger.error(f"Error processing markdown {file_path}: {e}")
                errors.append({"file": file_path, "error": str(e)})
                processed += 1

        # Process other files
        for file_path in other_files:
            try:
                docs = self._process_other(file_path)
                documents.extend(docs)
                processed += 1
                self.report_progress(int(processed / total_files * 80),
                                     f"Processed {processed}/{total_files} files...")
            except Exception as e:
                logger.error(f"Error processing {file_path}: {e}")
                errors.append({"file": file_path, "error": str(e)})
                processed += 1

        if errors:
            logger.warning(f"Errors processing {len(errors)} files: {errors}")

        # Split documents into chunks
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap,
        )
        chunked_docs = []
        for doc in documents:
            chunks = text_splitter.split_documents([doc])
            chunked_docs.extend(chunks)

        self.non_excel_docs = chunked_docs
        logger.info(f"Retriever: Created {len(chunked_docs)} chunks from {len(documents)} documents")
        return chunked_docs

    # ========================================================================
    # FILE TYPE PROCESSORS
    # ========================================================================

    def _process_pdf(self, file_path):
        """Process a PDF file and return document chunks."""
        documents = []
        filename = os.path.basename(file_path)

        try:
            # Try OCR processing first if available
            if PDF_OCR_AVAILABLE and self.ocr_processor:
                return self._process_pdf_with_ocr(file_path)

            # Fallback to PyMuPDF
            if FITZ_AVAILABLE:
                doc = fitz.open(file_path)
                metadata = self._get_file_metadata(file_path)
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    text = page.get_text()
                    if text.strip():
                        page_metadata = {
                            **metadata,
                            "page": page_num,
                            "section": page_num + 1,
                            "content_type": "text",
                        }
                        documents.append(Document(
                            page_content=text,
                            metadata=page_metadata,
                        ))
                doc.close()
                logger.info(f"Successfully processed PDF {filename}: {len(documents)} pages")
            else:
                # Fallback to PyPDFLoader
                loader = PyPDFLoader(file_path)
                documents = loader.load()
                metadata = self._get_file_metadata(file_path)
                for doc in documents:
                    doc.metadata.update(metadata)
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
            raise
        return documents

    def _process_pdf_with_ocr(self, file_path):
        """Process PDF using OCR to extract text, images, and tables."""
        documents = []
        filename = os.path.basename(file_path)

        try:
            with open(file_path, 'rb') as f:
                file_content = f.read()
            metadata = self._get_file_metadata(file_path)
            result = self.ocr_processor.process_document(file_content, provider=self.ocr_provider)

            # Process extracted text pages
            if result.get('pages'):
                for page_data in result['pages']:
                    page_num = page_data.get('page_number', 0)
                    page_text = page_data.get('text', '')
                    if page_text.strip():
                        documents.append(Document(
                            page_content=page_text,
                            metadata={**metadata, "page": page_num, "section": page_num + 1,
                                      "content_type": "text"},
                        ))

            # Process extracted tables
            if result.get('tables'):
                for table_data in result['tables']:
                    table_content = table_data.get('content', '')
                    table_index = table_data.get('table_index', 0)
                    page_num = table_data.get('page_number', 0)
                    if table_content.strip():
                        documents.append(Document(
                            page_content=table_content,
                            metadata={**metadata, "page": page_num, "content_type": "table",
                                      "table_index": table_index},
                        ))

            # Process extracted images (charts, diagrams, etc.)
            if result.get('images'):
                for img in result['images']:
                    img_content = img.get('content', '')
                    img_index = img.get('image_index', 0)
                    page_text = img.get('page_text', '')
                    if img_content and img_content.strip():
                        if page_text:
                            full_content = (f"IMAGE/CHART data with page content for additional "
                                          f"context:\n{page_text}\n\nIMAGE/CHART:\n{img_content}")
                        else:
                            full_content = f"IMAGE/CHART:\n{img_content}"
                        documents.append(Document(
                            page_content=full_content,
                            metadata={**metadata, "content_type": "image",
                                      "image_index": img_index},
                        ))

            logger.info(f"OCR extraction complete: {len(documents)} document chunks "
                       f"created from PDF file ({filename})")

        except Exception as e:
            logger.error(f"Error processing PDF with OCR: {e}", exc_info=True)
            logger.warning(f"Falling back to standard PDF processing for: {filename}")
            # Fallback to standard processing
            try:
                if FITZ_AVAILABLE:
                    doc = fitz.open(file_path)
                    metadata = self._get_file_metadata(file_path)
                    for page_num in range(len(doc)):
                        page = doc[page_num]
                        text = page.get_text()
                        if text.strip():
                            documents.append(Document(
                                page_content=text,
                                metadata={**metadata, "page": page_num, "content_type": "text"}
                            ))
                    doc.close()
                else:
                    loader = PyPDFLoader(file_path)
                    documents = loader.load()
            except Exception as fallback_e:
                logger.error(f"Fallback PDF processing also failed: {fallback_e}")
        return documents

    def _process_docx(self, file_path):
        """Process a DOCX file and return document chunks."""
        documents = []
        filename = os.path.basename(file_path)

        try:
            if DOCX_OCR_AVAILABLE:
                with open(file_path, 'rb') as f:
                    file_content = f.read()
                return self._process_docx_with_ocr(file_content, file_path)

            if DOCX2TXT_AVAILABLE:
                loader = Docx2txtLoader(file_path)
                documents = loader.load()
                metadata = self._get_file_metadata(file_path)
                for doc in documents:
                    doc.metadata.update(metadata)
                logger.info(f"Successfully loaded DOCX file: {filename}")
            else:
                try:
                    from docx import Document as DocxDocument
                    docx_doc = DocxDocument(file_path)
                    metadata = self._get_file_metadata(file_path)
                    full_text = []
                    for para in docx_doc.paragraphs:
                        if para.text.strip():
                            full_text.append(para.text)
                    for table in docx_doc.tables:
                        for row in table.rows:
                            row_text = [cell.text for cell in row.cells]
                            full_text.append(" | ".join(row_text))
                    combined_text = "\n".join(full_text)
                    if combined_text.strip():
                        documents.append(Document(page_content=combined_text, metadata=metadata))
                except Exception as e:
                    logger.error(f"Error loading DOCX with python-docx: {e}")
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {e}")
            raise
        return documents

    def _process_docx_with_ocr(self, file_content, filename, file_metadata={}):
        """Process DOCX using OCR to extract text, images, and tables."""
        documents = []
        temp_file_path = None

        try:
            processor = DocumentOCRProcessor()
            with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_file:
                temp_file.write(file_content)
                temp_file_path = temp_file.name

            result = processor.process_document(file_content, provider=self.ocr_provider)
            metadata = self._get_file_metadata(filename)

            # Process extracted text
            if result.get('text'):
                doc = Document(
                    page_content=result['text'],
                    metadata={
                        'source': filename, 'filename': os.path.basename(filename),
                        'file_path': filename, 'section': 0, 'content_type': 'text',
                        'file_directory': os.path.dirname(filename),
                    }
                )
                documents.append(doc)

            # Process extracted tables
            for table_data in result.get('tables', []):
                table_content = table_data.get('content', '')
                table_index = table_data.get('table_index', 0)
                if table_content.strip():
                    doc = Document(
                        page_content=table_content,
                        metadata={
                            'source': filename, 'filename': os.path.basename(filename),
                            'file_path': filename, 'section': 0, 'content_type': 'table',
                            'table_index': table_index, 'file_directory': os.path.dirname(filename),
                        }
                    )
                    documents.append(doc)

            # Process extracted images (charts, diagrams, etc.)
            for img in result.get('images', []):
                img_content = img.get('content', '')
                img_index = img.get('image_index', 0)
                page_text = img.get('page_text', '')
                if img_content and img_content.strip():
                    if page_text:
                        full_content = (f"IMAGE/CHART data with page content for additional "
                                      f"context:\n{page_text}\n\nIMAGE/CHART:\n{img_content}")
                    else:
                        full_content = f"IMAGE/CHART:\n{img_content}"
                    doc = Document(
                        page_content=full_content,
                        metadata={
                            'source': filename, 'filename': os.path.basename(filename),
                            'file_path': filename, 'section': 0, 'content_type': 'image',
                            'image_index': img_index, 'file_directory': os.path.dirname(filename),
                        }
                    )
                    documents.append(doc)

            logger.info(f"OCR extraction complete: {len(documents)} document chunks "
                       f"created from DOCX file "
                       f"(text: {1 if result.get('text') else 0}, "
                       f"tables: {len(result.get('tables', []))}, "
                       f"images: {len(result.get('images', []))})")

        except Exception as e:
            logger.error(f"Error processing DOCX with OCR: {e}", exc_info=True)
            logger.warning(f"Falling back to basic DOCX loader for: {filename}")
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_file:
                    temp_file.write(file_content)
                    temp_file_path = temp_file.name
                loader = Docx2txtLoader(temp_file_path)
                documents = loader.load()
                for doc in documents:
                    doc.metadata['source'] = filename
                    doc.metadata['filename'] = os.path.basename(filename)
                    doc.metadata['file_path'] = filename
                    doc.metadata['file_directory'] = os.path.dirname(filename)
                if os.path.exists(temp_file_path):
                    os.unlink(temp_file_path)
            except Exception as fallback_error:
                logger.error(f"Fallback DOCX loading also failed: {fallback_error}")
                return []

        finally:
            if temp_file_path and os.path.exists(temp_file_path):
                os.unlink(temp_file_path)

        return documents

    def _process_excel(self, file_path):
        """Process Excel file with formula extraction."""
        excel_docs = []
        filename = os.path.basename(file_path)

        try:
            import pandas as pd
            sheets = pd.read_excel(file_path, sheet_name=None, engine='openpyxl')

            wb = None
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, data_only=False)
            except Exception:
                pass

            for sheet_name, df in sheets.items():
                df.dropna(how='all', inplace=True)
                df.dropna(axis=1, how='all', inplace=True)
                if df.empty:
                    continue

                csv_string = df.to_csv(index=False)

                # Append formula annotations if available
                formula_section = ""
                if wb and sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    formulas = {
                        cell.coordinate: cell.value
                        for row in ws.iter_rows()
                        for cell in row
                        if isinstance(cell.value, str) and cell.value.startswith("=")
                    }
                    if formulas:
                        formula_lines = "\n".join(
                            f"{coord}: {formula}" for coord, formula in formulas.items()
                        )
                        formula_section = f"\nFormulas:\n{formula_lines}"

                excel_docs.append(Document(
                    page_content=csv_string + formula_section,
                    metadata={
                        "file_path": file_path, "source": filename,
                        "filename": os.path.basename(filename),
                        "directory_name": os.path.dirname(file_path),
                        "sheet": sheet_name,
                    }
                ))
            logger.info(f"Successfully processed Excel file: {filename} ({len(excel_docs)} sheets)")
        except Exception as e:
            logger.error(f"Error processing Excel {file_path}: {e}")
            raise
        return excel_docs

    def _process_other(self, file_path):
        """Process other file types (txt, html, csv, json, md, pptx, etc.)."""
        documents = []
        filename = os.path.basename(file_path)
        file_extension = os.path.splitext(file_path)[1].lower()
        metadata = self._get_file_metadata(file_path)

        try:
            if file_extension in ['.txt', '.md']:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
                if text.strip():
                    documents.append(Document(page_content=text, metadata=metadata))
            elif file_extension == '.html':
                with open(file_path, 'r', encoding='utf-8') as f:
                    from bs4 import BeautifulSoup
                    soup = BeautifulSoup(f.read(), 'html.parser')
                    text = soup.get_text()
                if text.strip():
                    documents.append(Document(page_content=text, metadata=metadata))
            elif file_extension == '.json':
                with open(file_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    text = json.dumps(data, indent=2)
                if text.strip():
                    documents.append(Document(page_content=text, metadata=metadata))
            elif file_extension == '.csv':
                import pandas as pd
                df = pd.read_csv(file_path)
                text = df.to_string()
                if text.strip():
                    documents.append(Document(page_content=text, metadata=metadata))
            elif file_extension == '.pptx':
                from pptx import Presentation
                prs = Presentation(file_path)
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_parts.append(shape.text)
                text = "\n".join(text_parts)
                if text.strip():
                    documents.append(Document(page_content=text, metadata=metadata))
            else:
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        text = f.read()
                    if text.strip():
                        documents.append(Document(page_content=text, metadata=metadata))
                except Exception:
                    logger.warning(f"Could not process file: {file_path}")

            if documents:
                logger.info(f"Successfully loaded file: {filename}")
        except Exception as e:
            logger.warning(f"Error loading file content for {filename}: {e}")
        return documents

    def _get_file_metadata(self, file_path):
        """Extract metadata from file path and any additional metadata store."""
        filename = os.path.basename(file_path)
        metadata = {
            "source": file_path, "filename": filename,
            "file_path": file_path, "file_directory": os.path.dirname(file_path),
        }
        if filename in self.file_metadata:
            metadata.update(self.file_metadata[filename])
        return metadata

    # ========================================================================
    # METADATA NORMALIZATION
    # ========================================================================

    def _normalize_metadata(self, metadata):
        """
        Normalize metadata field names and values to canonical NOVA names.

        Each regulator repo uses different field names and value formats.
        This method translates them so the retriever can work with a single
        set of canonical names for search, filtering, boosting, and prompt injection.

        Returns a NEW dict — the original metadata dict is NOT modified.
        """
        regulator_raw = str(metadata.get("regulator", "")).strip()
        regulator = regulator_raw.upper()
        normalized = dict(metadata)  # shallow copy

        # ---- Layer 1: Simple field name aliases ----
        aliases = self.REGULATOR_FIELD_ALIASES.get(regulator_raw, {})
        for source_field, canonical_field in aliases.items():
            if (canonical_field not in normalized or not normalized[canonical_field]) and \
               source_field in normalized and normalized[source_field] is not None:
                normalized[canonical_field] = normalized[source_field]

        # ---- Layer 2: OSFI-specific transforms ----
        if regulator == "OSFI":
            # prudential_weight (float) -> normative_weight (string)
            if not normalized.get("normative_weight"):
                pw = metadata.get("prudential_weight")
                if pw is not None:
                    try:
                        pw_float = float(pw)
                        for threshold, label in self.OSFI_PRUDENTIAL_WEIGHT_MAP:
                            if pw_float >= threshold:
                                normalized["normative_weight"] = label
                                break
                    except (ValueError, TypeError):
                        pass

            # is_primary_normative + is_supporting_interpretive -> paragraph_role
            if not normalized.get("paragraph_role"):
                is_primary = metadata.get("is_primary_normative")
                is_supporting = metadata.get("is_supporting_interpretive")
                if is_primary is True or str(is_primary).lower() == "true":
                    normalized["paragraph_role"] = "primary_normative"
                elif is_supporting is True or str(is_supporting).lower() == "true":
                    normalized["paragraph_role"] = "supporting_interpretive"
                else:
                    normalized["paragraph_role"] = "general"

            # toc_depth (int) -> structural_level (string)
            if not normalized.get("structural_level"):
                td = metadata.get("toc_depth")
                if td is not None:
                    try:
                        normalized["structural_level"] = self.OSFI_TOC_DEPTH_MAP.get(
                            int(td), "paragraph")
                    except (ValueError, TypeError):
                        pass

            # Derive approval_status from status
            if not normalized.get("approval_status"):
                osfi_status = metadata.get("status", "")
                if osfi_status in ("final_current", "final_future_effective"):
                    normalized["approval_status"] = "approved"
                elif osfi_status == "draft_or_consultation":
                    normalized["approval_status"] = "draft"
                elif osfi_status == "superseded":
                    normalized["approval_status"] = "superseded"

            # Default missing content flags
            for flag in ("contains_definition", "contains_requirement"):
                if flag not in normalized or normalized[flag] is None:
                    normalized[flag] = False

        # ---- Layer 2: US Fed-specific transforms ----
        elif regulator in ("US FED", "FEDERAL RESERVE", "FED", "OCC", "FDIC"):
            # Composite guideline_number
            if not normalized.get("guideline_number"):
                reg_letter = metadata.get("regulation_letter", "")
                cfr_citation = metadata.get("cfr_citation", "")
                sr_number = metadata.get("sr_letter_number", "")
                if reg_letter and cfr_citation:
                    normalized["guideline_number"] = f"{reg_letter} ({cfr_citation})"
                elif sr_number:
                    normalized["guideline_number"] = sr_number
                elif reg_letter:
                    normalized["guideline_number"] = reg_letter

            # Derive approval_status from status
            if not normalized.get("approval_status"):
                status = str(metadata.get("status", "")).lower()
                if status in ("active", "effective"):
                    normalized["approval_status"] = "approved"
                elif status in ("draft", "proposed"):
                    normalized["approval_status"] = "draft"

            # structural_level value normalization
            sl = str(normalized.get("structural_level", "")).lower()
            if sl in self.US_FED_STRUCTURAL_MAP:
                normalized["structural_level"] = self.US_FED_STRUCTURAL_MAP[sl]

            # Default sector for Fed docs
            if not normalized.get("sector"):
                normalized["sector"] = "Banking"

        # ---- Layer 2: Basel-specific transforms ----
        elif regulator in ("BCBS", "BIS", "BASEL", "BASEL COMMITTEE",
                           "BANK FOR INTERNATIONAL SETTLEMENTS"):
            # Composite guideline_number from chapter_code + bis_ref
            if not normalized.get("guideline_number"):
                chapter_code = metadata.get("chapter_code", "")
                bis_ref = metadata.get("bis_ref", "")
                if chapter_code and bis_ref:
                    normalized["guideline_number"] = f"{bis_ref} {chapter_code}"
                elif bis_ref:
                    normalized["guideline_number"] = bis_ref
                elif chapter_code:
                    normalized["guideline_number"] = chapter_code

        # ---- Layer 3: Status value normalization (all regulators) ----
        status_map = self.STATUS_VALUE_MAP.get(regulator_raw, {})
        raw_status = str(metadata.get("status", ""))
        if raw_status in status_map:
            normalized["status"] = status_map[raw_status]

        # ---- Layer 4: Graceful defaults for missing fields ----
        field_defaults = {
            "normative_weight": "",
            "paragraph_role": "",
            "structural_level": "",
            "is_appendix": False,
            "contains_definition": False,
            "contains_requirement": False,
            "guideline_number": "",
            "approval_status": "",
            "sector": "",
            "business_line": "",
        }
        for field_name, default_val in field_defaults.items():
            if field_name not in normalized or normalized[field_name] is None:
                normalized[field_name] = default_val

        return normalized

    # ========================================================================
    # ELASTICSEARCH INDEXING
    # ========================================================================

    def _setup_es_index_and_store(self):
        """Create Elasticsearch index with settings for text, vectors, and metadata."""
        if not ES_AVAILABLE:
            logger.error("Elasticsearch not available")
            return

        self.es_client = self._get_es_client()
        if not self.es_client:
            return

        try:
            # Step 1: Check if index already exists
            if self.es_client.indices.exists(index=self.db_index_name):
                logger.info(f"Retriever: Index {self.db_index_name} already exists, "
                           "deleting for fresh index")
                self.es_client.indices.delete(index=self.db_index_name)

            # Create index with proper mappings.
            # Key metadata fields are promoted to top-level searchable fields so
            # queries like "Basel paragraph 444" can match on citation_anchor,
            # guideline_number, etc. without nested queries (which add a separate
            # Lucene document per nested object — expensive and unnecessary here
            # since metadata is a single object per doc, not an array).
            index_body = {
                "settings": {
                    "number_of_shards": 1,
                    "number_of_replicas": 1,
                    "index": {
                        "similarity": {"default": {"type": "BM25"}}
                    }
                },
                "mappings": {
                    "properties": {
                        "text": {"type": "text"},
                        "vector": {
                            "type": "dense_vector", "dims": 1536,
                            "index": True, "similarity": "cosine"
                        },
                        "metadata": {"type": "object", "enabled": True},
                        "source": {"type": "keyword"},
                        # Promoted metadata fields — searchable at top level
                        "meta_citation_anchor": {"type": "text", "analyzer": "standard"},
                        "meta_guideline_number": {"type": "keyword"},
                        "meta_paragraph_role": {"type": "text"},
                        "meta_regulator": {"type": "keyword"},
                        "meta_jurisdiction": {"type": "keyword"},
                        "meta_authority_class": {"type": "keyword"},
                        "meta_nova_tier": {"type": "keyword"},
                        "meta_title": {"type": "text"},
                        "meta_status": {"type": "keyword"},
                        "meta_section": {"type": "keyword"},
                        # Temporal fields for as-of-date range queries (Rule 2 spec)
                        "meta_effective_date_start": {"type": "date", "format": "yyyy-MM-dd||epoch_millis", "ignore_malformed": True},
                        "meta_effective_date_end": {"type": "date", "format": "yyyy-MM-dd||epoch_millis", "ignore_malformed": True},
                        # Structural/semantic fields for gating and boosting
                        "meta_normative_weight": {"type": "keyword"},
                        "meta_structural_level": {"type": "keyword"},
                        "meta_is_appendix": {"type": "keyword"},
                        "meta_contains_definition": {"type": "keyword"},
                        "meta_contains_requirement": {"type": "keyword"},
                        "meta_document_class": {"type": "keyword"},
                    }
                }
            }
            self.es_client.indices.create(index=self.db_index_name, body=index_body)
            logger.info(f"Retriever: Created Elasticsearch index {self.db_index_name}")

            # Step 2: Configure index settings for bulk indexing
            self.es_client.indices.put_settings(
                index=self.db_index_name,
                body={"index": {"refresh_interval": "-1", "number_of_replicas": "0"}}
            )

            # Step 3: Bulk index directly
            # Promote key metadata fields to top-level so they're searchable
            # without nested queries. The full metadata dict is still stored
            # in the "metadata" object field for Rule 3 prompt injection.
            actions = []
            for doc, vector in zip(self.non_excel_docs, self.all_vectors):
                meta = self._normalize_metadata(doc.metadata)
                actions.append({
                    "_index": self.db_index_name,
                    "source": meta.get('source', ''),
                    "text": doc.page_content,
                    "vector": vector,
                    "metadata": doc.metadata,  # ORIGINAL preserved for provenance
                    # Promoted metadata fields
                    "meta_citation_anchor": str(meta.get('citation_anchor', '')),
                    "meta_guideline_number": str(meta.get('guideline_number', '')),
                    "meta_paragraph_role": str(meta.get('paragraph_role', '')),
                    "meta_regulator": str(meta.get('regulator', '')),
                    "meta_jurisdiction": str(meta.get('jurisdiction', '')),
                    "meta_authority_class": str(meta.get('authority_class', '')),
                    "meta_nova_tier": str(meta.get('nova_tier', '')),
                    "meta_title": str(meta.get('title', '')),
                    "meta_status": str(meta.get('status', '')),
                    "meta_section": str(meta.get('section', '')),
                    # Temporal fields for as-of-date filtering
                    "meta_effective_date_start": meta.get('effective_date_start') if meta.get('effective_date_start') else None,
                    "meta_effective_date_end": meta.get('effective_date_end') if meta.get('effective_date_end') else None,
                    # Structural/semantic fields
                    "meta_normative_weight": str(meta.get('normative_weight', '')),
                    "meta_structural_level": str(meta.get('structural_level', '')),
                    "meta_is_appendix": str(meta.get('is_appendix', '')).lower(),
                    "meta_contains_definition": str(meta.get('contains_definition', '')).lower(),
                    "meta_contains_requirement": str(meta.get('contains_requirement', '')).lower(),
                    "meta_document_class": str(meta.get('document_class', '')),
                })
            bulk(self.es_client, actions, chunk_size=500, request_timeout=120)

            logger.info(f"Retriever: Successfully indexed {len(actions)} documents to "
                       f"Elasticsearch index {self.db_index_name}")

            # Step 4: Re-enable refresh
            self.es_client.indices.put_settings(
                index=self.db_index_name,
                body={"index": {"refresh_interval": "1s", "number_of_replicas": "1"}}
            )

        except Exception as e:
            logger.error(f"Retriever: Error creating Elasticsearch vector store or connecting to "
                        f"the existing one: {e}")
            raise

    def connect_to_existing_store(self):
        """
        Connect to existing vector stores without re-indexing.
        Used for retrieval-only mode when indices already exist.
        """
        if self.model_provider_instance:
            embedding = self.model_provider_instance.get_embedding_model("retriever")
        else:
            embedding = self.embed_model

        self.es_client = self._get_es_client()

        try:
            if self.es_client and self.es_client.indices.exists(index=self.db_index_name):
                logger.info(f"Retriever: Connecting to existing Elasticsearch index "
                           f"{self.db_index_name}")

                vector_store = ElasticsearchVectorStore(
                    es_connection=self.es_client,
                    index_name=self.db_index_name,
                    embedding=embedding,
                )

                self.vector_store = vector_store

                bm25_retriever = ElasticsearchVectorStore(
                    es_connection=self.es_client,
                    index_name=self.db_index_name,
                    search_kwargs={"k": 5},
                )

                dense_retriever = vector_store.as_retriever()

                doc_count = self.es_client.count(index=self.db_index_name).get('count', 0)
                logger.info(f"Retriever: Connected to existing Elasticsearch index "
                           f"{self.db_index_name} with {doc_count} documents")
        except Exception as e:
            logger.error(f"Retriever: Error connecting to existing store: {e}")
            raise

    def build_index(self, documents=None):
        """Build the index from processed documents."""
        if documents is None:
            documents = self.non_excel_docs
        if not documents:
            logger.warning("No documents to index")
            return

        self.report_progress(85, "Generating embeddings...")

        self.all_vectors = []
        for doc in documents:
            try:
                vector = self.embed_model.get_text_embedding(doc.page_content)
                self.all_vectors.append(vector)
            except Exception as e:
                logger.error(f"Error generating embedding: {e}")
                self.all_vectors.append([0.0] * 1536)

        self.report_progress(90, "Indexing to Elasticsearch...")
        self._setup_es_index_and_store()

        # Smart wait for Elasticsearch to finish indexing
        self.report_progress(98, "Indexing documents in Elasticsearch...")
        max_wait = min(600, 60 + len(self.non_excel_docs) * 2)
        self._wait_for_indexing_complete(len(self.non_excel_docs), max_wait)

        self.report_progress(100, "Indexing complete")
        logger.info(f"Retriever: Successfully created and populated Elasticsearch index "
                   f"{self.db_index_name}")

    def _wait_for_indexing_complete(self, expected_doc_count, max_wait_time=120):
        """
        Smart wait for Elasticsearch indexing to complete.

        Args:
            expected_doc_count: Number of documents we expect to be indexed
            max_wait_time: Maximum time to wait in seconds
        """
        start_time = time.time()
        check_interval = 2

        while time.time() - start_time < max_wait_time:
            try:
                self.es_client.indices.refresh(index=self.db_index_name)
                doc_count = self.es_client.count(index=self.db_index_name)["count"]

                logger.info(f"Retriever: Index progress: {doc_count}/{expected_doc_count} documents indexed")

                if doc_count >= expected_doc_count:
                    logger.info(f"Retriever: Index ready with {doc_count} documents")
                    try:
                        search_result = self.es_client.search(
                            index=self.db_index_name,
                            body={"query": {"match_all": {}}, "size": 1}
                        )
                        if search_result["hits"]["total"]["value"] > 0:
                            logger.info("Retriever: Documents are searchable - indexing complete")
                            return True
                    except Exception as e:
                        logger.warning(f"Retriever: Search verification failed: {e}")

                time.sleep(check_interval)
            except Exception as e:
                logger.warning(f"Retriever: Index check failed: {e}")
                time.sleep(check_interval)

        logger.warning(f"Retriever: Indexing verification timed out after {max_wait_time}s, "
                      "proceeding anyway")
        return False

    def verify_indexing(self):
        """Verify that documents have been successfully indexed."""
        if not self.es_client:
            return False
        try:
            result = self.es_client.count(index=self.db_index_name)
            count = result.get('count', 0)
            logger.info(f"Retriever: Index {self.db_index_name} contains {count} documents")
            return count > 0
        except Exception as e:
            logger.error(f"Retriever: Error verifying index: {e}")
            return False

    # ========================================================================
    # RETRIEVAL WITH THREE-RULE METADATA ARCHITECTURE
    # ========================================================================

    def retrieve(self, query, k=10):
        """
        Retrieve relevant documents and format with NOVA metadata.

        Applies:
        - Rule 1: Metadata filtering (remove superseded/expired/draft)
        - Rule 2: Metadata boosting (rerank by authority/tier/jurisdiction)
        - Rule 3: Metadata injection (prompt-level headers for LLM reasoning)

        Args:
            query: User query string
            k: Number of results to return

        Returns:
            Dict with 'text' (formatted retrieved text) and 'sources' (source info)
        """
        if not self.es_client:
            raise ValueError("No Elasticsearch client available. Build or connect to an index first.")

        oversample_k = k * 3
        docs = self._search_elasticsearch(query, k=oversample_k)

        if not docs:
            logger.info("Retriever: No documents found")
            return {"text": "", "sources": []}

        logger.info(f"\n{'='*80}")
        logger.info(f"Retriever: Found {len(docs)} relevant document chunks")

        for i, doc in enumerate(docs):
            if 'source' in doc.metadata:
                source_path = doc.metadata['source']
                filename = os.path.basename(source_path)
            else:
                filename = 'unknown'
            content_preview = doc.page_content[:200].replace("\n", " ")
            logger.info(f"  Find Chunk #{i}: {filename}")
            logger.info(f"  Content: {content_preview}...")
        logger.info(f"{'='*80}\n")

        # RULE 2: Filter out superseded/expired/draft documents
        docs = self._apply_rule2_filters(docs, query)

        # RULE 2: Boost/rerank by authority, tier, and jurisdiction
        docs = self._apply_rule2_boosting(docs, query)

        # Temporal relevance reranking
        docs = self._rerank_by_temporal_relevance(docs, query)

        # Limit to requested k
        docs = docs[:k]

        # Build sources dict for grouped source references
        sources_dict = {}
        for i, doc in enumerate(docs):
            if 'source' in doc.metadata:
                source_path = doc.metadata['source']
                filename = os.path.basename(source_path)
            else:
                filename = 'unknown'
                source_path = 'unknown'

            # For PDF files, collect page information
            if filename.lower().endswith('.pdf') and 'page' in doc.metadata:
                page_num = doc.metadata.get('page', 0) + 1
                if filename not in sources_dict:
                    sources_dict[filename] = {
                        "pages": set(), "is_pdf": True, "is_md": False,
                        "file_path": doc.metadata.get('file_path', source_path),
                    }
                sources_dict[filename]['pages'].add(page_num)

            # For Markdown files, collect header information
            elif filename.lower().endswith('.md'):
                headers = {k: v for k, v in doc.metadata.items() if k.startswith("header")}
                if filename not in sources_dict:
                    sources_dict[filename] = {
                        "pages": set(), "is_pdf": False, "is_md": True,
                        "headers": set(),
                        "file_path": doc.metadata.get('file_path', source_path),
                    }
                if headers:
                    sorted_headers = sorted(headers.items(), key=lambda x: x[0])
                    cleaned_headers = []
                    for k, v in sorted_headers:
                        cleaned_value = re.sub(r'^\d+[\.\)]\s*', '', str(v).strip())
                        cleaned_value = cleaned_value.strip()
                        cleaned_headers.append(cleaned_value)
                    display_headers = cleaned_headers[-2:] if len(cleaned_headers) > 2 else cleaned_headers
                    header_path = " > ".join(display_headers)
                    sources_dict[filename]['headers'].add(header_path)
            else:
                if filename not in sources_dict:
                    sources_dict[filename] = {
                        "pages": set(), "is_pdf": False, "is_md": False,
                        "file_path": doc.metadata.get('file_path', source_path),
                    }

        # Format sources with grouped pages
        sources = []
        for filename, info in sources_dict.items():
            if filename in ['nova_master', 'glossary.txt']:
                continue
            file_path = info.get('file_path', filename)

            if info['is_pdf'] and info['pages']:
                sorted_pages = sorted(info['pages'])
                if len(sorted_pages) == 1:
                    sources.append({
                        "display_name": f"{filename} (page {sorted_pages[0]})",
                        "full_path": f"{file_path} (page {sorted_pages[0]})",
                    })
                else:
                    pages_str = ", ".join(str(p) for p in sorted_pages)
                    sources.append({
                        "display_name": f"{filename} (pages {pages_str})",
                        "full_path": f"{file_path} (pages {pages_str})",
                    })
            elif info.get('is_md') and info.get('headers'):
                sorted_headers = sorted(info['headers'])[:5]
                if len(sorted_headers) >= 1:
                    sources.append({
                        "display_name": f"{filename} ({sorted_headers[0]})",
                        "full_path": f"{file_path} ({sorted_headers[0]})",
                    })
                if len(sorted_headers) > 1:
                    headers_str = " | ".join(sorted_headers)
                    sources.append({
                        "display_name": f"{filename} ({headers_str})",
                        "full_path": file_path,
                    })
            else:
                sources.append({
                    "display_name": filename,
                    "full_path": file_path,
                })

        # RULE 3 METADATA
        # Format retrieved chunks with prompt-level metadata (Rule 3)
        formatted_chunks = []
        for i, doc in enumerate(docs[:k]):
            chunk_text = doc.page_content
            metadata = doc.metadata
            header_lines = []

            if metadata.get('title'):
                header_lines.append(f"Document: {metadata['title']}")
            if metadata.get('version_label'):
                header_lines.append(f"Version: {metadata['version_label']}")
            if metadata.get('status'):
                header_lines.append(f"Status: {metadata['status']}")

            effective_start = metadata.get('effective_date_start')
            effective_end = metadata.get('effective_date_end')
            if effective_start or effective_end:
                effective_str = f"Effective: {effective_start if effective_start else 'N/A'}"
                if effective_end:
                    effective_str += f" to {effective_end}"
                header_lines.append(effective_str)

            if metadata.get('authority_class'):
                header_lines.append(f"Authority: {metadata['authority_class']}")
            if metadata.get('authority_level'):
                header_lines.append(f"Authority Level: {metadata['authority_level']}")
            if metadata.get('nova_tier'):
                header_lines.append(f"NOVA Tier: {metadata['nova_tier']}")
            if metadata.get('jurisdiction'):
                header_lines.append(f"Jurisdiction: {metadata['jurisdiction']}")

            # Normative weight — critical for LLM to distinguish mandatory vs advisory
            # (Rule 3 spec: "mandatory" vs "advisory" fundamentally changes the answer)
            if metadata.get('normative_weight'):
                header_lines.append(f"Normative Weight: {metadata['normative_weight']}")

            # Paragraph role — tells LLM if this is a definition, exception, scope, etc.
            if metadata.get('paragraph_role'):
                header_lines.append(f"Paragraph Role: {metadata['paragraph_role']}")

            # Current version flag — LLM needs this to caveat superseded content
            current_flag = metadata.get('current_version_flag')
            if current_flag is not None:
                header_lines.append(f"Current Version: {current_flag}")

            # Approval status — draft vs approved changes how LLM presents the answer
            if metadata.get('approval_status'):
                header_lines.append(f"Approval Status: {metadata['approval_status']}")

            sector = metadata.get('sector')
            if sector:
                if isinstance(sector, list):
                    sector_str = ', '.join(str(s) for s in sector)
                else:
                    sector_str = str(sector)
                header_lines.append(f"Sector: {sector_str}")

            if metadata.get('regulator'):
                header_lines.append(f"Regulator: {metadata['regulator']}")
            if metadata.get('citation_anchor'):
                header_lines.append(f"Citation: {metadata['citation_anchor']}")

            if header_lines:
                header = "DOCUMENT METADATA:\n" + "\n".join(header_lines) + "\n" + "=" * 40 + "\n"
                formatted_chunks.append(header + chunk_text)
            else:
                formatted_chunks.append(chunk_text)

        retrieved_text = "\n\n".join(formatted_chunks)

        return {
            "text": retrieved_text,
            "sources": sources
        }

    def _extract_paragraph_ref(self, query):
        """
        Extract paragraph/section references from a query.
        Handles patterns like: 'paragraph 444', 'para 444', 'section 3.1',
        'article 12', 'clause 7.2', 'CRE31.4', 'MAR10.2'
        """
        patterns = [
            # "paragraph 444", "para 444", "para. 444"
            r'\b(?:paragraph|para\.?)\s*(\d+(?:\.\d+)*)',
            # "section 3.1", "sec 3.1"
            r'\b(?:section|sec\.?)\s*(\d+(?:\.\d+)*)',
            # "article 12"
            r'\b(?:article|art\.?)\s*(\d+(?:\.\d+)*)',
            # "clause 7.2"
            r'\b(?:clause|cl\.?)\s*(\d+(?:\.\d+)*)',
            # Basel-style codes: "CRE31.4", "MAR10.2", "SCO40.1"
            r'\b([A-Z]{2,4}\d+(?:\.\d+)*)\b',
        ]
        refs = []
        for pattern in patterns:
            matches = re.findall(pattern, query, re.IGNORECASE)
            refs.extend(matches)
        return refs

    def _search_elasticsearch(self, query, k=30):
        """
        Search Elasticsearch with hybrid vector + BM25 query.

        Implements Rule 2 of the NOVA architecture:
        - bool filter clauses gate results (status, jurisdiction, as-of-date)
        - should clauses boost results (authority, paragraph refs, content flags)
        - script_score provides dense vector similarity
        """
        docs = []
        try:
            query_vector = self.embed_model.get_text_embedding(query)

            # ---- SHOULD clauses (boosting, not gating) ----
            should_clauses = [
                {
                    "script_score": {
                        "query": {"match_all": {}},
                        "script": {
                            "source": "cosineSimilarity(params.query_vector, 'vector') + 1.0",
                            "params": {"query_vector": query_vector}
                        }
                    }
                },
                {
                    "match": {
                        "text": {
                            "query": query,
                            "boost": self.NOVA_MATCH_BOOST
                        }
                    }
                },
                # Boost requirements if query sounds like it's asking for obligations
                {"term": {"meta_normative_weight": {"value": "mandatory", "boost": 2.0}}},
            ]

            # ---- FILTER clauses (hard gates — Rule 2 spec) ----
            # These run in filter context (no scoring, cached, fast).
            filter_clauses = []

            # As-of-date filtering: only return docs effective at the current date
            # unless the query explicitly references a historical year.
            # Spec: "effective_date_start/end govern which version surfaces"
            query_year = self._extract_year_from_query(query)
            if query_year and query_year != datetime.now().year:
                # Historical query — use as-of-date from the query year
                as_of_date = f"{query_year}-12-31"
                logger.info(f"Retriever: As-of-date filter for historical query: {as_of_date}")
            else:
                as_of_date = None  # Don't gate by date for current queries

            if as_of_date:
                filter_clauses.append(
                    {"range": {"meta_effective_date_start": {"lte": as_of_date}}}
                )
                filter_clauses.append({
                    "bool": {
                        "should": [
                            {"bool": {"must_not": {"exists": {"field": "meta_effective_date_end"}}}},
                            {"range": {"meta_effective_date_end": {"gte": as_of_date}}}
                        ]
                    }
                })

            # Auto-detect regulator from query and add as filter if found
            auto_filters = auto_detect_filters_from_query(query)
            if auto_filters.get('regulator'):
                filter_clauses.append(
                    {"term": {"meta_regulator": auto_filters['regulator']}}
                )
                logger.info(f"Retriever: Auto-detected regulator filter: {auto_filters['regulator']}")

            # ---- Paragraph/section reference boosting ----
            para_refs = self._extract_paragraph_ref(query)
            if para_refs:
                logger.info(f"Retriever: Detected paragraph/section references: {para_refs}")
                for ref in para_refs:
                    should_clauses.append({
                        "match_phrase": {"text": {"query": f"paragraph {ref}", "boost": 5.0}}
                    })
                    should_clauses.append({
                        "match_phrase": {"text": {"query": ref, "boost": 3.0}}
                    })
                    should_clauses.append({
                        "match": {"meta_citation_anchor": {"query": ref, "boost": 10.0}}
                    })
                    should_clauses.append({
                        "match": {"meta_guideline_number": {"query": ref, "boost": 8.0}}
                    })
                    should_clauses.append({
                        "match": {"meta_paragraph_role": {"query": ref, "boost": 5.0}}
                    })
                    should_clauses.append({
                        "match": {"meta_section": {"query": ref, "boost": 4.0}}
                    })

            # ---- Content flag boosting from query intent ----
            if auto_filters.get('contains_definition'):
                should_clauses.append(
                    {"term": {"meta_contains_definition": {"value": "true", "boost": 3.0}}}
                )
            if auto_filters.get('contains_requirement'):
                should_clauses.append(
                    {"term": {"meta_contains_requirement": {"value": "true", "boost": 3.0}}}
                )

            # ---- Structural boosting: prefer sections over appendices ----
            should_clauses.append(
                {"term": {"meta_is_appendix": {"value": "false", "boost": 1.2}}}
            )

            # ---- Assemble final query ----
            bool_query = {"should": should_clauses}
            if filter_clauses:
                bool_query["filter"] = filter_clauses

            search_body = {
                "size": k,
                "query": {"bool": bool_query}
            }

            result = self.es_client.search(index=self.db_index_name, body=search_body)
            for hit in result['hits']['hits']:
                source = hit['_source']
                metadata = source.get('metadata', {})
                metadata = self._normalize_metadata(metadata)
                metadata['_score'] = hit['_score']
                doc = Document(page_content=source.get('text', ''), metadata=metadata)
                docs.append(doc)
        except Exception as e:
            logger.error(f"Retriever: Error searching Elasticsearch: {e}")
        return docs

    # ========================================================================
    # RULE 2: METADATA FILTERING
    # ========================================================================

    # Regulator versioning behaviour.
    # 'lifecycle'     – new guidance supersedes old; safe to filter out old versions.
    #                   Examples: OSFI, PRA, BOE, OCC
    # 'publication'   – docs reference/build on predecessors; old docs stay relevant.
    #                   Examples: BIS, BCBS (Basel), IOSCO, FSB
    # 'mixed'         – binding standards are versioned but other outputs accumulate.
    #                   Examples: US Fed, EBA, SEC
    REGULATOR_VERSION_MODEL = {
        # Lifecycle regulators — apply all filters
        'osfi': 'lifecycle',
        'pra': 'lifecycle',
        'bank of england': 'lifecycle',
        'boe': 'lifecycle',
        'occ': 'lifecycle',
        'fdic': 'lifecycle',
        # Publication-model regulators — never filter by version/status/expiry
        'bis': 'publication',
        'bcbs': 'publication',
        'basel': 'publication',
        'bank for international settlements': 'publication',
        'basel committee': 'publication',
        'iosco': 'publication',
        'fsb': 'publication',
        'financial stability board': 'publication',
        # Mixed regulators — filter superseded rules but keep circulars/guidance
        'federal reserve': 'mixed',
        'fed': 'mixed',
        'eba': 'mixed',
        'sec': 'mixed',
        'ecb': 'mixed',
    }

    # Document classes that should NOT be filtered even under lifecycle regulators.
    # These are reference/foundational materials, not versioned rules.
    KEEP_DOCUMENT_CLASSES = {
        'circular', 'consultation', 'discussion_paper', 'research',
        'working_paper', 'speech', 'report', 'faq',
    }

    def _get_regulator_version_model(self, metadata):
        """
        Determine the versioning model for a document based on its regulator.
        Returns 'lifecycle', 'publication', or 'mixed'.
        Defaults to 'lifecycle' for unknown regulators (safe — filters apply).
        """
        doc_regulator = str(metadata.get('regulator', '')).lower().strip()
        doc_source_type = str(metadata.get('source_type', '')).lower().strip()

        # Direct match
        if doc_regulator in self.REGULATOR_VERSION_MODEL:
            return self.REGULATOR_VERSION_MODEL[doc_regulator]

        # Partial match (e.g. "Basel Committee on Banking Supervision" contains "basel")
        for key, model in self.REGULATOR_VERSION_MODEL.items():
            if key in doc_regulator:
                return model

        # Source type fallback
        if doc_source_type in ('international_standard', 'publication', 'framework'):
            return 'publication'

        return 'lifecycle'

    def _apply_rule2_filters(self, docs: List[Document], query: str) -> List[Document]:
        """
        Apply Rule 2 metadata filtering to remove unwanted documents.

        Filtering is regulator-aware:
        - lifecycle regulators (OSFI, PRA): filter superseded/inactive/expired docs
        - publication regulators (Basel, BIS): keep all docs, rely on temporal ranking
        - mixed regulators (US Fed, EBA): filter superseded binding rules but keep
          circulars, consultations, and other accumulating document types

        Args:
            docs: List of retrieved documents
            query: User query (for context)

        Returns:
            Filtered list of documents
        """
        if not docs:
            return docs

        filtered_docs = []
        current_date = datetime.now()

        logger.info(f"\n{'='*80}")
        logger.info("RULE 2 FILTERING - Regulator-aware version filtering")
        logger.info(f"{'='*80}")

        for doc in docs:
            metadata = doc.metadata
            should_keep = True
            filter_reason = None
            version_model = self._get_regulator_version_model(metadata)
            doc_regulator = str(metadata.get('regulator', '')).strip()
            doc_class = str(metadata.get('document_class', '')).lower().strip()

            # ---- Publication model: keep everything ----
            if version_model == 'publication':
                logger.debug(f"  [publication] Keeping: "
                           f"{os.path.basename(metadata.get('source', 'unknown'))} "
                           f"(regulator={doc_regulator})")
                filtered_docs.append(doc)
                continue

            # ---- Mixed model: keep non-binding/accumulating doc types ----
            if version_model == 'mixed' and doc_class in self.KEEP_DOCUMENT_CLASSES:
                logger.debug(f"  [mixed] Keeping non-binding doc: "
                           f"{os.path.basename(metadata.get('source', 'unknown'))} "
                           f"(class={doc_class})")
                filtered_docs.append(doc)
                continue

            # ---- Lifecycle / mixed binding docs: apply all filters ----

            # Filter 1: Exclude superseded documents (prefer current versions)
            current_version = metadata.get('current_version_flag')
            if current_version is not None:
                is_not_current = (
                    current_version is False or
                    str(current_version).lower() in ['false', '0', 'no']
                )
                if is_not_current:
                    should_keep = False
                    filter_reason = f"Superseded version (current_version_flag={False})"

            # Filter 2: Exclude inactive/draft documents (prefer active)
            status = str(metadata.get('status', '')).lower()
            # Note: "future_effective" (from OSFI "final_future_effective") is
            # intentionally NOT filtered — approved but not yet in force docs
            # are relevant for planning queries.
            if status in ['draft', 'superseded', 'inactive', 'withdrawn']:
                should_keep = False
                filter_reason = f"Status: {status}"

            # Filter 3: Check effective date range (exclude expired documents)
            effective_end = metadata.get('effective_date_end')
            if effective_end is not None and str(effective_end).lower() not in ['na', 'none', 'null', '',
                'false']:
                try:
                    end_date = None
                    if re.match(r'^\d{4}-\d{2}-\d{2}', str(effective_end)):
                        end_date = datetime.strptime(effective_end[:10], '%Y-%m-%d')
                    else:
                        try:
                            from dateutil import parser as date_parser
                            end_date = date_parser.parse(str(effective_end))
                        except ImportError:
                            logger.debug("dateutil not available for date parsing")
                    if end_date and end_date < current_date:
                        should_keep = False
                        filter_reason = f"Expired (effective_date_end: {effective_end})"
                except Exception as e:
                    logger.debug(f"Could not parse date '{effective_end}': {e}")
                    pass

            if should_keep:
                filtered_docs.append(doc)
            else:
                source = os.path.basename(metadata.get('source', 'unknown'))
                logger.info(f"  [{version_model}] Filtered out {source}: {filter_reason}")

        logger.info(f"Filtered {len(docs)} -> {len(filtered_docs)} documents "
                    f"(removed {len(docs) - len(filtered_docs)})")
        logger.info(f"{'='*80}\n")
        return filtered_docs

    # ========================================================================
    # RULE 2: METADATA BOOSTING
    # ========================================================================

    def _apply_rule2_boosting(self, docs: List[Document], query: str) -> List[Document]:
        """
        Apply Rule 2 metadata boosting to rerank documents by authority and classification.

        Boosts: higher authority, better NOVA tier, jurisdiction match.

        Args:
            docs: List of filtered documents
            query: User query (for jurisdiction detection)

        Returns:
            Reranked list of documents
        """
        if not docs:
            return docs

        logger.info(f"\n{'='*80}")
        logger.info(f"RULE 2 BOOSTING - Ranking by authority, tier, and jurisdiction")
        logger.info(f"{'='*80}")

        # Authority hierarchy scoring.
        # Scores reflect cross-regulator ranking within the NOVA knowledge base.
        # The knowledge base contains docs from OSFI, Basel/BCBS, US Fed, PRA/BOE.
        #
        # Values sourced from actual authority_class fields across all 4 repos:
        #   OSFI (278 docs): primary_normative, official_support, official_interpretive,
        #                     contextual_summary, excluded
        #   US Fed eCFR (59 docs): primary_normative, procedural_administrative,
        #                          reference_interpretive
        #   US Fed FR (3037 docs): guidance_interpretive, primary_normative
        #   Basel (285 docs): guidance_interpretive, primary_normative,
        #                     reference_definitional
        #   PRA (677 docs): interpretive, primary_normative, context, consultative
        #
        # Tier A (100) — Primary binding instruments. These ARE the rules.
        #   OSFI: primary_normative (B-20, CAR guidelines)
        #   Basel: primary_normative (LCR standard, CRE framework)
        #   US Fed: primary_normative (eCFR regulations, FR Final Rules)
        #   PRA: primary_normative (PRA Rulebook rules)
        #
        # Tier B (85) — Official support / formal guidance.
        #   OSFI: official_support (transmittal letters, supporting docs)
        #   US Fed: procedural_administrative (rules of procedure)
        #   Basel: guidance_interpretive (supervisory guidelines)
        #   PRA: interpretive (supervisory statements — 446 docs, bulk of PRA)
        #
        # Tier C (60) — Interpretive / reference materials.
        #   OSFI: official_interpretive (assessment tools)
        #   US Fed: guidance_interpretive (FR interpretations — 2283 docs)
        #   US Fed: reference_interpretive (miscellaneous interpretations)
        #   Basel: reference_definitional (definitions, glossary)
        #
        # Tier D (35) — Consultative / contextual. Draft or reference only.
        #   OSFI: contextual_summary (guideline-at-a-glance summaries)
        #   PRA: context (glossary, guidance index)
        #   PRA: consultative (consultation papers — not yet finalized)
        #
        # Tier E (10) — Excluded. Not for retrieval.
        #   OSFI: excluded (marked for exclusion from corpus)
        authority_scores = {
            # Tier A: Primary binding instruments (score 100)
            'primary_normative': 100,
            'primary_law': 100,
            'statute': 100,
            'regulation': 100,
            'guideline': 100,
            # Tier B: Official support / formal guidance (score 85)
            'official_support': 85,
            'primary_guidance': 85,
            'guidance': 85,
            'directive': 85,
            'procedural_administrative': 85,
            # Tier C: Interpretive / reference (score 60)
            'guidance_interpretive': 60,
            'interpretive': 60,
            'official_interpretive': 60,
            'reference_interpretive': 60,
            'reference_definitional': 55,
            'secondary_normative': 65,
            'secondary_guidance': 60,
            'policy': 55,
            # Tier D: Consultative / contextual (score 35)
            'advisory': 35,
            'consultative': 35,
            'context': 35,
            'contextual_summary': 35,
            # Tier E: Informational / excluded (score 10-20)
            'informational': 20,
            'excluded': 10,
        }

        # NOVA Tier scoring - handle both integers and strings
        tier_scores = {
            1: 100, 'Tier 1': 100, '1': 100,
            2: 70, 'Tier 2': 70, '2': 70,
            3: 40, 'Tier 3': 40, '3': 40
        }

        # Detect jurisdiction from query (basic detection)
        query_lower = query.lower()

        jurisdiction_keywords = {
            'canada': ['canada', 'canadian', 'osfi', 'ontario', 'quebec', 'alberta', 'bc'],
            'us': ['us', 'usa', 'united states', 'federal reserve', 'fed', 'occ'],
            'uk': ['uk', 'united kingdom', 'britain', 'british', 'boe', 'pra'],
            'eu': ['eu', 'europe', 'european', 'eba', 'ecb', 'esma'],
            'global': ['basel', 'bcbs', 'international', 'global']
        }

        detected_jurisdiction = None
        for jurisdiction, keywords in jurisdiction_keywords.items():
            if any(kw in query_lower for kw in keywords):
                detected_jurisdiction = jurisdiction
                break

        if detected_jurisdiction:
            logger.info(f"Detected jurisdiction from query: {detected_jurisdiction}")

        scored_docs = []
        for doc in docs:
            metadata = doc.metadata
            boost_score = 0
            boost_details = []

            # Authority class boost
            authority_class_raw = metadata.get('authority_class', '')
            authority_class = str(authority_class_raw).strip().lower()
            auth_score = authority_scores.get(authority_class, 0)
            boost_score += auth_score
            boost_details.append(f"Authority({authority_class})={auth_score}")

            # NOVA Tier boost - handle integer or string
            nova_tier_raw = metadata.get('nova_tier')
            if nova_tier_raw is not None and nova_tier_raw != '':
                tier_score = tier_scores.get(nova_tier_raw)
                if tier_score is None:
                    tier_score = tier_scores.get(str(nova_tier_raw).lower(), 0)
                boost_score += tier_score
                boost_details.append(f"Tier({nova_tier_raw})={tier_score}")

            # Jurisdiction match boost
            doc_jurisdiction = metadata.get('jurisdiction', '')
            doc_jur_str = str(doc_jurisdiction).lower() if doc_jurisdiction else ''
            if detected_jurisdiction:
                if detected_jurisdiction in doc_jur_str or any(kw in doc_jur_str for kw in
                    jurisdiction_keywords.get(detected_jurisdiction, [])):
                    boost_score += 50
                    boost_details.append(f"Jurisdiction match={50}")

            # Authority level boost (if available)
            authority_level = metadata.get('authority_level')
            if authority_level:
                try:
                    level = int(authority_level)
                    level_boost = level * 10
                    boost_score += level_boost
                    boost_details.append(f"AuthLevel({level})={level_boost}")
                except (ValueError, TypeError):
                    pass

            scored_docs.append((doc, boost_score, boost_details))

        scored_docs.sort(key=lambda x: x[1], reverse=True)

        # Log boosting results
        for i, (doc, score, details) in enumerate(scored_docs[:10]):
            source = os.path.basename(doc.metadata.get('source', 'unknown'))
            details_str = ' | '.join(details) if details else "no boost"
            logger.info(f"  #{i}: Score: {score:3.0f} - {source} ({details_str})")

        logger.info(f"{'='*80}\n")
        return [doc for doc, score, details in scored_docs]

    # ========================================================================
    # TEMPORAL RELEVANCE
    # ========================================================================

    def _extract_year_from_query(self, query: str) -> Optional[int]:
        """Extract year from query using efficient regex patterns."""
        year_match = re.search(r'\b(19|20)\d{2}\b', query)
        if year_match:
            return int(year_match.group())

        quarter_match = re.search(r'Q[1-4]\s*(19|20)\d{2}', query)
        if quarter_match:
            year_match = re.search(r'(19|20)\d{2}', quarter_match.group())
            if year_match:
                return int(year_match.group())
        return None

    def _extract_year_from_metadata(self, doc: Document) -> Optional[int]:
        """Extract year from document metadata or content using efficient patterns."""
        metadata = doc.metadata

        for field in ['effective_date_start', 'effective_date', 'date', 'year', 'version_date',
            'publish_date']:
            if field in metadata:
                value = metadata[field]
                if value is None:
                    continue
                value_str = str(value)
                year_match = re.search(r'\b(19|20)\d{2}\b', value_str)
                if year_match:
                    return int(year_match.group())

        source = metadata.get('source', '')
        year_match = re.search(r'\b(19|20)\d{2}\b', source)
        if year_match:
            return int(year_match.group())

        content = doc.page_content[:1000]

        # Prioritize "effective [year]" patterns (case-insensitive)
        effective_match = re.search(
            r'\beffective\b.*?\b(19|20)\d{2}\b|'
            r'\b(january|february|march|april|may|june|july|august|'
            r'september|october|november|december)\b.*?\b(19|20)\d{2}\b',
            content, re.IGNORECASE
        )
        if effective_match:
            year_match = re.search(r'(19|20)\d{2}', effective_match.group())
            if year_match:
                return int(year_match.group())

        year_match = re.search(r'\b(19|20)\d{2}\b', content)
        if year_match:
            return int(year_match.group())
        return None

    def _rerank_by_temporal_relevance(self, docs: List[Document], query: str) -> List[Document]:
        """
        Rerank documents by temporal relevance to the query.

        Temporal decay is regulator-aware:
        - lifecycle regulators (OSFI, PRA): steep decay — 10 pts/year.
          A 2024 OSFI guideline should strongly outrank a 2019 version.
        - publication regulators (Basel, BIS): gentle decay — 3 pts/year.
          Basel II (2004) is still foundational even though Basel III (2017) exists.
          Old publications stay highly relevant because newer ones build on them.
        - mixed regulators (US Fed, EBA): moderate decay — 6 pts/year.
        """
        query_year = self._extract_year_from_query(query)
        current_year = datetime.now().year

        if not query_year:
            query_year = current_year

        # Decay rates per version model (points lost per year of distance from query year)
        DECAY_RATES = {
            'lifecycle': 10,     # Steep: old OSFI/PRA docs penalized heavily
            'mixed': 6,          # Moderate: old Fed rules penalized somewhat
            'publication': 3,    # Gentle: old Basel/BIS docs barely penalized
        }

        scored_docs = []
        for doc in docs:
            doc_year = self._extract_year_from_metadata(doc)
            version_model = self._get_regulator_version_model(doc.metadata)
            decay_rate = DECAY_RATES.get(version_model, 10)

            if doc_year:
                year_diff = abs(query_year - doc_year)
                temporal_score = max(100 - (year_diff * decay_rate), 0)
            else:
                # No year found — neutral score
                temporal_score = 50

            temporal_score = max(temporal_score, 0)
            scored_docs.append((doc, temporal_score))

        scored_docs.sort(key=lambda x: x[1], reverse=True)
        return [doc for doc, score in scored_docs]

    # ========================================================================
    # PROCESS AND INDEX (END-TO-END)
    # ========================================================================

    def process_and_index(self, file_paths=None):
        """
        End-to-end: load documents, process them, and build the index.
        """
        if file_paths:
            self.file_paths = file_paths
        if not self.file_paths:
            logger.warning("No file paths provided for processing")
            return
        try:
            documents = self.load_and_split_documents()
            if not documents:
                logger.warning("No documents were loaded")
                return
            self.build_index(documents)
            logger.info(f"Successfully processed and indexed {len(documents)} documents")
        except Exception as e:
            logger.error(f"Error in process_and_index: {e}")
            raise

    def get_index_stats(self):
        """Get statistics about the current index."""
        stats = {
            "has_index": self.index is not None,
            "has_es_client": self.es_client is not None,
            "index_name": self.db_index_name,
            "num_files": len(self.file_paths),
            "num_chunks": len(self.non_excel_docs),
            "num_excel_docs": len(self.excel_docs),
        }
        if self.es_client:
            try:
                result = self.es_client.count(index=self.db_index_name)
                stats["es_doc_count"] = result.get('count', 0)
            except Exception:
                stats["es_doc_count"] = "unknown"
        return stats
